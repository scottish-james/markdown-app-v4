import os
import tempfile
import streamlit as st
import requests
from markitdown import MarkItDown
from openai import OpenAI
from bs4 import BeautifulSoup
from urllib.parse import urlparse
from pptx import Presentation  # For PowerPoint hyperlink extraction
import fitz  # PyMuPDF for PDF processing
import re  # For regex pattern matching
import glob  # For directory file listing
import shutil  # For file operations
from pathlib import Path  # For path operations

# Set page configuration
st.set_page_config(
    page_title="Office to Markdown",
    page_icon="📄",
    layout="centered",
)

# Simplified styling that works reliably
st.markdown("""
<style>
    .main {
        background-color: #f5f7f9;
    }
    .stButton button {
        background-color: #4e54c8;
        color: white;
        font-weight: bold;
    }
</style>
""", unsafe_allow_html=True)


# Helper functions
def get_file_extension(filename):
    """Extract the file extension from a filename."""
    return filename.rsplit(".", 1)[1].lower() if "." in filename else ""


def fix_url(url):
    """
    Fix URLs by adding appropriate schemes if missing.

    Args:
        url: The URL to fix

    Returns:
        Properly formatted URL
    """
    if not url:
        return url

    # For email addresses
    if '@' in url and not url.startswith('mailto:'):
        return f"mailto:{url}"

    # For web URLs
    if not url.startswith(('http://', 'https://', 'mailto:', 'tel:', 'ftp://', '#')):
        if url.startswith('www.') or any(
                domain in url.lower() for domain in ['.com', '.org', '.net', '.edu', '.gov', '.io']):
            return f"https://{url}"

    return url


def extract_pptx_hyperlinks(pptx_path):
    """
    Extract hyperlinks from a PowerPoint file.

    Args:
        pptx_path: Path to the PowerPoint file

    Returns:
        List of dictionaries containing text, URL, and slide number
    """
    try:
        prs = Presentation(pptx_path)
        raw_hyperlinks = []

        def process_shape(shape, slide_num, parent_type=""):
            """Recursively process shapes, including those in groups"""
            try:
                # Check for hyperlinks in shape click actions
                if hasattr(shape, 'click_action') and shape.click_action is not None:
                    if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink is not None:
                        if shape.click_action.hyperlink.address:
                            shape_text = shape.text if hasattr(shape,
                                                               'text') else f"{parent_type} Shape on Slide {slide_num}"
                            raw_hyperlinks.append({
                                'text': shape_text,
                                'url': shape.click_action.hyperlink.address,
                                'slide': slide_num
                            })
            except Exception as e:
                # Just skip this shape if there's an error accessing click_action
                pass

            # Process text frames that might contain hyperlinks
            if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    # Create a unique ID for this paragraph
                    p_id = f"{parent_type}shape_{id(shape)}_p_{p_idx}"
                    paragraph_links = {}

                    for run in paragraph.runs:
                        if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                            url = run.hyperlink.address
                            if url not in paragraph_links:
                                paragraph_links[url] = []
                            paragraph_links[url].append(run.text)

                    # Add merged paragraph links
                    for url, text_parts in paragraph_links.items():
                        merged_text = "".join(text_parts)
                        raw_hyperlinks.append({
                            'text': merged_text,
                            'url': url,
                            'slide': slide_num
                        })

            # Process tables
            if hasattr(shape, 'table') and shape.table is not None:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        if hasattr(cell, 'text_frame') and cell.text_frame is not None:
                            for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                p_id = f"{parent_type}table_{id(shape)}_row_{r_idx}_cell_{c_idx}_p_{p_idx}"
                                paragraph_links = {}

                                for run in paragraph.runs:
                                    if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                        url = run.hyperlink.address
                                        if url not in paragraph_links:
                                            paragraph_links[url] = []
                                        paragraph_links[url].append(run.text)

                                # Add merged paragraph links
                                for url, text_parts in paragraph_links.items():
                                    merged_text = "".join(text_parts)
                                    raw_hyperlinks.append({
                                        'text': merged_text,
                                        'url': url,
                                        'slide': slide_num
                                    })

            # Recursively process group shapes
            if hasattr(shape, 'shapes') and shape.shapes is not None:
                try:
                    # For GroupShape objects
                    for i, child_shape in enumerate(shape.shapes):
                        process_shape(child_shape, slide_num, f"{parent_type}Group{i}_")
                except Exception as e:
                    # Skip if there's an issue accessing child shapes
                    pass

            # Handle SmartArt and Charts which might have text with hyperlinks
            if hasattr(shape, 'has_chart') and shape.has_chart:
                # Charts may have text with hyperlinks in their titles or labels
                try:
                    chart = shape.chart
                    if hasattr(chart, 'chart_title') and chart.chart_title is not None:
                        # Some charts have text properties in their titles
                        if hasattr(chart.chart_title, 'text_frame'):
                            for p_idx, paragraph in enumerate(chart.chart_title.text_frame.paragraphs):
                                for run in paragraph.runs:
                                    if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                        raw_hyperlinks.append({
                                            'text': run.text,
                                            'url': run.hyperlink.address,
                                            'slide': slide_num
                                        })
                except Exception:
                    # Skip if can't access chart data
                    pass

            # Handle SmartArt
            if hasattr(shape, 'has_smart_art') and shape.has_smart_art:
                # Some SmartArt structures may contain text with hyperlinks
                try:
                    # Process SmartArt nodes if they have text frames
                    if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            for run in paragraph.runs:
                                if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                    raw_hyperlinks.append({
                                        'text': run.text,
                                        'url': run.hyperlink.address,
                                        'slide': slide_num
                                    })
                except Exception:
                    # Skip if can't access SmartArt data
                    pass

        # Process all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            # Process all shapes in the slide
            for shape in slide.shapes:
                process_shape(shape, slide_num)

        # Final cleanup - remove duplicate links and empty text links
        hyperlinks = []
        seen_urls = {}  # Format: {(url, slide): best_text}

        for link in raw_hyperlinks:
            # Skip empty links
            text = link['text'].strip()
            if not text:
                continue

            # Keep track of the best text for each URL (prefer longer, more descriptive text)
            link_id = (link['url'], link['slide'])
            if link_id not in seen_urls or len(text) > len(seen_urls[link_id]):
                seen_urls[link_id] = text

        # Create final hyperlinks list with best text for each URL
        for (url, slide), text in seen_urls.items():
            hyperlinks.append({
                'text': text,
                'url': url,
                'slide': slide
            })

        return hyperlinks
    except Exception as e:
        st.warning(f"Error extracting PowerPoint hyperlinks: {str(e)}")
        return []


def extract_pdf_hyperlinks(pdf_path):
    """
    Extract hyperlinks from a PDF file using PyMuPDF (fitz).

    Args:
        pdf_path: Path to the PDF file

    Returns:
        List of dictionaries containing text, URL, and page number
    """
    try:
        hyperlinks = []

        # Open the PDF file with PyMuPDF
        doc = fitz.open(pdf_path)

        # Process each page
        for page_num, page in enumerate(doc, 1):
            # Get links from the page
            links = page.get_links()

            for link in links:
                # Process based on link type
                if link["kind"] == fitz.LINK_URI:  # External URI
                    url = link["uri"]

                    # Try to get text near the link location
                    rect = fitz.Rect(link["from"])
                    # Create an expanded rectangle instead of using inflate method
                    expanded_rect = fitz.Rect(rect.x0 - 10, rect.y0 - 10, rect.x1 + 10, rect.y1 + 10)
                    words = page.get_text("words", clip=expanded_rect)  # Get words near the link

                    # Combine words to form text
                    if words:
                        text = " ".join([w[4] for w in words])
                        # Clean up the text
                        text = text.strip()
                        if not text:
                            text = f"Link on page {page_num}"
                    else:
                        text = f"Link on page {page_num}"

                    hyperlinks.append({
                        'text': text,
                        'url': url,
                        'page': page_num
                    })

                elif link["kind"] == fitz.LINK_GOTO:  # Internal link
                    # The page number this links to
                    dest_page = link["page"] + 1  # 0-based to 1-based

                    # Try to get text near the link location
                    rect = fitz.Rect(link["from"])
                    # Create an expanded rectangle instead of using inflate method
                    expanded_rect = fitz.Rect(rect.x0 - 10, rect.y0 - 10, rect.x1 + 10, rect.y1 + 10)
                    words = page.get_text("words", clip=expanded_rect)

                    if words:
                        text = " ".join([w[4] for w in words])
                        text = text.strip()
                        if not text:
                            text = f"Go to page {dest_page}"
                    else:
                        text = f"Go to page {dest_page}"

                    hyperlinks.append({
                        'text': text,
                        'url': f"#page={dest_page}",
                        'page': page_num
                    })

            # Also look for URLs in the page text
            text = page.get_text()
            if text:
                # Find URLs using regex
                url_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+'
                urls = re.findall(url_pattern, text)

                for url in urls:
                    # Check if we already have this URL for this page
                    if not any(link.get('url') == url and link.get('page') == page_num for link in hyperlinks):
                        # Get some context around the URL
                        try:
                            # Find position of URL in text
                            pos = text.find(url)
                            # Get some text before and after
                            start = max(0, pos - 50)
                            end = min(len(text), pos + len(url) + 50)
                            context = text[start:end].replace(url, "").strip()

                            # If context is too long or empty, use a default
                            if len(context) > 100 or not context:
                                context = f"URL on page {page_num}"
                        except:
                            context = f"URL on page {page_num}"

                        hyperlinks.append({
                            'text': context,
                            'url': url,
                            'page': page_num
                        })

        return hyperlinks
    except Exception as e:
        st.warning(f"Error extracting PDF hyperlinks: {str(e)}")
        return []


def format_hyperlinks_section(hyperlinks, container_name="Document"):
    """
    Format hyperlinks into a markdown section.

    Args:
        hyperlinks: List of hyperlink dictionaries
        container_name: Name of the container (Presentation, Document, etc.)

    Returns:
        Formatted markdown string for hyperlinks section
    """
    if not hyperlinks:
        return ""

    # Group hyperlinks by page/slide
    pages_with_links = {}
    url_tracker = {}  # Track URLs we've seen per page

    for link in hyperlinks:
        # Get page number (slide for PowerPoint, page for PDF)
        page_num = link.get('page', link.get('slide', 0))
        if not page_num:
            continue

        # Get URL and text
        url = link.get('url', '')
        text = link.get('text', '').strip()

        # Skip empty entries
        if not url or not text:
            continue

        # Initialize page tracking if needed
        if page_num not in pages_with_links:
            pages_with_links[page_num] = []
            url_tracker[page_num] = set()

        # Fix URL format
        url = fix_url(url)

        # Check for duplicates on this page
        url_key = url.lower()
        if url_key not in url_tracker[page_num]:
            pages_with_links[page_num].append({
                'text': text,
                'url': url
            })
            url_tracker[page_num].add(url_key)
        else:
            # We already have this URL on this page
            # Check if the new text is better than the existing one
            for existing_link in pages_with_links[page_num]:
                if existing_link['url'].lower() == url_key:
                    # Replace with new text if it's better
                    if (len(text) > len(existing_link['text']) and
                        existing_link['text'] in text) or (
                            existing_link['text'].startswith(('Link on page', 'URL on page', 'Go to page'))
                    ):
                        existing_link['text'] = text

    # Exit if no valid links found
    if not any(links for links in pages_with_links.values()):
        return ""

    # Build the markdown
    page_title = "Slide" if container_name == "Presentation" else "Page"
    markdown = f"\n\n## Hyperlinks in {container_name}\n"

    for page_num in sorted(pages_with_links.keys()):
        links = pages_with_links[page_num]
        if not links:
            continue

        markdown += f"\n### {page_title} {page_num}\n"
        for link in links:
            markdown += f"* [{link['text']}]({link['url']})\n"

    return markdown


def enhance_markdown_formatting(markdown_content, api_key=None):
    """
    Enhance markdown formatting using OpenAI's GPT-4o model.

    Args:
        markdown_content: The raw markdown content to enhance
        api_key: OpenAI API key (optional if stored in environment)

    Returns:
        Tuple of (enhanced_markdown, error_message)
    """
    try:
        if not api_key:
            # Try to get API key from environment or Streamlit secrets
            api_key = os.getenv("OPENAI_API_KEY") or st.secrets.get("openai_api_key")

        if not api_key:
            return markdown_content, "No OpenAI API key provided"

        # Initialize OpenAI client
        client = OpenAI(api_key=api_key)

        # Create a prompt for markdown enhancement
        system_prompt = """You are a markdown formatting expert. Your task is to:
        1. Fix any syntax errors in the markdown
        2. Improve the structure and hierarchy of headers
        3. Ensure consistent formatting throughout
        4. Enhance bullet points and numbered lists
        5. Properly format tables and code blocks
        6. Add appropriate spacing between sections
        7. Maintain the original content without adding new information
        8. Preserve all links and references

        Return ONLY the enhanced markdown content without any explanations or additional text."""

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Please enhance this markdown content:\n\n{markdown_content}"}
            ],
            temperature=0.3,
            max_tokens=8000
        )

        enhanced_content = response.choices[0].message.content
        return enhanced_content, None

    except Exception as e:
        return markdown_content, str(e)


def convert_file_to_markdown(file_data, filename, enhance=True, api_key=None):
    """
    Convert a file to Markdown using MarkItDown and optionally enhance with ChatGPT-4o.
    Now with special handling for PowerPoint and PDF files to extract hyperlinks.

    Args:
        file_data: The binary content of the file
        filename: The name of the file
        enhance: Whether to enhance the markdown with ChatGPT-4o
        api_key: OpenAI API key (optional)

    Returns:
        Tuple of (markdown_content, error_message)
    """
    try:
        # Create a temporary file
        ext = get_file_extension(filename)
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{ext}") as tmp_file:
            tmp_file.write(file_data)
            tmp_file_path = tmp_file.name

        # Initialize MarkItDown and convert the file
        md = MarkItDown(enable_plugins=False)
        result = md.convert(tmp_file_path)
        markdown_content = result.text_content

        # Special handling for PowerPoint files
        if ext.lower() in ["pptx", "ppt"]:
            # Extract hyperlinks from PowerPoint
            hyperlinks = extract_pptx_hyperlinks(tmp_file_path)
            # Add hyperlinks to the markdown content
            markdown_content += format_hyperlinks_section(hyperlinks, "Presentation")

        # Special handling for PDF files
        elif ext.lower() == "pdf":
            try:
                # Extract hyperlinks from PDF
                hyperlinks = extract_pdf_hyperlinks(tmp_file_path)
                # Add hyperlinks to the markdown content
                markdown_content += format_hyperlinks_section(hyperlinks, "Document")
            except Exception as e:
                st.warning(f"PDF hyperlink extraction error: {str(e)}. Continuing with basic conversion.")

        # Clean up the temporary file
        os.unlink(tmp_file_path)

        # Enhance with ChatGPT-4o if enabled
        if enhance and api_key:
            enhanced_content, enhance_error = enhance_markdown_formatting(markdown_content, api_key)
            if enhance_error:
                st.warning(f"Enhancement error: {enhance_error}. Using original markdown.")
            else:
                markdown_content = enhanced_content

        return markdown_content, None
    except Exception as e:
        return "", str(e)


def convert_url_to_markdown(url, enhance=True, api_key=None):
    """
    Convert a website URL to Markdown by fetching HTML and using MarkItDown.

    Args:
        url: The website URL to convert
        enhance: Whether to enhance the markdown with ChatGPT-4o
        api_key: OpenAI API key (optional)

    Returns:
        Tuple of (markdown_content, error_message)
    """
    try:
        # Validate URL format
        parsed_url = urlparse(url)
        if not parsed_url.scheme or not parsed_url.netloc:
            return "", "Invalid URL format. Please include http:// or https://"

        # Add scheme if missing
        if not parsed_url.scheme:
            url = f"https://{url}"

        # Fetch the HTML content
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()  # Raise exception for 4XX/5XX status codes

        # Save HTML to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".html") as tmp_file:
            tmp_file.write(response.content)
            tmp_file_path = tmp_file.name

        # Initialize MarkItDown and convert the HTML file
        md = MarkItDown(enable_plugins=False)
        result = md.convert(tmp_file_path)

        # Clean up the temporary file
        os.unlink(tmp_file_path)

        markdown_content = result.text_content

        # Get page title for filename suggestion
        try:
            soup = BeautifulSoup(response.content, 'html.parser')
            page_title = soup.title.string if soup.title else parsed_url.netloc
            # Clean the title for use as a filename
            page_title = "".join(c if c.isalnum() or c in " -_" else "_" for c in page_title).strip()
            if len(page_title) > 50:
                page_title = page_title[:50]
            st.session_state.url_title = page_title
        except:
            st.session_state.url_title = parsed_url.netloc

        # Enhance with ChatGPT-4o if enabled
        if enhance and api_key:
            enhanced_content, enhance_error = enhance_markdown_formatting(markdown_content, api_key)
            if enhance_error:
                st.warning(f"Enhancement error: {enhance_error}. Using original markdown.")
            else:
                markdown_content = enhanced_content

        return markdown_content, None
    except requests.exceptions.RequestException as e:
        return "", f"Error fetching URL: {str(e)}"
    except Exception as e:
        return "", str(e)


def get_supported_formats():
    """
    Get a dictionary of supported file formats categorized by type.

    Returns:
        Dictionary of supported formats
    """
    return {
        "📝 Documents": {
            "formats": ["Word (.docx, .doc)", "PDF (with hyperlink extraction)", "EPub"],
            "extensions": ["docx", "doc", "pdf", "epub"],
        },
        "📊 Spreadsheets": {
            "formats": ["Excel (.xlsx, .xls)"],
            "extensions": ["xlsx", "xls"],
        },
        "📊 Presentations": {
            "formats": ["PowerPoint (.pptx, .ppt) with hyperlink extraction"],
            "extensions": ["pptx", "ppt"],
        },
        "🌐 Web": {"formats": ["HTML"], "extensions": ["html", "htm"]},
        "📁 Others": {
            "formats": ["CSV", "JSON", "XML", "ZIP (iterates over contents)"],
            "extensions": ["csv", "json", "xml", "zip"],
        },
    }


def process_folder(folder_path, output_folder=None, enhance=True, api_key=None):
    """
    Process all compatible files in a folder and convert them to markdown.

    Args:
        folder_path: Path to folder containing files to convert
        output_folder: Path to save markdown files (defaults to subfolder 'markdown' in input folder)
        enhance: Whether to enhance markdown with AI
        api_key: OpenAI API key

    Returns:
        Tuple of (success_count, error_count, errors_dict)
    """
    # Setup output folder
    if not output_folder:
        output_folder = os.path.join(folder_path, "markdown")

    # Create output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)

    # Get all compatible file types
    extensions = []
    formats = get_supported_formats()
    for category, info in formats.items():
        extensions.extend(info["extensions"])

    # Track results
    success_count = 0
    error_count = 0
    errors = {}

    # Get all compatible files
    files_to_process = []
    for ext in extensions:
        files_to_process.extend(glob.glob(os.path.join(folder_path, f"*.{ext}")))

    total_files = len(files_to_process)

    # Process each file
    for i, file_path in enumerate(files_to_process):
        file_name = os.path.basename(file_path)

        try:
            # Update progress
            progress = (i + 1) / total_files
            yield progress, f"Processing {file_name} ({i + 1}/{total_files})"

            # Read file content
            with open(file_path, 'rb') as file:
                file_data = file.read()

            # Convert to markdown
            markdown_content, error = convert_file_to_markdown(
                file_data,
                file_name,
                enhance=enhance,
                api_key=api_key
            )

            if error:
                error_count += 1
                errors[file_name] = error
                continue

            # Save markdown content
            output_file = os.path.join(output_folder, f"{os.path.splitext(file_name)[0]}.md")
            with open(output_file, 'w', encoding='utf-8') as md_file:
                md_file.write(markdown_content)

            success_count += 1

        except Exception as e:
            error_count += 1
            errors[file_name] = str(e)

    yield success_count, error_count, errors


def main():
    # Initialize session state
    if "markdown_content" not in st.session_state:
        st.session_state.markdown_content = ""
    if "file_name" not in st.session_state:
        st.session_state.file_name = ""
    if "url_title" not in st.session_state:
        st.session_state.url_title = ""
    if "folder_processing_results" not in st.session_state:
        st.session_state.folder_processing_results = None

    # App header
    st.title("Office to Markdown Converter")
    st.write("Convert your documents or websites to clean, structured Markdown with hyperlink extraction")

    # Sidebar
    with st.sidebar:
        st.header("Office to MD")
        st.write("Document Conversion Tool")

        # Supported formats in an expander
        with st.expander("Supported Formats"):
            formats = get_supported_formats()
            for category, info in formats.items():
                st.markdown(f"**{category}**")
                for format_name in info["formats"]:
                    st.markdown(f"- {format_name}")

            # Add website conversion info
            st.markdown("**🌐 Websites**")
            st.markdown("- Any URL (converts HTML to Markdown)")

        # Enhancement options
        st.subheader("Enhancement Options")
        enhance_markdown = st.checkbox("Enhance with AI", value=True,
                                       help="Use ChatGPT-4o to improve markdown formatting")

        # API key input
        openai_api_key = None
        if enhance_markdown:
            openai_api_key = st.text_input(
                "OpenAI API Key",
                type="password",
                help="Enter your OpenAI API key for enhancement"
            )
            if openai_api_key:
                # Store in environment for the session
                os.environ["OPENAI_API_KEY"] = openai_api_key

        # Developer info
        st.sidebar.markdown("---")
        st.sidebar.markdown("""
        **Developed by:** James Taylor
        """)

    # Main content area - Tabs
    tab1, tab2, tab3 = st.tabs(["File Upload", "Website URL", "Folder Processing"])

    with tab1:
        # File upload
        all_extensions = []
        formats = get_supported_formats()
        for category, info in formats.items():
            all_extensions.extend(info["extensions"])

        st.info(
            "Special feature: PDF and PowerPoint files will have their hyperlinks extracted and included in the markdown output.")

        uploaded_file = st.file_uploader(
            "Select a file to convert",
            type=all_extensions,
            help="Choose a file to convert to Markdown"
        )

        # Convert button for file
        if st.button("Convert File to Markdown", key="convert_file"):
            if not uploaded_file:
                st.error("Please upload a file to convert")
            else:
                with st.spinner("Converting to Markdown..."):
                    # Show progress bar
                    progress_bar = st.progress(0)
                    for percent_complete in range(100):
                        progress_bar.progress(percent_complete + 1)

                    # Convert uploaded file
                    markdown_content, error = convert_file_to_markdown(
                        uploaded_file.getbuffer(),
                        uploaded_file.name,
                        enhance=enhance_markdown,
                        api_key=openai_api_key
                    )

                    if error:
                        st.error(f"Error during conversion: {error}")
                    else:
                        st.session_state.markdown_content = markdown_content
                        st.session_state.file_name = uploaded_file.name
                        st.success("Conversion completed successfully!")

    with tab2:
        # Website URL input
        website_url = st.text_input(
            "Enter website URL",
            placeholder="https://example.com",
            help="Enter a complete URL to convert its content to Markdown"
        )

        # Convert button for URL
        if st.button("Convert Website to Markdown", key="convert_url"):
            if not website_url:
                st.error("Please enter a website URL to convert")
            else:
                with st.spinner("Fetching website and converting to Markdown..."):
                    # Show progress bar
                    progress_bar = st.progress(0)
                    for percent_complete in range(100):
                        progress_bar.progress(percent_complete + 1)

                    # Convert website URL
                    markdown_content, error = convert_url_to_markdown(
                        website_url,
                        enhance=enhance_markdown,
                        api_key=openai_api_key
                    )

                    if error:
                        st.error(f"Error during conversion: {error}")
                    else:
                        st.session_state.markdown_content = markdown_content
                        st.session_state.file_name = f"{st.session_state.url_title}.md"
                        st.success("Website conversion completed successfully!")

    with tab3:
        # Folder processing
        st.info(
            "Convert all supported files in a folder to markdown format. Each file will be processed and a .md file will be created.")

        # Input folder selection
        input_folder = st.text_input(
            "Enter path to folder with files to convert",
            placeholder="C:/Documents/MyFiles",
            help="Full path to the folder containing files to convert"
        )

        # Output folder selection (optional)
        output_folder = st.text_input(
            "Enter path to save markdown files (optional)",
            placeholder="Leave empty to create 'markdown' subfolder",
            help="Full path to save the converted markdown files. If left empty, files will be saved in a 'markdown' subfolder."
        )

        # Process folder button
        if st.button("Process Folder", key="process_folder"):
            if not input_folder or not os.path.isdir(input_folder):
                st.error("Please enter a valid folder path")
            else:
                # Process the folder
                progress_bar = st.progress(0)
                status_text = st.empty()

                try:
                    # Create a generator to process files and update progress
                    folder_processor = process_folder(
                        input_folder,
                        output_folder,
                        enhance=enhance_markdown,
                        api_key=openai_api_key
                    )

                    # Process files with progress updates
                    for progress, status in folder_processor:
                        progress_bar.progress(min(1.0, progress))
                        status_text.text(status)

                    # Get final results
                    success_count, error_count, errors = next(folder_processor)

                    # Save results to session state
                    st.session_state.folder_processing_results = {
                        "success_count": success_count,
                        "error_count": error_count,
                        "errors": errors,
                        "output_folder": output_folder if output_folder else os.path.join(input_folder, "markdown")
                    }

                    # Show success message
                    if success_count > 0:
                        st.success(f"Successfully converted {success_count} files.")

                    # Show error message if any
                    if error_count > 0:
                        st.warning(f"Failed to convert {error_count} files. See details below.")

                except Exception as e:
                    st.error(f"Error processing folder: {str(e)}")

        # Display folder processing results if available
        if st.session_state.folder_processing_results:
            results = st.session_state.folder_processing_results

            st.subheader("Folder Processing Results")
            st.markdown(f"**Output Location:** {results['output_folder']}")
            st.markdown(f"**Successfully Converted:** {results['success_count']} files")
            st.markdown(f"**Failed Conversions:** {results['error_count']} files")

            # Show errors if any
            if results['error_count'] > 0:
                with st.expander("View Conversion Errors"):
                    for file_name, error in results['errors'].items():
                        st.markdown(f"**{file_name}**: {error}")

            # Option to open output folder
            if st.button("Open Output Folder in File Explorer"):
                try:
                    # Try to open the folder based on platform
                    output_dir = results['output_folder']
                    if os.path.exists(output_dir):
                        if os.name == 'nt':  # Windows
                            os.startfile(output_dir)
                        elif os.name == 'posix':  # macOS, Linux
                            import subprocess
                            if os.path.exists('/usr/bin/open'):  # macOS
                                subprocess.call(['open', output_dir])
                            else:  # Linux
                                subprocess.call(['xdg-open', output_dir])
                        st.success(f"Opened folder: {output_dir}")
                    else:
                        st.error("Output folder not found")
                except Exception as e:
                    st.error(f"Failed to open folder: {str(e)}")

    # Display the converted markdown content (for single file or URL conversion)
    if st.session_state.markdown_content:
        st.subheader("Converted Markdown")
        st.text_area(
            "Markdown Content",
            value=st.session_state.markdown_content,
            height=400
        )

        # Download button
        st.download_button(
            label="Download Markdown",
            data=st.session_state.markdown_content,
            file_name=st.session_state.file_name.rsplit(".", 1)[0] + ".md"
            if "." in st.session_state.file_name else st.session_state.file_name + ".md",
            mime="text/markdown"
        )


if __name__ == "__main__":
    main()