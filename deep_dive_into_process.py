from pptx import Presentation

def extract_bullets_with_fallback(file_path):
    prs = Presentation(file_path)

    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"\n--- Slide {slide_idx} ---")
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                level = para.level

                # Skip empty lines
                if not text:
                    continue

                # Check real PowerPoint bullets (buChar or buAutoNum XML)
                is_real_bullet = (
                    para._element.xpath('./a:pPr/a:buChar') or
                    para._element.xpath('./a:pPr/a:buAutoNum')
                )

                # Check if it's a "visual" bullet typed manually
                looks_like_bullet = text.startswith(("•", "-", "→", "*"))

                if is_real_bullet or looks_like_bullet:
                    indent = "  " * level
                    print(f"{indent}- {text}")


extract_bullets_with_fallback("testing_powerpoint.pptx")