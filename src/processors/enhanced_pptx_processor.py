"""
PowerPoint Processor - Complete with Accessibility Reading Order Support
Maintains all original functionality while adding accessibility reading order extraction
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER
import json
import re
import os
from datetime import datetime
import xml.etree.ElementTree as ET


class PowerPointProcessor:
   """Complete PowerPoint processing with accessibility reading order and fixed bullet detection"""

   def __init__(self):
       self.supported_formats = ['.pptx', '.ppt']

       # Accessibility support - NEW
       self.namespaces = {
           'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
           'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
           'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
       }
       self.use_accessibility_order = True  # Set to False to use original positional order

   def convert_pptx_to_markdown_enhanced(self, file_path, convert_slide_titles=True):
       """
       Main entry point: v14 text extraction + v19 diagram detection + accessibility order

       Args:
           file_path (str): Path to the PowerPoint file
           convert_slide_titles (bool): Whether to convert slide titles from bullets to H1 headings
       """
       try:
           prs = Presentation(file_path)

           # Extract PowerPoint metadata first
           pptx_metadata = self.extract_pptx_metadata(prs, file_path)

           # Extract structured data with accessibility order
           structured_data = self.extract_presentation_data(prs)

           # Convert to basic markdown
           markdown = self.convert_structured_data_to_markdown(structured_data, convert_slide_titles)

           # Add PowerPoint metadata as comments for Claude to use
           markdown_with_metadata = self.add_pptx_metadata_for_claude(markdown, pptx_metadata)

           # APPEND v19 diagram analysis at the end
           diagram_analysis = self.analyze_structured_data_for_diagrams(structured_data)
           if diagram_analysis:
               markdown_with_metadata += "\n\n" + diagram_analysis

           return markdown_with_metadata
       except Exception as e:
           raise Exception(f"Error processing PowerPoint file: {str(e)}")

   def extract_presentation_data(self, presentation):
       """Extract all content with minimal processing"""
       data = {
           "total_slides": len(presentation.slides),
           "slides": []
       }

       for slide_idx, slide in enumerate(presentation.slides, 1):
           slide_data = self.extract_slide_data(slide, slide_idx)
           data["slides"].append(slide_data)

       return data

   def extract_slide_data(self, slide, slide_number):
       """
       ENHANCED: Extract slide content using true accessibility reading order with semantic hierarchy
       """
       if self.use_accessibility_order:
           try:
               # Try semantic accessibility order extraction
               ordered_shapes = self._get_semantic_accessibility_order(slide)
               extraction_method = "semantic_accessibility_order"
           except Exception as e:
               print(f"Accessibility extraction failed for slide {slide_number}, using document order: {e}")
               # Fall back to XML document order (NOT positional)
               ordered_shapes = self._get_xml_document_order(slide)
               extraction_method = "xml_document_order"
       else:
           # Use original positional method only if specifically disabled
           ordered_shapes = self._get_positional_ordered_shapes(slide)
           extraction_method = "positional_order"

       slide_data = {
           "slide_number": slide_number,
           "content_blocks": []
       }

       # Use existing extract_shape_content method (unchanged!)
       for shape in ordered_shapes:
           block = self.extract_shape_content(shape)
           if block:
               slide_data["content_blocks"].append(block)

       return slide_data

   def _get_positional_ordered_shapes(self, slide):
       """
       Original positional ordering method - preserved exactly
       """
       positioned_shapes = []
       for shape in slide.shapes:
           if hasattr(shape, 'top') and hasattr(shape, 'left'):
               positioned_shapes.append((shape.top, shape.left, shape))
           else:
               positioned_shapes.append((0, 0, shape))

       positioned_shapes.sort(key=lambda x: (x[0], x[1]))
       return [shape for _, _, shape in positioned_shapes]

   def _get_semantic_accessibility_order(self, slide):
       """
       NEW: Get shapes in semantic accessibility order - titles first, then document order
       """
       # Get all shapes in XML document order (true reading order)
       xml_ordered_shapes = self._get_xml_document_order(slide)

       # Separate shapes by semantic importance
       title_shapes = []
       subtitle_shapes = []
       content_shapes = []
       other_shapes = []

       for shape in xml_ordered_shapes:
           semantic_role = self._get_shape_semantic_role(shape)

           if semantic_role == "title":
               title_shapes.append(shape)
           elif semantic_role == "subtitle":
               subtitle_shapes.append(shape)
           elif semantic_role == "content":
               content_shapes.append(shape)
           else:
               other_shapes.append(shape)

       # Return in semantic hierarchy order: titles first, then content
       return title_shapes + subtitle_shapes + content_shapes + other_shapes

   def _get_xml_document_order(self, slide):
       """
       SIMPLIFIED: Get shapes in pure XML document order with better group handling
       """
       try:
           # Get the slide's XML
           slide_xml = self._get_slide_xml(slide)

           # Parse XML to get ALL shapes in document order
           xml_shape_info = self._parse_slide_xml_for_document_order(slide_xml)

           # SIMPLIFIED mapping - try to preserve exact order
           ordered_shapes = self._map_xml_to_pptx_shapes_simple(xml_shape_info, slide.shapes)

           return ordered_shapes
       except Exception as e:
           print(f"XML document order failed: {e}, falling back to slide.shapes order")
           # Ultimate fallback: use slide.shapes in their natural order (which is often document order)
           return list(slide.shapes)

   def _map_xml_to_pptx_shapes_simple(self, xml_shape_info, pptx_shapes):
       """
       SIMPLIFIED mapping that tries harder to preserve XML order
       """
       ordered_shapes = []
       used_shapes = []

       # Convert to list for easier manipulation
       available_shapes = list(pptx_shapes)

       for xml_info in xml_shape_info:
           best_match = None
           best_match_score = 0

           # Try to find the best match for this XML shape
           for shape in available_shapes:
               if shape in used_shapes:
                   continue

               match_score = 0

               # ID matching (highest priority)
               if xml_info.get('id') and xml_info['id'] == self._get_shape_id(shape):
                   match_score += 100

               # Name matching
               if xml_info.get('name') and xml_info['name'] == self._get_shape_name(shape):
                   match_score += 50

               # Type matching
               xml_type = xml_info.get('type', '')
               shape_type = str(shape.shape_type).split('.')[-1]

               if xml_type == 'grpSp' and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                   match_score += 30
                   # Additional group matching
                   xml_child_count = xml_info.get('group_child_count', 0)
                   shape_child_count = len(list(shape.shapes)) if hasattr(shape, 'shapes') else 0
                   if xml_child_count > 0 and xml_child_count == shape_child_count:
                       match_score += 20
               elif xml_type == 'sp' and shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                   match_score += 20
               elif xml_type == 'pic' and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                   match_score += 20

               # Text content matching
               if xml_info.get('text_content') and xml_info['text_content']:
                   shape_text = self._get_shape_text_preview(shape)
                   if shape_text and xml_info['text_content'] in shape_text:
                       match_score += 10

               # Track the best match
               if match_score > best_match_score:
                   best_match = shape
                   best_match_score = match_score

           # Use the best match we found
           if best_match:
               ordered_shapes.append(best_match)
               used_shapes.append(best_match)
           else:
               # No good match found - this might indicate an issue
               print(f"Warning: No match found for XML shape {xml_info.get('type', 'unknown')} with ID {xml_info.get('id', 'None')}")

       # Add any remaining shapes that weren't matched (preserve their relative order)
       for shape in available_shapes:
           if shape not in used_shapes:
               ordered_shapes.append(shape)
               print(f"Adding unmatched shape: {str(shape.shape_type).split('.')[-1]}")

       return ordered_shapes

   def _extract_shape_info_from_xml(self, shape_elem, order_index):
       """
       NEW: Extract identifying information from a shape's XML element
       Enhanced to better handle groups
       """
       shape_info = {
           'xml_order': order_index,
           'id': None,
           'name': None,
           'type': shape_elem.tag.split('}')[-1] if '}' in shape_elem.tag else shape_elem.tag,
           'has_text': False,
           'text_content': None,
           'is_group': False,
           'group_child_count': 0
       }

       # Get shape ID and name from non-visual properties - ENHANCED for groups
       nv_props = (shape_elem.find('.//p:nvSpPr', self.namespaces) or
                   shape_elem.find('.//p:nvPicPr', self.namespaces) or
                   shape_elem.find('.//p:nvGraphicFramePr', self.namespaces) or
                   shape_elem.find('.//p:nvGrpSpPr', self.namespaces) or
                   shape_elem.find('.//p:nvCxnSpPr', self.namespaces))

       if nv_props is not None:
           cnv_pr = nv_props.find('.//p:cNvPr', self.namespaces)
           if cnv_pr is not None:
               shape_info['id'] = cnv_pr.get('id')
               shape_info['name'] = cnv_pr.get('name', '')

       # Special handling for groups
       if shape_info['type'] == 'grpSp':
           shape_info['is_group'] = True
           # Count child shapes in group for better matching
           child_shapes = shape_elem.findall('.//p:sp', self.namespaces)
           child_shapes.extend(shape_elem.findall('.//p:pic', self.namespaces))
           child_shapes.extend(shape_elem.findall('.//p:graphicFrame', self.namespaces))
           shape_info['group_child_count'] = len(child_shapes)

           # For groups, also collect text from all children for matching
           group_texts = []
           for child in child_shapes[:3]:  # First 3 children for matching
               child_text_elements = child.findall('.//a:t', self.namespaces)
               for t_elem in child_text_elements:
                   if t_elem.text and t_elem.text.strip():
                       group_texts.append(t_elem.text.strip()[:20])

           if group_texts:
               shape_info['has_text'] = True
               shape_info['text_content'] = ' | '.join(group_texts)[:50]

       # Check if shape has text (for non-groups or as fallback)
       if not shape_info['has_text']:
           text_elements = shape_elem.findall('.//a:t', self.namespaces)
           if text_elements:
               all_text = ' '.join([t.text for t in text_elements if t.text])
               if all_text.strip():
                   shape_info['has_text'] = True
                   shape_info['text_content'] = all_text.strip()[:50]

       return shape_info

   def _get_shape_semantic_role(self, shape):
       """
       NEW: Determine the semantic role of a shape (title, subtitle, content, other)
       """
       # Check if it's a placeholder with semantic meaning
       try:
           if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
               ph_type = shape.placeholder_format.type

               # PowerPoint placeholder types that indicate semantic roles
               if hasattr(ph_type, 'name'):
                   ph_name = ph_type.name.upper()
                   if 'TITLE' in ph_name and 'SUBTITLE' not in ph_name:
                       return "title"
                   elif 'SUBTITLE' in ph_name:
                       return "subtitle"
                   elif any(content_type in ph_name for content_type in ['BODY', 'CONTENT', 'TEXT']):
                       return "content"
       except:
           pass

       # Check shape name for semantic indicators
       try:
           if hasattr(shape, 'name') and shape.name:
               name_lower = shape.name.lower()
               if any(title_word in name_lower for title_word in ['title', 'heading', 'header']):
                   return "title"
               elif 'subtitle' in name_lower:
                   return "subtitle"
       except:
           pass

       # Analyze text content for title characteristics
       try:
           text_content = self._get_shape_text_preview(shape)
           if text_content and self._is_title_like_text(text_content):
               return "title"
       except:
           pass

       # Default to content for text shapes, other for non-text
       if hasattr(shape, 'text_frame') or hasattr(shape, 'text'):
           return "content"
       else:
           return "other"

   def _is_title_like_text(self, text):
       """
       NEW: Analyze if text content looks like a title
       """
       if not text or len(text) > 100:
           return False

       # Title characteristics
       title_indicators = [
           len(text.split()) <= 10,  # Short phrases
           not text.endswith(('.', '!', '?')),  # No ending punctuation
           text.isupper() or text.istitle(),  # Proper capitalization
           len(text) < 80,  # Reasonable length
       ]

       # Must meet most criteria
       return sum(title_indicators) >= 3

   def _parse_slide_xml_for_document_order(self, slide_xml):
       """
       NEW: Parse slide XML to extract shapes in pure document order (not filtered)
       """
       # Parse the XML
       root = ET.fromstring(slide_xml)

       # Find the shape tree (spTree) which contains shapes in document order
       shape_tree = root.find('.//p:spTree', self.namespaces)
       if shape_tree is None:
           raise Exception("No shape tree found in slide XML")

       shape_order_info = []

       # Iterate through ALL children of spTree in exact document order
       for idx, elem in enumerate(shape_tree):
           # Check if this is any kind of shape element
           tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag

           # Include ALL shape types in document order
           if tag_name in ['sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp', 'AlternateContent']:
               shape_info = self._extract_shape_info_from_xml(elem, idx)
               if shape_info:
                   shape_order_info.append(shape_info)

       return shape_order_info

   def _map_xml_to_pptx_shapes_strict(self, xml_shape_info, pptx_shapes):
       """
       NEW: Strict mapping that preserves XML document order exactly
       Enhanced to properly handle groups in reading order
       """
       ordered_shapes = []
       used_shape_ids = set()

       # Create lookup for python-pptx shapes with their IDs
       pptx_shape_lookup = {}
       shape_id_to_shape = {}

       # Separate groups and non-groups for better matching
       pptx_groups = []
       pptx_non_groups = []

       for shape in pptx_shapes:
           if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
               pptx_groups.append(shape)
           else:
               pptx_non_groups.append(shape)

           # Get multiple identifiers for matching
           shape_id = self._get_shape_id(shape)
           shape_name = self._get_shape_name(shape)
           shape_text = self._get_shape_text_preview(shape)

           # Store shape by ID for tracking
           if shape_id:
               shape_id_to_shape[shape_id] = shape

           # Store with multiple keys for flexible matching
           if shape_id:
               pptx_shape_lookup[f"id_{shape_id}"] = shape
           if shape_name:
               pptx_shape_lookup[f"name_{shape_name}"] = shape
           if shape_text:
               pptx_shape_lookup[f"text_{shape_text}"] = shape

           # For groups, add special matching keys
           if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
               group_child_count = len(list(shape.shapes))
               pptx_shape_lookup[f"group_count_{group_child_count}"] = shape

               # Create a signature from group's text content
               group_text_signature = self._get_group_text_signature(shape)
               if group_text_signature:
                   pptx_shape_lookup[f"group_text_{group_text_signature}"] = shape

       # Match XML order to python-pptx shapes in EXACT order
       for xml_info in xml_shape_info:
           matched_shape = None

           # Enhanced group matching
           if xml_info.get('is_group', False):
               # Try matching groups by ID first
               if xml_info['id']:
                   matched_shape = pptx_shape_lookup.get(f"id_{xml_info['id']}")

               # Try matching by name
               if not matched_shape and xml_info['name']:
                   matched_shape = pptx_shape_lookup.get(f"name_{xml_info['name']}")

               # Try matching by child count
               if not matched_shape and xml_info.get('group_child_count', 0) > 0:
                   matched_shape = pptx_shape_lookup.get(f"group_count_{xml_info['group_child_count']}")

               # Try matching by text signature
               if not matched_shape and xml_info.get('text_content'):
                   matched_shape = pptx_shape_lookup.get(f"group_text_{xml_info['text_content']}")

               # Fallback: find first unused group
               if not matched_shape:
                   for group in pptx_groups:
                       group_id = self._get_shape_id(group)
                       if not group_id or group_id not in used_shape_ids:
                           if group not in ordered_shapes:  # Make sure not already added
                               matched_shape = group
                               break
           else:
               # Regular shape matching
               # Try matching by ID first (most reliable)
               if xml_info['id']:
                   matched_shape = pptx_shape_lookup.get(f"id_{xml_info['id']}")

               # Try matching by name
               if not matched_shape and xml_info['name']:
                   matched_shape = pptx_shape_lookup.get(f"name_{xml_info['name']}")

               # Try matching by text content
               if not matched_shape and xml_info['text_content']:
                   matched_shape = pptx_shape_lookup.get(f"text_{xml_info['text_content']}")

           # If we found a match and haven't used it yet
           if matched_shape:
               matched_shape_id = self._get_shape_id(matched_shape)

               # Check if already used
               already_used = False
               if matched_shape_id and matched_shape_id in used_shape_ids:
                   already_used = True
               elif matched_shape in ordered_shapes:  # Fallback check for shapes without IDs
                   already_used = True

               if not already_used:
                   ordered_shapes.append(matched_shape)
                   if matched_shape_id:
                       used_shape_ids.add(matched_shape_id)

                   # Remove from lookup to avoid duplicates
                   for key in list(pptx_shape_lookup.keys()):
                       if pptx_shape_lookup[key] is matched_shape:
                           del pptx_shape_lookup[key]

       # Add any remaining shapes that weren't matched
       for shape in pptx_shapes:
           shape_id = self._get_shape_id(shape)

           # Check if this shape is already included
           already_included = False
           if shape_id and shape_id in used_shape_ids:
               already_included = True
           elif not shape_id and shape in ordered_shapes:
               already_included = True

           if not already_included:
               ordered_shapes.append(shape)

       return ordered_shapes

   def _get_group_text_signature(self, group_shape):
       """
       NEW: Create a text signature for a group based on its child text content
       """
       try:
           group_texts = []
           for child_shape in list(group_shape.shapes)[:3]:  # First 3 children
               child_text = self._get_shape_text_preview(child_shape)
               if child_text:
                   group_texts.append(child_text[:20])

           if group_texts:
               return ' | '.join(group_texts)[:50]
       except:
           pass
       return None

   def _get_slide_xml(self, slide):
       """
       NEW: Extract the raw XML from a slide
       """
       if hasattr(slide, '_element') and hasattr(slide._element, 'xml'):
           return slide._element.xml
       elif hasattr(slide, 'element') and hasattr(slide.element, 'xml'):
           return slide.element.xml
       else:
           # Try alternative access methods
           slide_part = slide.part if hasattr(slide, 'part') else None
           if slide_part and hasattr(slide_part, '_element'):
               return slide_part._element.xml
           else:
               raise Exception("Cannot access slide XML")

   def _map_xml_to_pptx_shapes(self, xml_shape_info, pptx_shapes):
       """
       DEPRECATED: Use _map_xml_to_pptx_shapes_strict instead for true accessibility order
       This method is kept for backward compatibility only
       """
       return self._map_xml_to_pptx_shapes_strict(xml_shape_info, pptx_shapes)

   def _get_shape_id(self, shape):
       """
       NEW: Extract shape ID from python-pptx shape object
       """
       try:
           if hasattr(shape, '_element') and hasattr(shape._element, 'xml'):
               xml_str = shape._element.xml
               # Look for id attribute in cNvPr element
               match = re.search(r'<[^>]*:cNvPr[^>]+id="([^"]+)"', xml_str)
               if match:
                   return match.group(1)
       except:
           pass
       return None

   def _get_shape_name(self, shape):
       """
       NEW: Extract shape name from python-pptx shape object
       """
       try:
           if hasattr(shape, 'name') and shape.name:
               return shape.name
           elif hasattr(shape, '_element') and hasattr(shape._element, 'xml'):
               xml_str = shape._element.xml
               # Look for name attribute in cNvPr element
               match = re.search(r'<[^>]*:cNvPr[^>]+name="([^"]+)"', xml_str)
               if match:
                   return match.group(1)
       except:
           pass
       return None

   def _get_shape_text_preview(self, shape):
       """
       NEW: Get preview of shape text for matching
       """
       try:
           if hasattr(shape, 'text_frame') and shape.text_frame:
               text = ""
               for para in shape.text_frame.paragraphs:
                   if para.text.strip():
                       text += para.text.strip() + " "
               if text.strip():
                   return text.strip()[:50]
           elif hasattr(shape, 'text') and shape.text:
               return shape.text.strip()[:50]
       except:
           pass
       return None

   def debug_accessibility_order(self, file_path, slide_number=1):
       """
       ENHANCED: Debug method with detailed group tracking
       """
       prs = Presentation(file_path)
       if slide_number > len(prs.slides):
           print(f"Slide {slide_number} not found. Presentation has {len(prs.slides)} slides.")
           return

       slide = prs.slides[slide_number - 1]

       print(f"\n=== DEBUGGING SLIDE {slide_number} READING ORDER ===")

       # First, let's see what XML actually contains
       print("\n🔍 XML ANALYSIS:")
       try:
           slide_xml = self._get_slide_xml(slide)
           xml_shape_info = self._parse_slide_xml_for_document_order(slide_xml)

           print(f"XML found {len(xml_shape_info)} shapes:")
           for i, xml_info in enumerate(xml_shape_info):
               shape_type = xml_info.get('type', 'unknown')
               is_group = xml_info.get('is_group', False)
               text_content = xml_info.get('text_content', 'No text')
               shape_id = xml_info.get('id', 'No ID')
               name = xml_info.get('name', 'No name')

               group_info = ""
               if is_group:
                   child_count = xml_info.get('group_child_count', 0)
                   group_info = f" [GROUP-{child_count}]"

               print(f"  XML {i+1}: {shape_type}{group_info} | ID:{shape_id} | Name:{name} | Text:{text_content[:30]}...")
       except Exception as e:
           print(f"XML Analysis failed: {e}")

       # Now let's see what python-pptx shapes we have
       print(f"\n🐍 PYTHON-PPTX SHAPES ({len(slide.shapes)} total):")
       for i, shape in enumerate(slide.shapes):
           shape_type = str(shape.shape_type).split('.')[-1]
           shape_id = self._get_shape_id(shape)
           shape_name = self._get_shape_name(shape)
           text_preview = self._get_shape_text_preview(shape)

           group_info = ""
           if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
               child_count = len(list(shape.shapes))
               group_info = f" [GROUP-{child_count}]"

           print(f"  PPTX {i+1}: {shape_type}{group_info} | ID:{shape_id} | Name:{shape_name} | Text:{text_preview[:30] if text_preview else 'No text'}...")

       # Test the mapping
       print(f"\n🔗 MAPPING TEST:")
       try:
           xml_shape_info = self._parse_slide_xml_for_document_order(slide_xml)
           mapped_shapes = self._map_xml_to_pptx_shapes_strict(xml_shape_info, slide.shapes)

           print(f"Successfully mapped {len(mapped_shapes)} shapes:")
           for i, shape in enumerate(mapped_shapes):
               shape_type = str(shape.shape_type).split('.')[-1]
               text_preview = self._get_shape_text_preview(shape)

               group_info = ""
               if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                   child_count = len(list(shape.shapes))
                   group_info = f" [GROUP-{child_count}]"

               print(f"  Mapped {i+1}: {shape_type}{group_info} | Text:{text_preview[:30] if text_preview else 'No text'}...")
       except Exception as e:
           print(f"Mapping test failed: {e}")

       # Get semantic accessibility order
       try:
           semantic_ordered = self._get_semantic_accessibility_order(slide)
           semantic_success = True
       except Exception as e:
           print(f"Semantic accessibility extraction failed: {e}")
           semantic_ordered = []
           semantic_success = False

       # Get positional order for comparison
       positional_ordered = self._get_positional_ordered_shapes(slide)

       print(f"\n📊 FINAL RESULTS:")
       print(f"Semantic Success: {semantic_success}")
       print(f"Semantic Order Count: {len(semantic_ordered)}")
       print(f"Positional Order Count: {len(positional_ordered)}")

       # Compare the orders
       print("\n🎯 SEMANTIC ACCESSIBILITY ORDER (titles first):")
       for i, shape in enumerate(semantic_ordered):
           text_preview = self._get_shape_text_preview(shape) or "No text"
           shape_type = str(shape.shape_type).split('.')[-1] if hasattr(shape, 'shape_type') else "unknown"
           semantic_role = self._get_shape_semantic_role(shape)

           # Special handling for groups
           if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
               child_count = len(list(shape.shapes))
               group_text = self._get_group_text_signature(shape) or text_preview
               print(f"  {i+1}. [{semantic_role.upper()}] [GROUP-{child_count}] {group_text[:40]}...")
           else:
               print(f"  {i+1}. [{semantic_role.upper()}] [{shape_type}] {text_preview[:40]}...")

       print("\n📍 POSITIONAL ORDER (old method):")
       for i, shape in enumerate(positional_ordered):
           text_preview = self._get_shape_text_preview(shape) or "No text"
           shape_type = str(shape.shape_type).split('.')[-1] if hasattr(shape, 'shape_type') else "unknown"

           # Special handling for groups
           if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
               child_count = len(list(shape.shapes))
               group_text = self._get_group_text_signature(shape) or text_preview
               print(f"  {i+1}. [GROUP-{child_count}] {group_text[:40]}...")
           else:
               print(f"  {i+1}. [{shape_type}] {text_preview[:40]}...")

       # Show where groups ended up
       groups_in_semantic = [(i, s) for i, s in enumerate(semantic_ordered) if s.shape_type == MSO_SHAPE_TYPE.GROUP]
       groups_in_positional = [(i, s) for i, s in enumerate(positional_ordered) if s.shape_type == MSO_SHAPE_TYPE.GROUP]

       if groups_in_semantic or groups_in_positional:
           print(f"\n📦 GROUP POSITION ANALYSIS:")
           print(f"Groups in semantic order: {[i+1 for i, s in groups_in_semantic]}")
           print(f"Groups in positional order: {[i+1 for i, s in groups_in_positional]}")

           if groups_in_semantic != groups_in_positional:
               print("⚠️  GROUP POSITIONS DIFFER! This indicates the accessibility order is working.")
           else:
               print("❌ GROUP POSITIONS SAME - accessibility order may not be working for groups.")


   def extract_shape_content(self, shape):
       """Extract shape content with proper type detection - v14 approach BUT capture shape info for diagram analysis"""
       # Capture basic shape info for later diagram analysis (with safe error handling)
       shape_info = {
           "shape_type": "unknown",
           "auto_shape_type": None,
           "position": {
               "top": getattr(shape, 'top', 0),
               "left": getattr(shape, 'left', 0),
               "width": getattr(shape, 'width', 0),
               "height": getattr(shape, 'height', 0)
           }
       }

       # Safely get shape type
       try:
           if hasattr(shape, 'shape_type'):
               shape_info["shape_type"] = str(shape.shape_type).split('.')[-1]  # Get just the name part
       except:
           shape_info["shape_type"] = "unknown"

       # Check for auto shape type (for arrows and special shapes)
       try:
           if hasattr(shape, 'auto_shape_type'):
               shape_info["auto_shape_type"] = str(shape.auto_shape_type).split('.')[-1]
       except:
           pass

       # MAIN EXTRACTION - v14 approach that works
       content_block = None

       try:
           if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
               content_block = self.extract_image(shape)
           elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
               content_block = self.extract_table(shape.table)
           elif hasattr(shape, 'has_chart') and shape.has_chart:
               content_block = self.extract_chart(shape)
           elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
               content_block = self.extract_group(shape)
           elif hasattr(shape, 'text_frame') and shape.text_frame:
               content_block = self.extract_text_frame_fixed(shape.text_frame, shape)
           elif hasattr(shape, 'text') and shape.text:
               content_block = self.extract_plain_text(shape)
       except Exception as e:
           print(f"Warning: Error extracting shape content: {e}")
           return None

       # DIAGRAM ANALYSIS - add shape info for diagram detection (with safe checks)
       if not content_block:
           # For shapes without text content, create minimal blocks for diagram analysis
           try:
               if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                   content_block = {"type": "line", "line_type": "simple"}
               elif shape.shape_type == MSO_SHAPE_TYPE.CONNECTOR:
                   content_block = {"type": "line", "line_type": "connector"}
               elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                   content_block = {"type": "line", "line_type": "freeform"}
               elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                   if self.is_arrow_shape(shape_info["auto_shape_type"]):
                       content_block = {"type": "arrow", "arrow_type": shape_info["auto_shape_type"]}
                   else:
                       content_block = {"type": "shape", "shape_subtype": "auto_shape"}
               else:
                   content_block = {"type": "shape", "shape_subtype": "generic"}
           except Exception as e:
               # Fallback for any shape type issues
               content_block = {"type": "shape", "shape_subtype": "unknown"}

       # Add shape analysis info to content block for diagram detection
       if content_block:
           try:
               content_block.update(shape_info)
           except Exception as e:
               print(f"Warning: Error adding shape info: {e}")

       return content_block

   def analyze_structured_data_for_diagrams(self, structured_data):
       """
       v19 diagram analysis system - analyze extracted structured data
       """
       try:
           diagram_slides = []

           for slide_idx, slide in enumerate(structured_data["slides"]):
               score_analysis = self.score_slide_for_diagram(slide)
               if score_analysis["probability"] >= 40:  # 40%+ probability threshold
                   diagram_slides.append({
                       "slide": slide_idx + 1,
                       "analysis": score_analysis
                   })

           # Generate detailed summary
           if diagram_slides:
               summary = "## DIAGRAM ANALYSIS (v19 Scoring System)\n\n"
               summary += "**Slides with potential diagrams:**\n\n"

               for slide_info in diagram_slides:
                   analysis = slide_info["analysis"]
                   summary += f"- **Slide {slide_info['slide']}**: {analysis['probability']}% probability "
                   summary += f"(Score: {analysis['total_score']}) - {', '.join(analysis['reasons'])}\n"
                   summary += f"  - Shapes: {analysis['shape_count']}, Lines: {analysis['line_count']}, Arrows: {analysis['arrow_count']}\n\n"

               return summary

           return None

       except Exception as e:
           return f"\n\n<!-- v19 Diagram analysis error: {e} -->"

   def score_slide_for_diagram(self, slide_data):
       """
       v19 scoring system: Score a slide for diagram probability using sophisticated rules
       """
       content_blocks = slide_data.get("content_blocks", [])

       # Collect all shapes and lines from structured data
       shapes = []
       lines = []
       arrows = []
       text_blocks = []

       for block in content_blocks:
           if block.get("type") == "line":
               lines.append(block)
           elif block.get("type") == "arrow":
               arrows.append(block)
           elif block.get("type") == "text":
               text_blocks.append(block)
               shapes.append(block)
           elif block.get("type") in ["shape", "image", "chart"]:
               shapes.append(block)
           elif block.get("type") == "group":
               # Recursively analyze group contents
               group_analysis = self._analyze_group_contents(block)
               shapes.extend(group_analysis["shapes"])
               lines.extend(group_analysis["lines"])
               arrows.extend(group_analysis["arrows"])
               text_blocks.extend(group_analysis["text_blocks"])

       # Calculate score based on v19 rules
       score = 0
       reasons = []

       # Rule 1: Line/Arrow threshold (20+ points each)
       if len(arrows) > 0:
           score += 20
           reasons.append(f"block_arrows:{len(arrows)}")

       if len(lines) >= 3:
           score += 20
           reasons.append(f"connector_lines:{len(lines)}")

       # Rule 2: Line-to-shape ratio (15 points)
       total_lines = len(lines) + len(arrows)
       if len(shapes) > 0:
           line_ratio = total_lines / len(shapes)
           if line_ratio >= 0.5:
               score += 15
               reasons.append(f"line_ratio:{line_ratio:.1f}")

       # Rule 3: Spatial layout analysis (10-15 points)
       layout_score = self._analyze_spatial_layout(shapes)
       score += layout_score["score"]
       if layout_score["score"] > 0:
           reasons.append(f"layout:{layout_score['type']}")

       # Rule 4: Shape variety (10-15 points)
       variety_score = self._analyze_shape_variety(shapes)
       score += variety_score
       if variety_score > 0:
           reasons.append(f"variety:{variety_score}")

       # Rule 5: Text density analysis (10 points)
       text_score = self._analyze_text_density(text_blocks)
       score += text_score
       if text_score > 0:
           reasons.append(f"short_text:{text_score}")

       # Rule 6: Flow patterns (20 points)
       flow_score = self._analyze_flow_patterns(shapes, lines, arrows, text_blocks)
       score += flow_score
       if flow_score > 0:
           reasons.append(f"flow_pattern:{flow_score}")

       # Negative indicators
       negative_score = self._analyze_negative_indicators(text_blocks, shapes)
       score += negative_score  # negative_score will be negative or 0
       if negative_score < 0:
           reasons.append(f"negatives:{negative_score}")

       # Convert score to probability
       if score >= 60:
           probability = 95
       elif score >= 40:
           probability = 75
       elif score >= 20:
           probability = 40
       else:
           probability = 10

       return {
           "total_score": score,
           "probability": probability,
           "reasons": reasons,
           "shape_count": len(shapes),
           "line_count": len(lines),
           "arrow_count": len(arrows)
       }

   def _analyze_group_contents(self, group_block):
       """Recursively analyze group contents for diagram elements"""
       result = {"shapes": [], "lines": [], "arrows": [], "text_blocks": []}

       for extracted_block in group_block.get("extracted_blocks", []):
           if extracted_block.get("type") == "line":
               result["lines"].append(extracted_block)
           elif extracted_block.get("type") == "arrow":
               result["arrows"].append(extracted_block)
           elif extracted_block.get("type") == "text":
               result["text_blocks"].append(extracted_block)
               result["shapes"].append(extracted_block)
           elif extracted_block.get("type") in ["shape", "image", "chart"]:
               result["shapes"].append(extracted_block)

       return result

   def _analyze_spatial_layout(self, shapes):
       """Analyze spatial layout patterns"""
       if len(shapes) < 3:
           return {"score": 0, "type": "insufficient"}

       positions = []
       for shape in shapes:
           pos = shape.get("position")
           if pos:
               positions.append((pos["top"], pos["left"]))

       if len(positions) < 3:
           return {"score": 0, "type": "no_position_data"}

       # Calculate spread
       tops = [p[0] for p in positions]
       lefts = [p[1] for p in positions]

       top_range = max(tops) - min(tops) if tops else 0
       left_range = max(lefts) - min(lefts) if lefts else 0

       # Check for grid-like arrangement
       unique_tops = len(set(round(t / 100000) for t in tops))  # Group by approximate position
       unique_lefts = len(set(round(l / 100000) for l in lefts))

       if unique_tops >= 2 and unique_lefts >= 2:
           return {"score": 15, "type": "grid_layout"}
       elif top_range > 1000000 and left_range > 1000000:
           return {"score": 10, "type": "spread_layout"}
       else:
           return {"score": 0, "type": "linear_layout"}

   def _analyze_shape_variety(self, shapes):
       """Analyze variety in shape types and sizes"""
       if len(shapes) < 2:
           return 0

       shape_types = set()
       sizes = []

       for shape in shapes:
           shape_types.add(shape.get("type", "unknown"))
           pos = shape.get("position")
           if pos:
               size = pos["width"] * pos["height"]
               sizes.append(size)

       score = 0

       # Multiple shape types
       if len(shape_types) >= 3:
           score += 15
       elif len(shape_types) >= 2:
           score += 10

       # Consistent sizing (indicates process flow)
       if len(sizes) >= 3:
           avg_size = sum(sizes) / len(sizes)
           variations = [abs(size - avg_size) / avg_size for size in sizes if avg_size > 0]
           if variations and max(variations) < 0.5:  # Less than 50% variation
               score += 5

       return score

   def _analyze_text_density(self, text_blocks):
       """Analyze text characteristics for diagram indicators"""
       if not text_blocks:
           return 0

       short_text_count = 0
       total_blocks = len(text_blocks)

       for block in text_blocks:
           # Count average words per paragraph
           total_words = 0
           para_count = 0

           for para in block.get("paragraphs", []):
               clean_text = para.get("clean_text", "")
               if clean_text:
                   words = len(clean_text.split())
                   total_words += words
                   para_count += 1

           if para_count > 0:
               avg_words = total_words / para_count
               if avg_words <= 5:  # Short labels
                   short_text_count += 1

       # Score based on percentage of short text blocks
       if total_blocks > 0:
           short_ratio = short_text_count / total_blocks
           if short_ratio >= 0.7:  # 70%+ short text
               return 10
           elif short_ratio >= 0.5:  # 50%+ short text
               return 5

       return 0

   def _analyze_flow_patterns(self, shapes, lines, arrows, text_blocks):
       """Analyze for flow patterns and process keywords"""
       score = 0

       # Check for start/end keywords
       flow_keywords = ["start", "begin", "end", "finish", "process", "step", "decision"]
       action_words = ["create", "update", "check", "verify", "send", "receive", "analyze"]

       all_text = ""
       for block in text_blocks:
           for para in block.get("paragraphs", []):
               all_text += " " + para.get("clean_text", "").lower()

       flow_matches = sum(1 for keyword in flow_keywords if keyword in all_text)
       action_matches = sum(1 for keyword in action_words if keyword in all_text)

       if flow_matches >= 2:
           score += 20
       elif flow_matches >= 1:
           score += 10

       if action_matches >= 3:
           score += 10

       # Bonus for having both shapes and connecting elements
       if len(shapes) >= 3 and (len(lines) > 0 or len(arrows) > 0):
           score += 15

       return score

   def _analyze_negative_indicators(self, text_blocks, shapes):
       """Check for negative indicators that suggest NOT a diagram"""
       score = 0

       # Check for long paragraphs
       long_text_count = 0
       bullet_count = 0

       for block in text_blocks:
           for para in block.get("paragraphs", []):
               clean_text = para.get("clean_text", "")
               if clean_text:
                   word_count = len(clean_text.split())
                   if word_count > 20:  # Long paragraph
                       long_text_count += 1

                   # Check for bullet points
                   if para.get("hints", {}).get("is_bullet", False):
                       bullet_count += 1

       # Penalize long text
       if long_text_count >= 2:
           score -= 15

       # Penalize if mostly bullet points
       total_paras = sum(len(block.get("paragraphs", [])) for block in text_blocks)
       if total_paras > 0 and bullet_count / total_paras > 0.8:
           score -= 10

       # Penalize single column layout (all shapes vertically aligned)
       if len(shapes) >= 3:
           positions = [s.get("position") for s in shapes if s.get("position")]
           if len(positions) >= 3:
               lefts = [p["left"] for p in positions]
               left_variance = max(lefts) - min(lefts) if lefts else 0
               if left_variance < 500000:  # Very narrow horizontal spread
                   score -= 10

       return score

   def is_arrow_shape(self, auto_shape_type):
       """Check if an auto shape type is an arrow"""
       if not auto_shape_type:
           return False

       arrow_types = [
           "LEFT_ARROW", "DOWN_ARROW", "UP_ARROW", "RIGHT_ARROW",
           "LEFT_RIGHT_ARROW", "UP_DOWN_ARROW", "QUAD_ARROW",
           "LEFT_RIGHT_UP_ARROW", "BENT_ARROW", "U_TURN_ARROW",
           "CURVED_LEFT_ARROW", "CURVED_RIGHT_ARROW",
           "CURVED_UP_ARROW", "CURVED_DOWN_ARROW",
           "STRIPED_RIGHT_ARROW", "NOTCHED_RIGHT_ARROW",
           "BLOCK_ARC"
       ]

       return any(arrow_type in auto_shape_type for arrow_type in arrow_types)

   def extract_pptx_metadata(self, presentation, file_path):
       """Extract comprehensive metadata from PowerPoint file"""
       metadata = {}

       try:
           # Core properties from PowerPoint
           core_props = presentation.core_properties

           # Basic file info
           metadata['filename'] = os.path.basename(file_path)
           metadata['file_size'] = os.path.getsize(file_path) if os.path.exists(file_path) else None

           # Document properties
           metadata['title'] = getattr(core_props, 'title', '') or ''
           metadata['author'] = getattr(core_props, 'author', '') or ''
           metadata['subject'] = getattr(core_props, 'subject', '') or ''
           metadata['keywords'] = getattr(core_props, 'keywords', '') or ''
           metadata['comments'] = getattr(core_props, 'comments', '') or ''
           metadata['category'] = getattr(core_props, 'category', '') or ''
           metadata['content_status'] = getattr(core_props, 'content_status', '') or ''
           metadata['language'] = getattr(core_props, 'language', '') or ''
           metadata['version'] = getattr(core_props, 'version', '') or ''

           # Dates
           metadata['created'] = getattr(core_props, 'created', None)
           metadata['modified'] = getattr(core_props, 'modified', None)
           metadata['last_modified_by'] = getattr(core_props, 'last_modified_by', '') or ''
           metadata['last_printed'] = getattr(core_props, 'last_printed', None)

           # Revision and identifier
           metadata['revision'] = getattr(core_props, 'revision', None)
           metadata['identifier'] = getattr(core_props, 'identifier', '') or ''

           # Presentation-specific info
           metadata['slide_count'] = len(presentation.slides)

           # Try to extract slide master themes/layouts
           try:
               slide_masters = presentation.slide_masters
               if slide_masters:
                   metadata['slide_master_count'] = len(slide_masters)
                   # Get layout names if available
                   layout_names = []
                   for master in slide_masters:
                       for layout in master.slide_layouts:
                           if hasattr(layout, 'name') and layout.name:
                               layout_names.append(layout.name)
                   metadata['layout_types'] = ', '.join(set(layout_names)) if layout_names else ''
           except:
               metadata['slide_master_count'] = 0
               metadata['layout_types'] = ''

           # Application that created the file
           try:
               app_props = presentation.app_properties if hasattr(presentation, 'app_properties') else None
               if app_props:
                   metadata['application'] = getattr(app_props, 'application', '') or ''
                   metadata['app_version'] = getattr(app_props, 'app_version', '') or ''
                   metadata['company'] = getattr(app_props, 'company', '') or ''
                   metadata['doc_security'] = getattr(app_props, 'doc_security', None)
           except:
               metadata['application'] = ''
               metadata['app_version'] = ''
               metadata['company'] = ''
               metadata['doc_security'] = None

       except Exception as e:
           print(f"Warning: Could not extract some metadata: {e}")

       return metadata

   def add_pptx_metadata_for_claude(self, markdown_content, metadata):
       """Add PowerPoint metadata as comments for Claude to incorporate"""
       # Format metadata for Claude
       metadata_comments = "\n<!-- POWERPOINT METADATA FOR CLAUDE:\n"

       if metadata.get('title'):
           metadata_comments += f"Document Title: {metadata['title']}\n"
       if metadata.get('author'):
           metadata_comments += f"Author: {metadata['author']}\n"
       if metadata.get('subject'):
           metadata_comments += f"Subject: {metadata['subject']}\n"
       if metadata.get('keywords'):
           metadata_comments += f"Keywords: {metadata['keywords']}\n"
       if metadata.get('category'):
           metadata_comments += f"Category: {metadata['category']}\n"
       if metadata.get('comments'):
           metadata_comments += f"Document Comments: {metadata['comments']}\n"
       if metadata.get('created'):
           metadata_comments += f"Created Date: {metadata['created']}\n"
       if metadata.get('modified'):
           metadata_comments += f"Last Modified: {metadata['modified']}\n"
       if metadata.get('last_modified_by'):
           metadata_comments += f"Last Modified By: {metadata['last_modified_by']}\n"
       if metadata.get('version'):
           metadata_comments += f"Version: {metadata['version']}\n"
       if metadata.get('application'):
           metadata_comments += f"Created With: {metadata['application']}\n"
       if metadata.get('company'):
           metadata_comments += f"Company: {metadata['company']}\n"
       if metadata.get('language'):
           metadata_comments += f"Language: {metadata['language']}\n"
       if metadata.get('content_status'):
           metadata_comments += f"Content Status: {metadata['content_status']}\n"

       # File info
       metadata_comments += f"Filename: {metadata.get('filename', 'unknown')}\n"
       metadata_comments += f"Slide Count: {metadata.get('slide_count', 0)}\n"
       if metadata.get('slide_size'):
           metadata_comments += f"Slide Size: {metadata['slide_size']}\n"
       if metadata.get('layout_types'):
           metadata_comments += f"Layout Types: {metadata['layout_types']}\n"

       metadata_comments += "-->\n"

       # Add metadata at the beginning
       return metadata_comments + markdown_content

   def extract_text_frame_fixed(self, text_frame, shape):
       """Fixed text extraction with proper bullet detection"""
       if not text_frame.paragraphs:
           return None

       block = {
           "type": "text",
           "paragraphs": [],
           "shape_hyperlink": self.extract_shape_hyperlink(shape)
       }

       for para_idx, para in enumerate(text_frame.paragraphs):
           if not para.text.strip():
               continue

           para_data = self.process_paragraph_fixed(para)
           if para_data:
               block["paragraphs"].append(para_data)

       return block if block["paragraphs"] else None

   def process_paragraph_fixed(self, para):
       """Fixed paragraph processing with reliable bullet detection"""
       raw_text = para.text
       if not raw_text.strip():
           return None

       # First, check if PowerPoint knows this is a bullet
       ppt_level = getattr(para, 'level', None)

       # Check XML for bullet formatting
       is_ppt_bullet = False
       xml_level = None

       try:
           if hasattr(para, '_p') and para._p is not None:
               xml_str = str(para._p.xml)
               # Look for bullet indicators
               if any(indicator in xml_str for indicator in ['buChar', 'buAutoNum', 'buFont']):
                   is_ppt_bullet = True
                   # Try to extract level
                   import re
                   level_match = re.search(r'lvl="(\d+)"', xml_str)
                   if level_match:
                       xml_level = int(level_match.group(1))
       except:
           pass

       # Determine final bullet level
       bullet_level = -1
       if is_ppt_bullet:
           bullet_level = xml_level if xml_level is not None else (ppt_level if ppt_level is not None else 0)
       elif ppt_level is not None:
           # PowerPoint says it has a level, trust it
           bullet_level = ppt_level

       # Check for manual bullets and numbered lists
       clean_text = raw_text.strip()
       manual_bullet = self.is_manual_bullet(clean_text)
       numbered = self.is_numbered_list(clean_text)

       # If we found a manual bullet but PowerPoint didn't recognize it
       if manual_bullet and bullet_level < 0:
           # Estimate level from indentation
           leading_spaces = len(raw_text) - len(raw_text.lstrip())
           bullet_level = min(leading_spaces // 2, 6)
           clean_text = self.remove_bullet_char(clean_text)
       elif bullet_level >= 0:
           # Remove any manual bullet chars if PowerPoint formatted it
           clean_text = self.remove_bullet_char(clean_text)
       elif numbered:
           clean_text = self.remove_number_prefix(clean_text)

       # Extract formatted runs - THIS IS KEY!
       formatted_runs = self.extract_runs_with_text_preservation(para.runs, clean_text, bullet_level >= 0 or numbered)

       para_data = {
           "raw_text": raw_text,
           "clean_text": clean_text,
           "formatted_runs": formatted_runs,
           "hints": {
               "has_powerpoint_level": ppt_level is not None,
               "powerpoint_level": ppt_level,
               "bullet_level": bullet_level,
               "is_bullet": bullet_level >= 0,
               "is_numbered": numbered,
               "starts_with_bullet": manual_bullet,
               "starts_with_number": numbered,
               "short_text": len(clean_text) < 100,
               "all_caps": clean_text.isupper() if clean_text else False,
               "likely_heading": self.is_likely_heading(clean_text)
           }
       }

       return para_data

   def extract_runs_with_text_preservation(self, runs, clean_text, has_prefix_removed):
       """Extract runs while preserving formatting after bullet/number removal"""
       if not runs:
           return [{"text": clean_text, "bold": False, "italic": False, "hyperlink": None}]

       formatted_runs = []

       # If we removed a prefix (bullet/number), we need to adjust the runs
       if has_prefix_removed:
           # Find where the clean text starts in the original runs
           full_text = "".join(run.text for run in runs)

           # Find the start position of clean_text in full_text
           # This is tricky because clean_text has had prefixes removed
           start_pos = -1
           for i in range(len(full_text)):
               remaining = full_text[i:].strip()
               if remaining == clean_text:
                   start_pos = i
                   break

           if start_pos == -1:
               # Fallback: just process runs normally
               start_pos = 0

           # Now process runs, skipping content before start_pos
           char_count = 0
           for run in runs:
               run_text = run.text
               run_start = char_count
               run_end = char_count + len(run_text)

               # Skip if this run is entirely before our clean text
               if run_end <= start_pos:
                   char_count += len(run_text)
                   continue

               # Adjust text if run spans the start position
               if run_start < start_pos < run_end:
                   run_text = run_text[start_pos - run_start:]

               if run_text:
                   formatted_runs.append(self.extract_run_formatting(run, run_text))

               char_count += len(run.text)
       else:
           # No prefix removed, process runs normally
           for run in runs:
               if run.text:
                   formatted_runs.append(self.extract_run_formatting(run, run.text))

       return formatted_runs

   def extract_run_formatting(self, run, text_override=None):
       """Extract formatting from a single run"""
       run_data = {
           "text": text_override if text_override is not None else run.text,
           "bold": False,
           "italic": False,
           "hyperlink": None
       }

       # Get formatting
       try:
           font = run.font
           if hasattr(font, 'bold') and font.bold:
               run_data["bold"] = True
           if hasattr(font, 'italic') and font.italic:
               run_data["italic"] = True
       except:
           pass

       # Get hyperlinks
       try:
           if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
               run_data["hyperlink"] = self.fix_url(run.hyperlink.address)
       except:
           pass

       return run_data

   def is_manual_bullet(self, text):
       """Check if text starts with a manual bullet character"""
       if not text:
           return False
       bullet_chars = '•◦▪▫‣·○■□→►✓✗-*+※◆◇'
       return text[0] in bullet_chars

   def is_numbered_list(self, text):
       """Check if text starts with a number pattern"""
       patterns = [
           r'^\d+[\.\)]\s+',  # 1. or 1)
           r'^[a-zA-Z][\.\)]\s+',  # a. or A)
           r'^[ivxlcdm]+[\.\)]\s+',  # Roman numerals (lowercase)
           r'^[IVXLCDM]+[\.\)]\s+',  # Roman numerals (uppercase)
       ]
       return any(re.match(pattern, text) for pattern in patterns)

   def remove_bullet_char(self, text):
       """Remove bullet characters from start of text"""
       if not text:
           return text
       # Remove common bullet chars and following spaces
       return re.sub(r'^[•◦▪▫‣·○■□→►✓✗\-\*\+※◆◇]\s*', '', text)

   def remove_number_prefix(self, text):
       """Remove number prefix from text"""
       return re.sub(r'^[^\s]+\s+', '', text)

   def is_likely_heading(self, text):
       """Determine if text is likely a heading"""
       if not text or len(text) > 150:
           return False

       # All caps
       if text.isupper() and len(text) > 2:
           return True

       # Short text without ending punctuation
       if len(text) < 80 and not text.endswith(('.', '!', '?', ';', ':', ',')):
           return True

       return False

   def extract_plain_text(self, shape):
       """Extract plain text from shape"""
       if not hasattr(shape, 'text') or not shape.text:
           return None

       return {
           "type": "text",
           "paragraphs": [{
               "raw_text": shape.text,
               "clean_text": shape.text.strip(),
               "formatted_runs": [{"text": shape.text, "bold": False, "italic": False, "hyperlink": None}],
               "hints": self._analyze_plain_text_hints(shape.text)
           }],
           "shape_hyperlink": self.extract_shape_hyperlink(shape)
       }

   def _analyze_plain_text_hints(self, text):
       """Analyze plain text for formatting hints"""
       if not text:
           return {}

       stripped = text.strip()

       # Check each line for bullets
       lines = text.split('\n')
       has_bullets = False
       for line in lines:
           if line.strip() and self.is_manual_bullet(line.strip()):
               has_bullets = True
               break

       return {
           "has_powerpoint_level": False,
           "powerpoint_level": None,
           "bullet_level": -1,
           "is_bullet": has_bullets,
           "is_numbered": any(self.is_numbered_list(line.strip()) for line in lines if line.strip()),
           "starts_with_bullet": stripped and self.is_manual_bullet(stripped),
           "starts_with_number": bool(re.match(r'^\s*\d+[\.\)]\s', text)),
           "short_text": len(stripped) < 100,
           "all_caps": stripped.isupper() if stripped else False,
           "likely_heading": self.is_likely_heading(stripped)
       }

   def extract_image(self, shape):
       """Extract image info with proper alt text extraction"""
       alt_text = "Image"

       try:
           # Try multiple methods to get alt text
           if hasattr(shape, 'alt_text') and shape.alt_text:
               alt_text = shape.alt_text
           elif hasattr(shape, 'image') and hasattr(shape.image, 'alt_text') and shape.image.alt_text:
               alt_text = shape.image.alt_text
           elif hasattr(shape, '_element'):
               # Try to extract from XML
               try:
                   xml_str = str(shape._element.xml) if hasattr(shape._element, 'xml') else ""
                   if xml_str:
                       root = ET.fromstring(xml_str)
                       # Look for description attributes
                       for elem in root.iter():
                           if 'descr' in elem.attrib and elem.attrib['descr']:
                               alt_text = elem.attrib['descr']
                               break
                           elif 'title' in elem.attrib and elem.attrib['title']:
                               alt_text = elem.attrib['title']
                               break
               except:
                   pass
       except:
           pass

       return {
           "type": "image",
           "alt_text": alt_text.strip() if alt_text else "Image",
           "hyperlink": self.extract_shape_hyperlink(shape)
       }

   def extract_table(self, table):
       """Extract table data"""
       if not table.rows:
           return None

       table_data = []
       for row in table.rows:
           row_data = []
           for cell in row.cells:
               # Extract cell text with formatting
               cell_content = ""
               if hasattr(cell, 'text_frame') and cell.text_frame:
                   cell_paras = []
                   for para in cell.text_frame.paragraphs:
                       if para.text.strip():
                           # Process paragraph for bullets
                           para_processed = self.process_paragraph_fixed(para)
                           if para_processed and para_processed['hints']['is_bullet']:
                               level = para_processed['hints']['bullet_level']
                               indent = "  " * level
                               cell_paras.append(f"{indent}• {para_processed['clean_text']}")
                           elif para_processed:
                               cell_paras.append(para_processed['clean_text'])
                   cell_content = " ".join(cell_paras)
               else:
                   cell_content = cell.text.strip() if hasattr(cell, 'text') else ""
               row_data.append(cell_content)
           table_data.append(row_data)

       return {
           "type": "table",
           "data": table_data
       }

   def extract_chart(self, shape):
       """Extract chart/diagram information"""
       try:
           chart = shape.chart
           chart_data = {
               "type": "chart",
               "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
               "title": "",
               "data_points": [],
               "categories": [],
               "series": [],
               "hyperlink": self.extract_shape_hyperlink(shape)
           }

           # Try to get chart title
           try:
               if hasattr(chart, 'chart_title') and chart.chart_title and hasattr(chart.chart_title, 'text_frame'):
                   chart_data["title"] = chart.chart_title.text_frame.text.strip()
           except:
               pass

           # Try to extract data for potential Mermaid conversion
           try:
               if hasattr(chart, 'plots') and chart.plots:
                   plot = chart.plots[0]
                   if hasattr(plot, 'categories') and plot.categories:
                       chart_data["categories"] = [cat.label for cat in plot.categories if hasattr(cat, 'label')]

                   if hasattr(plot, 'series') and plot.series:
                       for series in plot.series:
                           series_data = {
                               "name": series.name if hasattr(series, 'name') else "",
                               "values": []
                           }
                           if hasattr(series, 'values'):
                               try:
                                   series_data["values"] = [val for val in series.values if val is not None]
                               except:
                                   pass
                           chart_data["series"].append(series_data)
           except:
               pass

           return chart_data

       except Exception:
           # Fallback for charts we can't parse
           return {
               "type": "chart",
               "chart_type": "unknown",
               "title": "Chart",
               "data_points": [],
               "categories": [],
               "series": [],
               "hyperlink": self.extract_shape_hyperlink(shape)
           }

   def extract_group(self, shape):
       """Extract content from grouped shapes - EXACT v14 approach that was working"""
       try:
           # For grouped shapes, extract text from all child shapes
           extracted_blocks = []

           for child_shape in shape.shapes:
               # Extract text directly from each child shape - EXACTLY like v14
               if hasattr(child_shape, 'text_frame') and child_shape.text_frame:
                   text_block = self.extract_text_frame_fixed(child_shape.text_frame, child_shape)
                   if text_block:
                       extracted_blocks.append(text_block)
               elif hasattr(child_shape, 'text') and child_shape.text:
                   text_block = self.extract_plain_text(child_shape)
                   if text_block:
                       extracted_blocks.append(text_block)
               elif child_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                   image_block = self.extract_image(child_shape)
                   if image_block:
                       extracted_blocks.append(image_block)
               elif child_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                   table_block = self.extract_table(child_shape.table)
                   if table_block:
                       extracted_blocks.append(table_block)
               elif hasattr(child_shape, 'has_chart') and child_shape.has_chart:
                   chart_block = self.extract_chart(child_shape)
                   if chart_block:
                       extracted_blocks.append(chart_block)
               # Handle nested groups recursively
               elif child_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                   nested_group = self.extract_group(child_shape)
                   if nested_group and nested_group.get("extracted_blocks"):
                       extracted_blocks.extend(nested_group["extracted_blocks"])

           # Return a simplified group structure
           if extracted_blocks:
               return {
                   "type": "group",
                   "extracted_blocks": extracted_blocks,
                   "hyperlink": self.extract_shape_hyperlink(shape)
               }

           return None

       except Exception as e:
           print(f"Error extracting group: {e}")
           return None

   def extract_shape_hyperlink(self, shape):
       """Extract shape-level hyperlink"""
       try:
           if hasattr(shape, 'click_action') and shape.click_action:
               if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink:
                   if shape.click_action.hyperlink.address:
                       return self.fix_url(shape.click_action.hyperlink.address)
       except:
           pass
       return None

   def fix_url(self, url):
       """Fix URLs by adding schemes if missing"""
       if not url:
           return url

       if '@' in url and not url.startswith('mailto:'):
           return f"mailto:{url}"

       if not url.startswith(('http://', 'https://', 'mailto:', 'tel:', 'ftp://', '#')):
           if url.startswith('www.') or any(
                   domain in url.lower() for domain in ['.com', '.org', '.net', '.edu', '.gov', '.io']):
               return f"https://{url}"

       return url

   def convert_structured_data_to_markdown(self, data, convert_slide_titles=True):
       """Convert structured data to markdown"""
       markdown_parts = []

       for slide in data["slides"]:
           # Add slide marker
           markdown_parts.append(f"\n<!-- Slide {slide['slide_number']} -->\n")

           for block in slide["content_blocks"]:
               if block["type"] == "text":
                   markdown_parts.append(self.convert_text_block_to_markdown(block))
               elif block["type"] == "table":
                   markdown_parts.append(self.convert_table_to_markdown(block))
               elif block["type"] == "image":
                   markdown_parts.append(self.convert_image_to_markdown(block))
               elif block["type"] == "chart":
                   markdown_parts.append(self.convert_chart_to_markdown(block))
               elif block["type"] == "group":
                   markdown_parts.append(self.convert_group_to_markdown(block))

       markdown_content = "\n\n".join(filter(None, markdown_parts))

       # Post-process to convert slide titles from bullets to H1 headings if requested
       if convert_slide_titles:
           markdown_content = self.convert_slide_titles_to_headings(markdown_content)

       return markdown_content

   def convert_slide_titles_to_headings(self, markdown_content):
       """
       Post-process markdown to convert slide titles from bullet points to H1 headings.

       This function identifies likely slide titles by looking for bullet points that appear
       immediately after slide markers and have title-like characteristics.
       """
       lines = markdown_content.split('\n')
       processed_lines = []

       i = 0
       while i < len(lines):
           line = lines[i]
           processed_lines.append(line)

           # Check if this is a slide marker
           if line.strip().startswith('<!-- Slide ') and line.strip().endswith(' -->'):
               # Look ahead for the first non-empty content line
               j = i + 1
               while j < len(lines) and not lines[j].strip():
                   processed_lines.append(lines[j])
                   j += 1

               # Check if the next content line is a bullet that looks like a title
               if j < len(lines):
                   next_line = lines[j].strip()
                   if self.is_likely_slide_title(next_line):
                       # Convert bullet to H1 heading
                       title_text = self.extract_title_from_bullet(next_line)
                       processed_lines.append(f"\n# {title_text}")
                       i = j  # Skip the original bullet line
                   else:
                       i = j - 1  # Process the next line normally
               else:
                   break

           i += 1

       return '\n'.join(processed_lines)

   def is_likely_slide_title(self, line):
       """
       Determine if a line is likely a slide title based on formatting and content.

       Args:
           line (str): The line to evaluate

       Returns:
           bool: True if the line appears to be a slide title
       """
       if not line.strip():
           return False

       # Must be a bullet point to be converted
       if not line.startswith('- '):
           return False

       # Extract the text content
       text_content = line[2:].strip()

       # Title characteristics
       title_indicators = [
           len(text_content) <= 150,  # Reasonable title length
           not text_content.endswith(('.', '!', '?', ';', ':')),  # Titles typically don't end with punctuation
           not self._contains_multiple_sentences(text_content),  # Titles are usually single phrases
           not text_content.lower().startswith(('the following', 'here are', 'this slide', 'key points')),
           # Avoid descriptive text
       ]

       # Additional positive indicators
       positive_indicators = [
           text_content.isupper(),  # All caps suggests title
           text_content.istitle(),  # Title case suggests title
           len(text_content.split()) <= 10,  # Short phrases are more likely titles
           any(word in text_content.lower() for word in
               ['overview', 'introduction', 'conclusion', 'agenda', 'objectives']),  # Common title words
       ]

       # Must meet basic criteria and have at least one positive indicator
       basic_criteria_met = all(title_indicators)
       has_positive_indicator = any(positive_indicators)

       return basic_criteria_met and (has_positive_indicator or len(text_content.split()) <= 6)

   def extract_title_from_bullet(self, bullet_line):
       """
       Extract clean title text from a bullet point line.

       Args:
           bullet_line (str): The bullet point line (e.g., "- Title Text")

       Returns:
           str: Clean title text
       """
       # Remove bullet prefix
       title_text = bullet_line[2:].strip()

       # Clean up common title artifacts
       title_text = title_text.strip('*_`')  # Remove markdown formatting artifacts

       return title_text

   def _contains_multiple_sentences(self, text):
       """
       Check if text contains multiple sentences.

       Args:
           text (str): Text to check

       Returns:
           bool: True if text appears to contain multiple sentences
       """
       # Simple heuristic: look for sentence-ending punctuation followed by space and capital letter
       sentence_pattern = r'[.!?]\s+[A-Z]'
       return bool(re.search(sentence_pattern, text))

   def convert_text_block_to_markdown(self, block):
       """Convert text block to markdown with proper formatting"""
       lines = []

       for para in block["paragraphs"]:
           line = self.convert_paragraph_to_markdown(para)
           if line:
               lines.append(line)

       # If entire shape is a hyperlink, wrap it
       result = "\n".join(lines)
       if block.get("shape_hyperlink") and result:
           result = f"[{result}]({block['shape_hyperlink']})"

       return result

   def convert_paragraph_to_markdown(self, para):
       """Convert paragraph to markdown with correct formatting"""
       if not para.get("clean_text"):
           return ""

       # Build formatted text from runs
       formatted_text = self.build_formatted_text_from_runs(para["formatted_runs"], para["clean_text"])

       # Now apply structural formatting based on hints
       hints = para.get("hints", {})

       # Bullets
       if hints.get("is_bullet", False):
           level = hints.get("bullet_level", 0)
           if level < 0:
               level = 0
           indent = "  " * level
           return f"{indent}- {formatted_text}"

       # Numbered lists
       elif hints.get("is_numbered", False):
           return f"1. {formatted_text}"

       # Headings
       elif hints.get("likely_heading", False):
           # Determine heading level
           if hints.get("all_caps") or len(para["clean_text"]) < 30:
               return f"## {formatted_text}"
           else:
               return f"### {formatted_text}"

       # Regular paragraph
       else:
           return formatted_text

   def build_formatted_text_from_runs(self, runs, clean_text):
       """Build formatted text from runs, handling edge cases"""
       if not runs:
           return clean_text

       # First check if we have any formatting at all
       has_formatting = any(
           run.get("bold") or run.get("italic") or run.get("hyperlink")
           for run in runs
       )

       if not has_formatting:
           return clean_text

       # Build formatted text preserving run boundaries
       formatted_parts = []

       for run in runs:
           text = run["text"]
           if not text:
               continue

           # Apply formatting
           if run.get("bold") and run.get("italic"):
               text = f"***{text}***"
           elif run.get("bold"):
               text = f"**{text}**"
           elif run.get("italic"):
               text = f"*{text}*"

           # Apply hyperlink
           if run.get("hyperlink"):
               text = f"[{text}]({run['hyperlink']})"

           formatted_parts.append(text)

       return "".join(formatted_parts)

   def convert_table_to_markdown(self, block):
       """Convert table to markdown"""
       if not block["data"]:
           return ""

       markdown = ""
       for i, row in enumerate(block["data"]):
           markdown += "| " + " | ".join(cell.replace("|", "\\|") for cell in row) + " |\n"

           # Add separator after header
           if i == 0:
               markdown += "| " + " | ".join("---" for _ in row) + " |\n"

       return markdown

   def convert_image_to_markdown(self, block):
       """Convert image to markdown"""
       image_md = f"![{block['alt_text']}](image)"

       if block.get("hyperlink"):
           image_md = f"[{image_md}]({block['hyperlink']})"

       return image_md

   def convert_chart_to_markdown(self, block):
       """Convert chart to markdown"""
       chart_md = f"**Chart: {block.get('title', 'Untitled Chart')}**\n"
       chart_md += f"*Chart Type: {block.get('chart_type', 'unknown')}*\n\n"

       # Add data if available
       if block.get('categories') and block.get('series'):
           chart_md += "Data:\n"
           for series in block['series']:
               if series.get('name'):
                   chart_md += f"- {series['name']}: "
                   if series.get('values'):
                       chart_md += ", ".join(map(str, series['values'][:5]))
                       if len(series['values']) > 5:
                           chart_md += "..."
                   chart_md += "\n"

       # Add comment for diagram conversion
       chart_md += f"\n<!-- DIAGRAM_CANDIDATE: chart, type={block.get('chart_type', 'unknown')} -->\n"

       if block.get("hyperlink"):
           chart_md = f"[{chart_md}]({block['hyperlink']})"

       return chart_md

   def convert_group_to_markdown(self, block):
       """Convert grouped shapes to markdown - handle all shape types"""
       # Get the extracted blocks from the group
       extracted_blocks = block.get("extracted_blocks", [])

       if not extracted_blocks:
           return ""

       # Convert each extracted block to markdown
       content_parts = []

       for extracted_block in extracted_blocks:
           if extracted_block["type"] == "text":
               content = self.convert_text_block_to_markdown(extracted_block)
               if content:
                   content_parts.append(content)
           elif extracted_block["type"] == "image":
               content = self.convert_image_to_markdown(extracted_block)
               if content:
                   content_parts.append(content)
           elif extracted_block["type"] == "table":
               content = self.convert_table_to_markdown(extracted_block)
               if content:
                   content_parts.append(content)
           elif extracted_block["type"] == "chart":
               content = self.convert_chart_to_markdown(extracted_block)
               if content:
                   content_parts.append(content)
           elif extracted_block["type"] == "line":
               # Lines don't produce visible content but are tracked for diagram analysis
               pass
           elif extracted_block["type"] == "arrow":
               # Arrows don't produce visible content but are tracked for diagram analysis
               pass
           elif extracted_block["type"] == "shape":
               # Generic shapes might have minimal content
               content = f"[Shape: {extracted_block.get('shape_subtype', 'unknown')}]"
               content_parts.append(content)

       # Join all content together
       group_md = "\n\n".join(content_parts) if content_parts else ""

       # Add shape-level hyperlink if present
       if block.get("hyperlink") and group_md:
           group_md = f"[{group_md}]({block['hyperlink']})"

       return group_md


# Convenience functions for backward compatibility and ease of use

def convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles=True):
   """
   Convenience function to maintain backward compatibility

   Args:
       file_path (str): Path to the PowerPoint file
       convert_slide_titles (bool): Whether to convert slide titles from bullets to H1 headings
   """
   processor = PowerPointProcessor()
   return processor.convert_pptx_to_markdown_enhanced(file_path, convert_slide_titles)


def process_powerpoint_file(file_path, output_format="markdown", convert_slide_titles=True):
   """
   Convenience function for complete file processing

   Args:
       file_path (str): Path to the PowerPoint file
       output_format (str): "markdown", "json", "text", or "summary"
       convert_slide_titles (bool): Whether to convert slide titles from bullets to H1 headings

   Returns:
       dict: Processed content and metadata
   """
   processor = PowerPointProcessor()
   return processor.process_file_complete(file_path, output_format, convert_slide_titles)


# USAGE EXAMPLES:
"""
# Basic usage (with semantic accessibility by default - titles read first):
processor = PowerPointProcessor()
result = processor.convert_pptx_to_markdown_enhanced("presentation.pptx")

# Debug semantic accessibility order (shows titles first, then content):
processor.debug_accessibility_order("presentation.pptx", slide_number=1)

# Disable accessibility and use original positional order:
processor.use_accessibility_order = False
result = processor.convert_pptx_to_markdown_enhanced("presentation.pptx")

# Check extraction method used:
prs = Presentation("presentation.pptx")
slide_data = processor.extract_slide_data(prs.slides[0], 1)
print(f"Extraction method: {slide_data.get('extraction_method', 'not_specified')}")

# The new approach ensures:
# 1. Titles are read first (semantic hierarchy)
# 2. Content follows in XML document order (true reading order)
# 3. No coordinate-based positioning unless explicitly disabled
# 4. Respects PowerPoint's accessibility features and placeholder semantics
"""