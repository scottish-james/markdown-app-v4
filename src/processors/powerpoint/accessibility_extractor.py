"""
Fixed Accessibility Order Extractor - DUPLICATE ELIMINATION
Key fixes:
1. Deduplicate shapes by ID in XML parsing
2. More selective XML element processing
3. Better shape tree navigation
4. Improved debugging to track duplicates
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
import re


class AccessibilityOrderExtractor:
    def __init__(self, use_accessibility_order=True):
        self.use_accessibility_order = use_accessibility_order
        self.last_extraction_method = "not_extracted"

        # XML namespaces for PowerPoint OOXML processing
        self.namespaces = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
            'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        }

    def get_slide_reading_order(self, slide, slide_number):
        """
        QUICK TEST: Force simple approach to see if it fixes duplicates.
        Replace your existing get_slide_reading_order method with this temporarily.
        """
        print(f"DEBUG: FORCING SIMPLE APPROACH for testing")

        # Get shapes directly from python-pptx (no XML parsing)
        original_shapes = list(slide.shapes)
        print(f"DEBUG: Original slide has {len(original_shapes)} shapes")

        # Expand groups if present
        final_shapes = []
        for shape in original_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                print(f"DEBUG: Expanding group: {getattr(shape, 'name', 'unnamed')}")
                group_children = list(shape.shapes)
                final_shapes.extend(group_children)
            else:
                final_shapes.append(shape)

        print(f"DEBUG: After group expansion: {len(final_shapes)} shapes")

        # Apply semantic roles WITHOUT XML parsing for order
        shapes_with_roles = []
        for shape in final_shapes:
            semantic_role = self._get_semantic_role_from_xml(shape)
            shapes_with_roles.append((shape, semantic_role))

            # Debug info
            shape_type = str(shape.shape_type).split('.')[-1]
            text_preview = ""
            try:
                if hasattr(shape, 'text') and shape.text:
                    text_preview = shape.text.strip()[:30]
            except:
                pass
            print(f"DEBUG: Shape {shape_type}, Role: {semantic_role}, Text: '{text_preview}'")

        # Sort by semantic priority
        priority_order = {"title": 1, "subtitle": 2, "content": 3, "other": 4}
        shapes_with_roles.sort(key=lambda x: priority_order.get(x[1], 4))

        result = [shape for shape, role in shapes_with_roles]

        self.last_extraction_method = "simple_test_approach"
        print(f"DEBUG: Simple approach returning {len(result)} shapes")

        return result

    def _get_semantic_accessibility_order(self, slide):
        """
        FIXED: Enhanced semantic ordering with duplicate elimination.
        """
        print(f"DEBUG: Starting semantic accessibility order extraction")

        # Step 1: Get all shapes in XML document order (now deduplicated)
        xml_ordered_shapes = self._get_xml_document_order_deduplicated(slide)
        print(f"DEBUG: XML document order returned {len(xml_ordered_shapes)} shapes (after deduplication)")

        # Step 2: Process groups by extracting children individually
        final_ordered_shapes = []
        group_count = 0
        expanded_children_count = 0

        for shape in xml_ordered_shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_count += 1
                print(f"DEBUG: Found group {group_count}: {getattr(shape, 'name', 'unnamed')}")
                group_children = self.get_reading_order_of_grouped_by_shape(shape)
                print(f"DEBUG: Group produced {len(group_children)} children")
                final_ordered_shapes.extend(group_children)
                expanded_children_count += len(group_children)
            else:
                final_ordered_shapes.append(shape)

        print(f"DEBUG: After group expansion: {len(final_ordered_shapes)} shapes total")
        print(f"DEBUG: Found {group_count} groups, expanded to {expanded_children_count} children")

        # Step 3: CRITICAL FIX - Deduplicate final shapes by object ID
        deduplicated_shapes = self._deduplicate_shapes_by_object_id(final_ordered_shapes)
        print(f"DEBUG: After final deduplication: {len(deduplicated_shapes)} shapes")

        # Step 4: Separate by semantic importance for final ordering
        title_shapes = []
        subtitle_shapes = []
        content_shapes = []
        other_shapes = []

        for shape in deduplicated_shapes:
            semantic_role = self._get_semantic_role_from_xml(shape)

            if semantic_role == "title":
                title_shapes.append(shape)
            elif semantic_role == "subtitle":
                subtitle_shapes.append(shape)
            elif semantic_role == "content":
                content_shapes.append(shape)
            else:
                other_shapes.append(shape)

        print(
            f"DEBUG: Semantic classification: {len(title_shapes)} titles, {len(subtitle_shapes)} subtitles, {len(content_shapes)} content, {len(other_shapes)} other")

        # Step 5: Return in semantic priority order
        result = title_shapes + subtitle_shapes + content_shapes + other_shapes
        print(f"DEBUG: Final semantic order: {len(result)} shapes")

        return deduplicated_shapes

    def _get_xml_document_order_deduplicated(self, slide):
        """
        FIXED: Extract shapes in XML document order with deduplication.
        """
        try:
            # Step 1: Get the slide's raw XML
            slide_xml = self._get_slide_xml(slide)

            # Step 2: Parse XML to get shapes in document order (now deduplicated)
            xml_shape_info = self._parse_slide_xml_for_document_order_deduplicated(slide_xml)
            print(f"DEBUG: XML parsing found {len(xml_shape_info)} unique shapes")

            # Step 3: Map XML order to python-pptx shapes
            ordered_shapes = self._map_xml_to_pptx_shapes_deduplicated(xml_shape_info, slide.shapes)

            return ordered_shapes
        except Exception as e:
            raise Exception(f"XML document order extraction failed: {e}")

    def _parse_slide_xml_for_document_order_deduplicated(self, slide_xml):
        """
        FIXED: Parse slide XML with deduplication by shape ID.
        """
        root = ET.fromstring(slide_xml)

        # Find the PRIMARY shape tree - more selective than before
        shape_tree = root.find('.//p:cSld/p:spTree', self.namespaces)
        if shape_tree is None:
            # Fallback to any shape tree
            shape_tree = root.find('.//p:spTree', self.namespaces)

        if shape_tree is None:
            raise Exception("No shape tree found in slide XML")

        shape_order_info = []
        seen_shape_ids = set()  # Track IDs to prevent duplicates

        print(f"DEBUG: Processing shape tree with {len(list(shape_tree))} direct children")

        # Process ONLY direct children of the shape tree to avoid nested duplicates
        for idx, elem in enumerate(shape_tree):
            tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag

            # Only process actual shape elements at the top level
            if tag_name in ['sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp']:
                shape_info = self._extract_shape_info_from_xml(elem, idx)

                if shape_info and shape_info.get('id'):
                    shape_id = shape_info['id']

                    # CRITICAL FIX: Only add if not seen before
                    if shape_id not in seen_shape_ids:
                        shape_order_info.append(shape_info)
                        seen_shape_ids.add(shape_id)
                        print(f"DEBUG: Added unique shape ID {shape_id}: {shape_info.get('name', 'unnamed')}")
                    else:
                        print(f"DEBUG: Skipped duplicate shape ID {shape_id}: {shape_info.get('name', 'unnamed')}")
                elif shape_info:
                    # Shape without ID - add with warning
                    shape_order_info.append(shape_info)
                    print(f"DEBUG: Added shape without ID: {shape_info.get('name', 'unnamed')}")

        print(f"DEBUG: Final unique shapes after XML deduplication: {len(shape_order_info)}")
        return shape_order_info

    def _map_xml_to_pptx_shapes_deduplicated(self, xml_shape_info, pptx_shapes):
        """
        FIXED: Map XML shape information to python-pptx shapes with better deduplication.
        """
        ordered_shapes = []
        used_shape_object_ids = set()  # Track by object ID

        print(f"DEBUG: Mapping {len(xml_shape_info)} XML shapes to {len(pptx_shapes)} python-pptx shapes")

        # Create efficient lookup tables for matching
        shape_lookup = {}
        for shape in pptx_shapes:
            shape_id = self._get_shape_id(shape)
            shape_name = self._get_shape_name(shape)

            # Use prefixed keys to avoid ID/name collisions
            if shape_id:
                shape_lookup[f"id_{shape_id}"] = shape
            if shape_name:
                # Only use name if ID not available (names can be duplicate)
                if not shape_id:
                    shape_lookup[f"name_{shape_name}"] = shape

        # Match XML order to python-pptx shapes
        for xml_info in xml_shape_info:
            matched_shape = None

            # Priority 1: Try ID matching (most reliable)
            if xml_info['id']:
                matched_shape = shape_lookup.get(f"id_{xml_info['id']}")

            # Priority 2: Try name matching only if no ID match
            if not matched_shape and xml_info['name']:
                matched_shape = shape_lookup.get(f"name_{xml_info['name']}")

            # Add if found and not already used
            if matched_shape:
                shape_obj_id = id(matched_shape)
                if shape_obj_id not in used_shape_object_ids:
                    ordered_shapes.append(matched_shape)
                    used_shape_object_ids.add(shape_obj_id)
                    print(f"DEBUG: Mapped XML shape {xml_info.get('id', 'no-id')} to python-pptx shape")
                else:
                    print(f"DEBUG: Skipped already-used shape {xml_info.get('id', 'no-id')}")

        # Add any remaining unmatched shapes at the end
        for shape in pptx_shapes:
            shape_obj_id = id(shape)
            if shape_obj_id not in used_shape_object_ids:
                ordered_shapes.append(shape)
                print(f"DEBUG: Added unmatched shape: {getattr(shape, 'name', 'unnamed')}")

        print(f"DEBUG: Final mapped shapes: {len(ordered_shapes)}")
        return ordered_shapes

    def _deduplicate_shapes_by_object_id(self, shapes):
        """
        FINAL SAFETY NET: Remove any remaining duplicates by object identity.
        """
        seen_object_ids = set()
        deduplicated = []

        for shape in shapes:
            shape_obj_id = id(shape)
            if shape_obj_id not in seen_object_ids:
                deduplicated.append(shape)
                seen_object_ids.add(shape_obj_id)
            else:
                print(f"DEBUG: Removed final duplicate: {getattr(shape, 'name', 'unnamed')}")

        return deduplicated

    # [Include all the other existing methods from the original class]
    def _has_xml_access(self, slide):
        """XML availability check with multiple access pattern attempts."""
        try:
            slide_xml = self._get_slide_xml(slide)
            return slide_xml is not None and len(slide_xml) > 0
        except Exception:
            return False

    def _get_semantic_role_from_xml(self, shape):
        """XML-first semantic role detection using PowerPoint's internal XML structure."""
        shape_xml = self._get_shape_xml_content(shape)
        if not shape_xml:
            return "other"

        try:
            root = ET.fromstring(shape_xml)

            # Priority 1: Check XML placeholder types (most reliable)
            xml_role = self._extract_placeholder_type_from_xml(root)
            if xml_role:
                return xml_role

            # Priority 2: Check XML shape properties and names
            xml_role = self._extract_role_from_xml_properties(root)
            if xml_role:
                return xml_role

            # Priority 3: Analyse XML text content and positioning
            xml_role = self._extract_role_from_xml_content_analysis(root)
            return xml_role

        except Exception as e:
            return "other"

    def _get_shape_xml_content(self, shape):
        """Extract raw XML content from shape using XML-first approach."""
        try:
            if hasattr(shape, '_element') and hasattr(shape._element, 'xml'):
                return shape._element.xml
            elif hasattr(shape, 'element') and hasattr(shape.element, 'xml'):
                return shape.element.xml
            return None
        except Exception:
            return None

    def _extract_placeholder_type_from_xml(self, xml_root):
        """Parse XML to extract PowerPoint placeholder type information directly."""
        placeholder_elements = xml_root.findall('.//p:ph', self.namespaces)

        for ph_elem in placeholder_elements:
            ph_type = ph_elem.get('type', '').lower()

            if ph_type in ['sldnum', 'ftr', 'dt']:
                return "other"

            if ph_type in ['title', 'ctrtitle', 'centertitle']:
                return "title"

            if ph_type in ['subtitle', 'subhead']:
                return "subtitle"

            if ph_type in ['body', 'obj', 'tbl', 'chart', 'media']:
                return "content"

        return None

    def _extract_role_from_xml_properties(self, xml_root):
        """Parse XML shape properties for semantic role indicators."""
        cnv_pr_elements = xml_root.findall('.//p:cNvPr', self.namespaces)
        for cnv_pr in cnv_pr_elements:
            shape_name = cnv_pr.get('name', '').lower()

            if shape_name:
                if any(exclude_term in shape_name for exclude_term in [
                    'slide number', 'slide_number', 'page number', 'page_number',
                    'footer', 'date time', 'datetime'
                ]):
                    return "other"

                title_terms = ['title', 'heading', 'header']
                exclude_terms = ['subtitle', 'sub-title', 'sub_title']

                if any(title_term in shape_name for title_term in title_terms):
                    if not any(exclude_term in shape_name for exclude_term in exclude_terms):
                        if self._validate_title_from_xml_content(xml_root):
                            return "title"

                if any(subtitle_term in shape_name for subtitle_term in exclude_terms):
                    return "subtitle"

        return None

    def _extract_role_from_xml_content_analysis(self, xml_root):
        """Analyse XML text content and positioning to determine semantic role."""
        text_elements = xml_root.findall('.//a:t', self.namespaces)
        if text_elements:
            all_text = ' '.join([elem.text for elem in text_elements if elem.text])

            if all_text.strip():
                if self._is_slide_number_pattern_in_xml(all_text.strip()):
                    return "other"

                if self._is_positioned_like_footer_in_xml(xml_root):
                    return "other"

                return "content"

        return "other"

    def _get_slide_xml(self, slide):
        """Extract raw XML from slide with multiple access pattern attempts."""
        try:
            if hasattr(slide, '_element') and hasattr(slide._element, 'xml'):
                return slide._element.xml
            elif hasattr(slide, 'element') and hasattr(slide.element, 'xml'):
                return slide.element.xml
            else:
                slide_part = slide.part if hasattr(slide, 'part') else None
                if slide_part and hasattr(slide_part, '_element'):
                    return slide_part._element.xml
                else:
                    raise Exception("Cannot access slide XML")
        except Exception:
            return None

    def _get_shape_id(self, shape):
        """Extract PowerPoint's internal shape ID from python-pptx shape."""
        try:
            if hasattr(shape, '_element') and hasattr(shape._element, 'xml'):
                xml_str = shape._element.xml
                match = re.search(r'<[^>]*:cNvPr[^>]+id="([^"]+)"', xml_str)
                if match:
                    return match.group(1)
        except:
            pass
        return None

    def _get_shape_name(self, shape):
        """Extract shape name with error handling."""
        try:
            if hasattr(shape, 'name') and shape.name:
                return shape.name
        except:
            pass
        return None

    def _extract_shape_info_from_xml(self, shape_elem, order_index):
        """Extract identifying information from XML shape element."""
        shape_info = {
            'xml_order': order_index,
            'id': None,
            'name': None,
            'type': shape_elem.tag.split('}')[-1] if '}' in shape_elem.tag else shape_elem.tag,
            'has_text': False,
            'text_content': None
        }

        nv_props = (shape_elem.find('.//p:nvSpPr', self.namespaces) or
                    shape_elem.find('.//p:nvPicPr', self.namespaces) or
                    shape_elem.find('.//p:nvGraphicFramePr', self.namespaces) or
                    shape_elem.find('.//p:nvGrpSpPr', self.namespaces) or
                    shape_elem.find('.//p:nvCxnSpPr', self.namespaces))

        if nv_props is not None:
            cnv_pr = nv_props.find('.//p:cNvPr', self.namespaces)
            if cnv_pr is not None:
                shape_info['id'] = cnv_pr.get('id')
                shape_info['name'] = cnv_pr.get('name', '')

        text_elements = shape_elem.findall('.//a:t', self.namespaces)
        if text_elements:
            all_text = ' '.join([t.text for t in text_elements if t.text])
            if all_text.strip():
                shape_info['has_text'] = True
                shape_info['text_content'] = all_text.strip()[:50]

        return shape_info

    def _get_positional_ordered_shapes(self, slide):
        """Simple fallback: positional ordering (top-to-bottom, left-to-right)."""
        positioned_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                positioned_shapes.append((shape.top, shape.left, shape))
            else:
                positioned_shapes.append((0, 0, shape))

        positioned_shapes.sort(key=lambda x: (x[0], x[1]))
        return [shape for _, _, shape in positioned_shapes]

    def get_last_extraction_method(self):
        """Get the method used in the last extraction for debugging/monitoring."""
        return self.last_extraction_method

    def get_reading_order_of_grouped_by_shape(self, group_shape):
        """Extract reading order of shapes within a group using XML or z-axis order."""
        try:
            xml_ordered_children = self._get_group_xml_reading_order(group_shape)
            if xml_ordered_children:
                return xml_ordered_children
        except Exception as e:
            print(f"XML group reading order failed: {e}")

        return self._get_group_z_axis_order(group_shape)

    def _get_group_xml_reading_order(self, group_shape):
        """Extract child shapes from group XML in document order."""
        try:
            group_xml = group_shape._element.xml
            root = ET.fromstring(group_xml)

            child_elements = []
            for elem in root.iter():
                tag_name = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
                if tag_name in ['sp', 'pic', 'graphicFrame', 'cxnSp']:
                    child_info = self._extract_child_shape_info(elem)
                    if child_info:
                        child_elements.append(child_info)

            if child_elements:
                return self._map_xml_children_to_pptx_children(child_elements, group_shape.shapes)

        except Exception as e:
            print(f"Group XML parsing failed: {e}")

        return None

    def _get_group_z_axis_order(self, group_shape):
        """Get child shapes ordered by z-axis (stacking order) as fallback."""
        try:
            children_with_z_order = []

            for child_shape in group_shape.shapes:
                z_order = self._get_shape_z_order(child_shape)
                children_with_z_order.append((z_order, child_shape))

            children_with_z_order.sort(key=lambda x: x[0])
            return [shape for z_order, shape in children_with_z_order]

        except Exception as e:
            print(f"Z-axis ordering failed: {e}")
            return list(group_shape.shapes)

    def _get_shape_z_order(self, shape):
        """Extract z-order (stacking order) from shape XML with fallbacks."""
        try:
            if hasattr(shape, '_element') and hasattr(shape._element, 'xml'):
                xml_str = shape._element.xml

                z_order_match = re.search(r'z-?order["\s]*[:=]["\s]*(\d+)', xml_str, re.IGNORECASE)
                if z_order_match:
                    return int(z_order_match.group(1))

                id_match = re.search(r'id["\s]*=["\s]*["\'](\d+)["\']', xml_str)
                if id_match:
                    return int(id_match.group(1))

        except Exception:
            pass

        return 0

    def _extract_child_shape_info(self, shape_elem):
        """Extract information about a child shape from group XML."""
        child_info = {
            'id': None,
            'name': None,
            'type': shape_elem.tag.split('}')[-1] if '}' in shape_elem.tag else shape_elem.tag,
            'z_order': 0
        }

        nv_props = (shape_elem.find('.//p:nvSpPr', self.namespaces) or
                    shape_elem.find('.//p:nvPicPr', self.namespaces) or
                    shape_elem.find('.//p:nvGraphicFramePr', self.namespaces) or
                    shape_elem.find('.//p:nvCxnSpPr', self.namespaces))

        if nv_props is not None:
            cnv_pr = nv_props.find('.//p:cNvPr', self.namespaces)
            if cnv_pr is not None:
                child_info['id'] = cnv_pr.get('id')
                child_info['name'] = cnv_pr.get('name', '')

        return child_info

    def _map_xml_children_to_pptx_children(self, xml_children, pptx_children):
        """Map XML child shape info to python-pptx child shapes."""
        ordered_children = []
        used_child_ids = set()

        child_lookup = {}
        for child in pptx_children:
            child_id = self._get_shape_id(child)
            child_name = self._get_shape_name(child)

            if child_id:
                child_lookup[f"id_{child_id}"] = child
            if child_name:
                child_lookup[f"name_{child_name}"] = child

        for xml_child in xml_children:
            matched_child = None

            if xml_child['id']:
                matched_child = child_lookup.get(f"id_{xml_child['id']}")

            if not matched_child and xml_child['name']:
                matched_child = child_lookup.get(f"name_{xml_child['name']}")

            if matched_child:
                child_obj_id = id(matched_child)
                if child_obj_id not in used_child_ids:
                    ordered_children.append(matched_child)
                    used_child_ids.add(child_obj_id)

        for child in pptx_children:
            child_obj_id = id(child)
            if child_obj_id not in used_child_ids:
                ordered_children.append(child)

        return ordered_children

    # Add all missing helper methods for completeness
    def _validate_title_from_xml_content(self, xml_root):
        """Validate title candidate using XML text content analysis."""
        try:
            text_elements = xml_root.findall('.//a:t', self.namespaces)
            if text_elements:
                all_text = ' '.join([elem.text for elem in text_elements if elem.text])
                text_stripped = all_text.strip()

                if text_stripped.isdigit():
                    return False

                if len(text_stripped) <= 3 and any(char.isdigit() for char in text_stripped):
                    return False

            if self._is_positioned_like_footer_in_xml(xml_root):
                return False

            if self._is_too_small_for_title_in_xml(xml_root):
                return False

            return True

        except Exception:
            return True

    def _is_slide_number_pattern_in_xml(self, text_content):
        """Check if XML text content matches slide number patterns."""
        text_lower = text_content.lower().strip()

        slide_number_patterns = [
            r'^\d+$',
            r'^\d+\s*/\s*\d+$',
            r'^slide\s+\d+$',
            r'^page\s+\d+$',
            r'^\d+\s+of\s+\d+$',
        ]

        for pattern in slide_number_patterns:
            if re.match(pattern, text_lower):
                return True

        return False

    def _is_positioned_like_footer_in_xml(self, xml_root):
        """Check XML positioning attributes to identify footer-like placement."""
        try:
            xfrm_elements = xml_root.findall('.//a:xfrm', self.namespaces)

            for xfrm in xfrm_elements:
                off_elem = xfrm.find('a:off', self.namespaces)
                ext_elem = xfrm.find('a:ext', self.namespaces)

                if off_elem is not None and ext_elem is not None:
                    y_pos = int(off_elem.get('y', 0))
                    height = int(ext_elem.get('cy', 0))

                    slide_height = 6858000
                    relative_position = (y_pos + height) / slide_height

                    if relative_position > 0.85:
                        return True

            return False

        except Exception:
            return False

    def _is_too_small_for_title_in_xml(self, xml_root):
        """Check XML size attributes to exclude very small shapes from being titles."""
        try:
            xfrm_elements = xml_root.findall('.//a:xfrm', self.namespaces)

            for xfrm in xfrm_elements:
                ext_elem = xfrm.find('a:ext', self.namespaces)

                if ext_elem is not None:
                    height = int(ext_elem.get('cy', 0))
                    min_title_height = 254000

                    if height < min_title_height:
                        return True

            return False

        except Exception:
            return False