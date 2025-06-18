"""
Enhanced Diagram Analyzer - v19 Scoring System with Direct Shape Access
Analyzes slide content to identify potential diagrams with probability scoring

ENHANCEMENT OVERVIEW:
This enhanced version can extract its own diagram-relevant data directly from
PowerPoint slide objects, bypassing content extractor filtering that may remove
shapes crucial for diagram analysis (lines, connectors, basic shapes, etc.).

The core v19 scoring algorithm remains unchanged - only the data extraction
has been enhanced to ensure all diagram-relevant shapes are captured.

DUAL ANALYSIS MODES:
1. Structured Data Mode: Uses pre-processed content from ContentExtractor
2. Direct Slide Mode: Extracts shapes directly from PowerPoint slide objects

The analyzer automatically chooses the best mode and falls back gracefully.
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
import re


class DiagramAnalyzer:
    """
    Analyzes slide content to identify potential diagrams with enhanced shape extraction.

    COMPONENT RESPONSIBILITIES:
    - Score individual slides for diagram probability using v19 system
    - Extract diagram-relevant shapes directly from slides when needed
    - Categorize slide elements (shapes, lines, arrows, text) comprehensively
    - Analyze spatial layouts and flow patterns
    - Generate detailed analysis reports with reasoning
    - Convert raw scores to meaningful probability percentages

    ENHANCED FEATURES:
    - Direct PowerPoint slide object analysis
    - Comprehensive shape extraction bypassing content filtering
    - Fallback to structured data when direct access unavailable
    - Preserves all original v19 scoring logic
    """

    def analyze_slides_for_diagrams(self, slides=None, structured_data=None):
        """
        Enhanced main analysis entry point supporting both direct slides and structured data.

        ANALYSIS STRATEGY:
        1. If PowerPoint slide objects available: Use direct extraction for complete shape data
        2. If only structured data available: Use existing structured data analysis
        3. Compare results and use the most comprehensive analysis

        SHAPE COMPLETENESS:
        Direct slide analysis captures ALL shapes including:
        - Lines and connectors (often filtered by content extractor)
        - Basic geometric shapes without text
        - Arrows and flow indicators
        - Empty text boxes and placeholders

        These shapes are crucial for diagram detection but may be filtered
        out by content extraction for markdown generation.

        Args:
            slides (list): List of python-pptx Slide objects (preferred)
            structured_data (dict): Pre-processed structured data (fallback)

        Returns:
            str: Diagram analysis summary or None if no diagrams found
        """
        try:
            diagram_slides = []

            # Method 1: Direct slide analysis (preferred for comprehensive shape data)
            if slides:
                print("🎯 Using direct slide analysis for comprehensive diagram detection")
                for slide_idx, slide in enumerate(slides):
                    slide_data = self._extract_slide_data_for_diagram_analysis(slide)
                    score_analysis = self.score_slide_for_diagram(slide_data)

                    if score_analysis["probability"] >= 40:  # 40%+ probability threshold
                        diagram_slides.append({
                            "slide": slide_idx + 1,
                            "analysis": score_analysis,
                            "method": "direct_slide_analysis"
                        })

            # Method 2: Structured data analysis (fallback)
            elif structured_data:
                print("📄 Using structured data analysis (some shapes may be filtered)")
                for slide_idx, slide in enumerate(structured_data["slides"]):
                    score_analysis = self.score_slide_for_diagram(slide)

                    if score_analysis["probability"] >= 40:
                        diagram_slides.append({
                            "slide": slide_idx + 1,
                            "analysis": score_analysis,
                            "method": "structured_data_analysis"
                        })
            else:
                print("❌ No slide data provided for diagram analysis")
                return None

            # Generate detailed summary if diagrams found
            if diagram_slides:
                summary = "## DIAGRAM ANALYSIS (v19 Enhanced Scoring System)\n\n"
                summary += "**Slides with potential diagrams:**\n\n"

                for slide_info in diagram_slides:
                    analysis = slide_info["analysis"]
                    method = slide_info.get("method", "unknown")

                    summary += f"- **Slide {slide_info['slide']}**: {analysis['probability']}% probability "
                    summary += f"(Score: {analysis['total_score']}) - {', '.join(analysis['reasons'])}\n"
                    summary += f"  - Shapes: {analysis['shape_count']}, Lines: {analysis['line_count']}, Arrows: {analysis['arrow_count']}\n"
                    summary += f"  - Analysis method: {method}\n\n"

                return summary

            return None

        except Exception as e:
            # Analysis errors shouldn't stop processing - return error comment
            return f"\n\n<!-- Enhanced v19 Diagram analysis error: {e} -->"

    def _extract_slide_data_for_diagram_analysis(self, slide):
        """
        Extract comprehensive slide data directly from PowerPoint slide object.

        COMPREHENSIVE EXTRACTION:
        This method captures ALL shapes on the slide, including those that
        might be filtered out by content extraction for markdown generation:
        - Lines and connectors without text
        - Basic geometric shapes (rectangles, circles, etc.)
        - Arrows and directional indicators
        - Empty or placeholder text boxes
        - Grouped shapes (extracted individually)

        SHAPE PROCESSING:
        1. Iterate through all slide.shapes
        2. Expand groups to process individual components
        3. Categorize each shape by type and characteristics
        4. Extract minimal text content for flow analysis
        5. Preserve positional information for layout analysis

        Args:
            slide: python-pptx Slide object

        Returns:
            dict: Slide data optimized for diagram analysis
        """
        slide_data = {"content_blocks": []}

        try:
            # Process all shapes on the slide comprehensively
            all_shapes = self._get_all_shapes_including_groups(slide)
            print(f"DEBUG: Direct extraction found {len(all_shapes)} total shapes")

            for shape in all_shapes:
                content_block = self._create_diagram_content_block(shape)
                if content_block:
                    slide_data["content_blocks"].append(content_block)

        except Exception as e:
            print(f"Warning: Error in direct slide extraction: {e}")

        return slide_data

    def _get_all_shapes_including_groups(self, slide):
        """
        Extract all individual shapes from slide, expanding groups.

        GROUP EXPANSION:
        PowerPoint groups can hide important diagram elements like
        connecting lines and arrows. This method ensures all individual
        shapes are captured for analysis.

        Args:
            slide: python-pptx Slide object

        Returns:
            list: All individual shapes (groups expanded)
        """
        all_shapes = []

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                # Expand group to get individual shapes
                try:
                    group_shapes = list(shape.shapes)
                    all_shapes.extend(group_shapes)
                    print(f"DEBUG: Expanded group with {len(group_shapes)} shapes")
                except:
                    # Group expansion failed, treat as single shape
                    all_shapes.append(shape)
            else:
                all_shapes.append(shape)

        return all_shapes

    def _create_diagram_content_block(self, shape):
        """
        Create content block optimized for diagram analysis.

        DIAGRAM-SPECIFIC EXTRACTION:
        Unlike content extraction for markdown, this focuses on
        capturing all information relevant to diagram detection:
        - Shape type and characteristics
        - Positional information for layout analysis
        - Basic text content for flow keyword detection
        - Directional indicators (arrows, connectors)

        DEFENSIVE PROCESSING:
        All shape property access is wrapped in try/catch to handle
        various PowerPoint file formats and shape types gracefully.

        Args:
            shape: python-pptx Shape object

        Returns:
            dict: Content block optimized for diagram analysis
        """
        try:
            shape_type = shape.shape_type
            shape_type_name = str(shape_type).split('.')[-1] if hasattr(shape_type, '__str__') else 'unknown'

            # Basic content block structure
            content_block = {
                "type": self._determine_diagram_type(shape, shape_type_name),
                "position": self._extract_position_info(shape),
                "shape_info": {
                    "shape_type": shape_type_name,
                    "auto_shape_type": self._get_auto_shape_type(shape)
                }
            }

            # Extract text content if present (important for flow analysis)
            text_content = self._extract_basic_text_content(shape)
            if text_content:
                content_block["text_content"] = text_content
                content_block["paragraphs"] = [{"clean_text": text_content}]

            return content_block

        except Exception as e:
            print(f"DEBUG: Error creating diagram content block: {e}")
            return None

    # Replace the _determine_diagram_type method in your DiagramAnalyzer class with this:

    def _determine_diagram_type(self, shape, shape_type_name):
        """
        Determine content block type for diagram analysis.
        ENHANCED DEBUG: Show exactly why LINE shapes aren't being detected.
        """
        try:
            print(f"    -> DEBUG: Input shape_type_name = '{shape_type_name}'")

            # Extract just the shape type without numbers/parentheses
            clean_shape_type = shape_type_name.split('(')[0].strip() if '(' in shape_type_name else shape_type_name
            print(f"    -> DEBUG: Cleaned shape type = '{clean_shape_type}'")

            # Lines and connectors (crucial for diagram detection)
            if clean_shape_type in ['LINE', 'CONNECTOR', 'FREEFORM']:
                print(f"    -> ✅ MATCHED LINE PATTERN! Categorized as 'line'")
                return "line"

            # Arrows (directional flow indicators)
            elif clean_shape_type == 'AUTO_SHAPE':
                auto_shape_type = self._get_auto_shape_type(shape)
                print(f"    -> DEBUG: AUTO_SHAPE type = '{auto_shape_type}'")
                if self._is_arrow_shape(auto_shape_type):
                    print(f"    -> ✅ MATCHED ARROW PATTERN! Categorized as 'arrow'")
                    return "arrow"
                else:
                    print(f"    -> AUTO_SHAPE not an arrow, categorized as 'shape'")
                    return "shape"

            # Charts (often represent diagrams)
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                print(f"    -> Categorized as 'chart'")
                return "chart"

            # Tables (structured data representation)
            elif clean_shape_type == 'TABLE':
                print(f"    -> Categorized as 'table'")
                return "table"

            # Text elements (for flow keyword analysis)
            elif hasattr(shape, 'text_frame') or hasattr(shape, 'text'):
                print(f"    -> Categorized as 'text'")
                return "text"

            # Everything else as generic shape
            else:
                print(f"    -> ❌ NO PATTERN MATCHED! Categorized as generic 'shape'")
                return "shape"

        except Exception as e:
            print(f"    -> ❌ Error in categorization: {e}")
            return "shape"



    def _extract_position_info(self, shape):
        """
        Extract positional information for layout analysis.

        LAYOUT ANALYSIS IMPORTANCE:
        Position data is crucial for detecting grid layouts,
        spatial distributions, and organized arrangements
        that indicate diagram content.

        Args:
            shape: python-pptx Shape object

        Returns:
            dict: Position and size information
        """
        try:
            return {
                "top": getattr(shape, 'top', 0),
                "left": getattr(shape, 'left', 0),
                "width": getattr(shape, 'width', 0),
                "height": getattr(shape, 'height', 0)
            }
        except:
            return {"top": 0, "left": 0, "width": 0, "height": 0}

    def _get_auto_shape_type(self, shape):
        """
        Get auto shape type for arrow detection.

        Args:
            shape: python-pptx Shape object

        Returns:
            str: Auto shape type or None
        """
        try:
            if hasattr(shape, 'auto_shape_type'):
                auto_shape_type = shape.auto_shape_type
                return str(auto_shape_type).split('.')[-1] if hasattr(auto_shape_type, '__str__') else None
        except:
            pass
        return None

    def _extract_basic_text_content(self, shape):
        """
        Extract basic text content for flow keyword analysis.

        FLOW ANALYSIS PURPOSE:
        Text content is analyzed for workflow vocabulary
        (process, step, decision, etc.) which indicates
        diagram content even when shapes are basic.

        Args:
            shape: python-pptx Shape object

        Returns:
            str: Basic text content or None
        """
        try:
            # Try text_frame first (most common)
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    return text

            # Try direct text attribute
            elif hasattr(shape, 'text'):
                text = shape.text.strip()
                if text:
                    return text

        except:
            pass
        return None

    # ORIGINAL v19 SCORING METHODS (UNCHANGED)
    # ==========================================

    def analyze_structured_data_for_diagrams(self, structured_data):
        """
        Original analysis method for backward compatibility.

        This method preserves the original interface for existing code
        while the enhanced analyze_slides_for_diagrams provides better results.
        """
        return self.analyze_slides_for_diagrams(structured_data=structured_data)

    def score_slide_for_diagram(self, slide_data):
        """
        Core scoring algorithm - implements v19 sophisticated rules.

        *** UNCHANGED FROM ORIGINAL - PRESERVES WORKING LOGIC ***
        """
        content_blocks = slide_data.get("content_blocks", [])

        # Categorize all slide elements for analysis
        shapes, lines, arrows, text_blocks = self._categorize_slide_elements(content_blocks)

        # Initialize scoring system
        score = 0
        reasons = []

        # Rule 1: Arrow/Line threshold analysis (20+ points each)
        if len(arrows) > 0:
            score += 20
            reasons.append(f"block_arrows:{len(arrows)}")

        if len(lines) >= 3:
            score += 20
            reasons.append(f"connector_lines:{len(lines)}")

        # Rule 2: Line-to-shape ratio analysis (15 points)
        total_lines = len(lines) + len(arrows)
        if len(shapes) > 0:
            line_ratio = total_lines / len(shapes)
            if line_ratio >= 0.5:  # 50%+ lines relative to shapes
                score += 15
                reasons.append(f"line_ratio:{line_ratio:.1f}")

        # Rule 3: Spatial layout analysis (10-15 points)
        layout_score = self._analyze_spatial_layout(shapes)
        score += layout_score["score"]
        if layout_score["score"] > 0:
            reasons.append(f"layout:{layout_score['type']}")

        # Rule 4: Shape variety analysis (10-15 points)
        variety_score = self._analyze_shape_variety(shapes)
        score += variety_score
        if variety_score > 0:
            reasons.append(f"variety:{variety_score}")

        # Rule 5: Text density analysis (10 points)
        text_score = self._analyze_text_density(text_blocks)
        score += text_score
        if text_score > 0:
            reasons.append(f"short_text:{text_score}")

        # Rule 6: Flow pattern analysis (20 points)
        flow_score = self._analyze_flow_patterns(shapes, lines, arrows, text_blocks)
        score += flow_score
        if flow_score > 0:
            reasons.append(f"flow_pattern:{flow_score}")

        # Rule 7: Negative indicators (subtract points)
        negative_score = self._analyze_negative_indicators(text_blocks, shapes)
        score += negative_score  # negative_score will be ≤ 0
        if negative_score < 0:
            reasons.append(f"negatives:{negative_score}")

        # Convert raw score to probability percentage
        probability = self._calculate_probability_from_score(score)

        return {
            "total_score": score,
            "probability": probability,
            "reasons": reasons,
            "shape_count": len(shapes),
            "line_count": len(lines),
            "arrow_count": len(arrows)
        }

    def _categorize_slide_elements(self, content_blocks):
        """*** UNCHANGED FROM ORIGINAL ***"""
        shapes = []
        lines = []
        arrows = []
        text_blocks = []

        for block in content_blocks:
            block_type = block.get("type")

            if block_type == "line":
                lines.append(block)
            elif block_type == "arrow":
                arrows.append(block)
            elif block_type == "text":
                text_blocks.append(block)
                shapes.append(block)  # Text boxes are also shapes
            elif block_type in ["shape", "image", "chart"]:
                shapes.append(block)
            elif block_type == "group":
                # Recursively process group contents
                group_analysis = self._analyze_group_contents(block)
                shapes.extend(group_analysis["shapes"])
                lines.extend(group_analysis["lines"])
                arrows.extend(group_analysis["arrows"])
                text_blocks.extend(group_analysis["text_blocks"])

        return shapes, lines, arrows, text_blocks

    def _analyze_group_contents(self, group_block):
        """*** UNCHANGED FROM ORIGINAL ***"""
        result = {"shapes": [], "lines": [], "arrows": [], "text_blocks": []}

        for extracted_block in group_block.get("extracted_blocks", []):
            block_type = extracted_block.get("type")

            if block_type == "line":
                result["lines"].append(extracted_block)
            elif block_type == "arrow":
                result["arrows"].append(extracted_block)
            elif block_type == "text":
                result["text_blocks"].append(extracted_block)
                result["shapes"].append(extracted_block)
            elif block_type in ["shape", "image", "chart"]:
                result["shapes"].append(extracted_block)

        return result

    def _analyze_spatial_layout(self, shapes):
        """*** UNCHANGED FROM ORIGINAL ***"""
        if len(shapes) < 3:
            return {"score": 0, "type": "insufficient"}

        positions = []
        for shape in shapes:
            pos = shape.get("position")
            if pos:
                positions.append((pos["top"], pos["left"]))

        if len(positions) < 3:
            return {"score": 0, "type": "no_position_data"}

        # Calculate spatial distribution metrics
        tops = [p[0] for p in positions]
        lefts = [p[1] for p in positions]

        top_range = max(tops) - min(tops) if tops else 0
        left_range = max(lefts) - min(lefts) if lefts else 0

        # Analyze alignment patterns
        # Group positions by approximate location (100K EMU tolerance)
        unique_tops = len(set(round(t / 100000) for t in tops))
        unique_lefts = len(set(round(l / 100000) for l in lefts))

        # Classify layout pattern
        if unique_tops >= 2 and unique_lefts >= 2:
            # Grid-like arrangement: multiple rows and columns
            return {"score": 15, "type": "grid_layout"}
        elif top_range > 1000000 and left_range > 1000000:
            # Spread arrangement: wide spatial distribution
            return {"score": 10, "type": "spread_layout"}
        else:
            # Linear arrangement: single row or column
            return {"score": 0, "type": "linear_layout"}

    def _analyze_shape_variety(self, shapes):
        """*** UNCHANGED FROM ORIGINAL ***"""
        if len(shapes) < 2:
            return 0

        shape_types = set()
        sizes = []

        for shape in shapes:
            shape_types.add(shape.get("type", "unknown"))
            pos = shape.get("position")
            if pos:
                size = pos["width"] * pos["height"]
                sizes.append(size)

        score = 0

        # Shape type diversity scoring
        if len(shape_types) >= 3:
            score += 15  # High diversity
        elif len(shape_types) >= 2:
            score += 10  # Moderate diversity

        # Size consistency analysis (process flow indicator)
        if len(sizes) >= 3:
            avg_size = sum(sizes) / len(sizes)
            if avg_size > 0:
                # Calculate coefficient of variation
                variations = [abs(size - avg_size) / avg_size for size in sizes]
                if variations and max(variations) < 0.5:  # <50% variation
                    score += 5  # Consistent sizing bonus

        return score

    def _analyze_text_density(self, text_blocks):
        """*** UNCHANGED FROM ORIGINAL ***"""
        if not text_blocks:
            return 0

        short_text_count = 0
        total_blocks = len(text_blocks)

        for block in text_blocks:
            # Analyze average words per paragraph in this block
            total_words = 0
            para_count = 0

            for para in block.get("paragraphs", []):
                clean_text = para.get("clean_text", "")
                if clean_text:
                    words = len(clean_text.split())
                    total_words += words
                    para_count += 1

            # Classify block based on average paragraph length
            if para_count > 0:
                avg_words = total_words / para_count
                if avg_words <= 5:  # Short labels (≤5 words per paragraph)
                    short_text_count += 1

        # Score based on percentage of short text blocks
        if total_blocks > 0:
            short_ratio = short_text_count / total_blocks
            if short_ratio >= 0.7:  # 70%+ short text
                return 10
            elif short_ratio >= 0.5:  # 50%+ short text
                return 5

        return 0

    def _analyze_flow_patterns(self, shapes, lines, arrows, text_blocks):
        """*** UNCHANGED FROM ORIGINAL ***"""
        score = 0

        # Define workflow vocabulary
        flow_keywords = ["start", "begin", "end", "finish", "process", "step", "decision"]
        action_words = ["create", "update", "check", "verify", "send", "receive", "analyze"]

        # Extract all text content for analysis
        all_text = ""
        for block in text_blocks:
            for para in block.get("paragraphs", []):
                all_text += " " + para.get("clean_text", "").lower()

        # Count keyword matches
        flow_matches = sum(1 for keyword in flow_keywords if keyword in all_text)
        action_matches = sum(1 for keyword in action_words if keyword in all_text)

        # Score based on workflow vocabulary
        if flow_matches >= 2:
            score += 20  # Strong workflow indicator
        elif flow_matches >= 1:
            score += 10  # Moderate workflow indicator

        if action_matches >= 3:
            score += 10  # Action-heavy content

        # Bonus for structural flow indicators
        if len(shapes) >= 3 and (len(lines) > 0 or len(arrows) > 0):
            score += 15  # Shapes connected by lines/arrows

        return score

    def _analyze_negative_indicators(self, text_blocks, shapes):
        """*** UNCHANGED FROM ORIGINAL ***"""
        score = 0

        # Analyze text characteristics for document patterns
        long_text_count = 0
        bullet_count = 0
        total_paras = 0

        for block in text_blocks:
            for para in block.get("paragraphs", []):
                clean_text = para.get("clean_text", "")
                if clean_text:
                    total_paras += 1
                    word_count = len(clean_text.split())

                    # Count long paragraphs (document indicator)
                    if word_count > 20:
                        long_text_count += 1

                    # Count bullet points (list indicator)
                    if para.get("hints", {}).get("is_bullet", False):
                        bullet_count += 1

        # Penalize document-style long text
        if long_text_count >= 2:
            score -= 15

        # Penalize bullet-heavy content
        if total_paras > 0 and bullet_count / total_paras > 0.8:
            score -= 10  # 80%+ bullet points

        # Analyze layout for single-column arrangement
        if len(shapes) >= 3:
            positions = [s.get("position") for s in shapes if s.get("position")]
            if len(positions) >= 3:
                lefts = [p["left"] for p in positions]
                left_variance = max(lefts) - min(lefts) if lefts else 0
                # Very narrow horizontal spread suggests single column
                if left_variance < 500000:  # 500K EMU threshold
                    score -= 10

        return score

    def _calculate_probability_from_score(self, score):
        """*** UNCHANGED FROM ORIGINAL ***"""
        if score >= 60:
            return 95  # Very high confidence
        elif score >= 40:
            return 75  # High confidence
        elif score >= 20:
            return 40  # Moderate confidence (threshold)
        else:
            return 10  # Low confidence

    def _is_arrow_shape(self, auto_shape_type):
        """Enhanced arrow detection with debugging."""
        if not auto_shape_type:
            print(f"      -> No auto_shape_type provided")
            return False

        try:
            auto_shape_str = str(auto_shape_type).upper()
            print(f"      -> Testing arrow pattern against: '{auto_shape_str}'")

            arrow_types = [
                "LEFT_ARROW", "DOWN_ARROW", "UP_ARROW", "RIGHT_ARROW",
                "LEFT_RIGHT_ARROW", "UP_DOWN_ARROW", "QUAD_ARROW",
                "LEFT_RIGHT_UP_ARROW", "BENT_ARROW", "U_TURN_ARROW",
                "CURVED_LEFT_ARROW", "CURVED_RIGHT_ARROW",
                "CURVED_UP_ARROW", "CURVED_DOWN_ARROW",
                "STRIPED_RIGHT_ARROW", "NOTCHED_RIGHT_ARROW",
                "BLOCK_ARC"
            ]

            for arrow_type in arrow_types:
                if arrow_type in auto_shape_str:
                    print(f"      -> ✅ Matched arrow type: {arrow_type}")
                    return True

            print(f"      -> ❌ No arrow pattern matched")
            return False

        except Exception as e:
            print(f"      -> ❌ Error checking arrow shape: {e}")
            return False