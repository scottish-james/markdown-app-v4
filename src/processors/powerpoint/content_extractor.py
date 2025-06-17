"""
Content Extractor - ENHANCED: Now preserves meaningful alt text from images
Updated to check for meaningful alt text before ignoring shapes
"""

from pptx.enum.shapes import MSO_SHAPE_TYPE
import xml.etree.ElementTree as ET
import re


class ContentExtractor:
    """
    Extracts content from various PowerPoint shape types with semantic role preservation.
    ENHANCED: Now checks for meaningful alt text before deciding to ignore shapes.
    """

    def extract_shape_content(self, shape, text_processor, accessibility_extractor=None, groups_already_expanded=False):
        """
        Main extraction router - delegates based on shape type and captures semantic role.
        ENHANCED: Now checks for meaningful alt text before ignoring shapes.

        Args:
            groups_already_expanded: If True, skip group processing as shapes already expanded
        """
        # Capture basic shape info for diagram analysis
        shape_info = self._get_shape_analysis_info(shape)

        # Capture semantic role from XML analysis
        semantic_role = "other"
        if accessibility_extractor:
            try:
                semantic_role = accessibility_extractor._get_semantic_role_from_xml(shape)
            except Exception as e:
                print(f"DEBUG: Error getting semantic role: {e}")

        content_block = None

        try:
            # Get shape type safely
            shape_type = shape.shape_type
            shape_type_name = str(shape_type).split('.')[-1] if hasattr(shape_type, '__str__') else 'unknown'

            # Route based on shape type using string comparison for safety
            if shape_type_name == 'PICTURE':
                content_block = self.extract_image(shape)
            elif shape_type_name == 'TABLE':
                content_block = self.extract_table(shape.table, text_processor)
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                content_block = self.extract_chart(shape)
            elif shape_type_name == 'GROUP':
                # CRITICAL FIX: Only process groups if they haven't been expanded already
                if not groups_already_expanded:
                    content_block = self.extract_group(shape, text_processor, accessibility_extractor)
                else:
                    # Groups already expanded by accessibility extractor - skip to avoid double processing
                    print(f"DEBUG: Skipping group '{getattr(shape, 'name', 'unnamed')}' - already expanded")
                    return None
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                content_block = text_processor.extract_text_frame(shape.text_frame, shape)
            elif hasattr(shape, 'text') and shape.text:
                content_block = text_processor.extract_plain_text(shape)
        except Exception as e:
            print(f"Warning: Error extracting shape content: {e}")
            return None

        # Handle shapes without text content - but preserve meaningful alt text
        if not content_block:
            # NEW: Check for meaningful alt text first - this takes precedence
            if self._has_meaningful_alt_text(shape):
                print(f"DEBUG: Shape has meaningful alt text - extracting as image")
                content_block = self.extract_image(shape)
            elif self._is_meaningful_non_text_shape(shape, shape_info):
                content_block = self._create_non_text_content_block(shape, shape_info)
            else:
                # Skip meaningless shapes (lines, basic auto-shapes, etc.)
                return None

        # Add shape analysis info and semantic role for downstream processing
        if content_block:
            try:
                content_block.update(shape_info)
                content_block["semantic_role"] = semantic_role
            except Exception as e:
                print(f"Warning: Error adding shape info: {e}")

        return content_block

    def _has_meaningful_alt_text(self, shape):
        """
        Check if shape has meaningful alt text that's worth preserving.

        MEANINGFUL ALT TEXT CRITERIA:
        - Not empty or just whitespace
        - Not generic like "Image", "Picture", "image1.png", etc.
        - Not just numbers or short meaningless strings
        - Actually describes the content

        Args:
            shape: python-pptx Shape object

        Returns:
            bool: True if shape has meaningful alt text
        """
        alt_text = self._extract_alt_text_from_shape(shape)

        if not alt_text or not alt_text.strip():
            return False

        alt_text = alt_text.strip()

        # Check for generic/meaningless alt text patterns
        meaningless_patterns = [
            r'^image\d*\.?(png|jpg|jpeg|gif|bmp|svg|webp)?$',  # image123.png, image.jpg, etc.
            r'^picture\d*$',  # picture, picture1, etc.
            r'^img\d*$',  # img, img1, etc.
            r'^graphic\d*$',  # graphic, graphic1, etc.
            r'^shape\d*$',  # shape, shape1, etc.
            r'^slide\d+image\d*$',  # slide1image1, etc.
            r'^\d+$',  # just numbers
            r'^[a-z]{1,3}$',  # very short generic strings
        ]

        # Check against meaningless patterns (case insensitive)
        alt_text_lower = alt_text.lower()
        for pattern in meaningless_patterns:
            if re.match(pattern, alt_text_lower):
                print(f"DEBUG: Alt text '{alt_text}' matches meaningless pattern '{pattern}'")
                return False

        # Check for very short text that's likely meaningless
        if len(alt_text) < 3:
            print(f"DEBUG: Alt text '{alt_text}' too short to be meaningful")
            return False

        # Check for generic words that suggest auto-generated content
        generic_words = ['image', 'picture', 'graphic', 'shape', 'photo', 'diagram']
        if alt_text_lower in generic_words:
            print(f"DEBUG: Alt text '{alt_text}' is generic word")
            return False

        # If we get here, it's likely meaningful
        print(f"DEBUG: Alt text '{alt_text}' appears meaningful")
        return True

    def _extract_alt_text_from_shape(self, shape):
        """
        Extract alt text from shape using multiple methods.
        This is similar to extract_image but just returns the alt text string.

        Args:
            shape: python-pptx Shape object

        Returns:
            str: Alt text or None if not found
        """
        try:
            # Method 1: Direct alt_text attribute
            if hasattr(shape, 'alt_text') and shape.alt_text:
                return shape.alt_text.strip()

            # Method 2: Image alt_text
            if hasattr(shape, 'image') and hasattr(shape.image, 'alt_text') and shape.image.alt_text:
                return shape.image.alt_text.strip()

            # Method 3: XML extraction
            if hasattr(shape, '_element'):
                try:
                    xml_str = str(shape._element.xml) if hasattr(shape._element, 'xml') else ""
                    if xml_str:
                        root = ET.fromstring(xml_str)
                        for elem in root.iter():
                            # Check description attribute
                            if 'descr' in elem.attrib and elem.attrib['descr']:
                                return elem.attrib['descr'].strip()
                            # Check title attribute
                            elif 'title' in elem.attrib and elem.attrib['title']:
                                return elem.attrib['title'].strip()
                except:
                    pass
        except:
            pass

        return None

    def _is_meaningful_non_text_shape(self, shape, shape_info):
        """
        Determine if a non-text shape is worth including in output.
        FIXED: More defensive shape type checking to avoid enum errors.

        Args:
            shape: python-pptx Shape object
            shape_info: Shape analysis info dict

        Returns:
            bool: True if shape should be included in output
        """
        try:
            shape_type = shape.shape_type
            shape_type_name = str(shape_type).split('.')[-1] if hasattr(shape_type, '__str__') else 'unknown'

            # Always include images and charts (already handled above)
            if shape_type_name in ['PICTURE', 'TABLE']:
                return True

            # Skip basic lines and connectors unless they're part of a larger diagram
            if shape_type_name in ['LINE', 'CONNECTOR', 'FREEFORM']:
                return False  # Usually just decorative

            # Include meaningful auto-shapes (arrows, but not basic rectangles/circles)
            if shape_type_name == 'AUTO_SHAPE':
                auto_shape_type = shape_info.get("auto_shape_type", "")
                # Include arrows as they often indicate flow/relationships
                if self._is_arrow_shape(auto_shape_type):
                    return True
                # Skip basic geometric shapes
                elif any(basic_type in str(auto_shape_type) for basic_type in [
                    "RECTANGLE", "OVAL", "CIRCLE", "TRIANGLE", "DIAMOND", "HEXAGON"
                ]):
                    return False
                # Include other auto-shapes (might be meaningful icons/symbols)
                else:
                    return True

            # Include other shape types by default (might be meaningful)
            return True

        except Exception as e:
            print(f"DEBUG: Error checking shape meaningfulness: {e}")
            # If we can't determine, err on the side of inclusion
            return True

    def extract_group(self, shape, text_processor, accessibility_extractor=None):
        """
        Extract content from grouped shapes using proper ordering.
        Only called when groups haven't been pre-expanded by accessibility extractor.
        """
        try:
            extracted_blocks = []

            # Use accessibility extractor for proper ordering if available
            if accessibility_extractor:
                print(f"DEBUG: Using accessibility extractor for group '{getattr(shape, 'name', 'unnamed')}' ordering")
                ordered_children = accessibility_extractor.get_reading_order_of_grouped_by_shape(shape)
            else:
                print(f"DEBUG: No accessibility extractor - using default group order")
                ordered_children = list(shape.shapes)

            print(f"DEBUG: Group has {len(ordered_children)} children")

            # Process each child shape
            for child_shape in ordered_children:
                # Recursively extract content from each child
                child_block = self.extract_shape_content(
                    child_shape,
                    text_processor,
                    accessibility_extractor,
                    groups_already_expanded=False  # Child shapes not pre-expanded
                )
                if child_block:
                    extracted_blocks.append(child_block)

            # Return group container if any content was extracted
            if extracted_blocks:
                print(
                    f"DEBUG: Group '{getattr(shape, 'name', 'unnamed')}' produced {len(extracted_blocks)} content blocks")
                return {
                    "type": "group",
                    "extracted_blocks": extracted_blocks,
                    "hyperlink": self._extract_shape_hyperlink(shape),
                    "semantic_role": "group"
                }

            print(f"DEBUG: Group '{getattr(shape, 'name', 'unnamed')}' produced no content")
            return None

        except Exception as e:
            print(f"Error extracting group: {e}")
            return None

    def extract_image(self, shape):
        """Extract image information with comprehensive alt text detection."""
        alt_text = "Image"

        try:
            if hasattr(shape, 'alt_text') and shape.alt_text:
                alt_text = shape.alt_text
            elif hasattr(shape, 'image') and hasattr(shape.image, 'alt_text') and shape.image.alt_text:
                alt_text = shape.image.alt_text
            elif hasattr(shape, '_element'):
                try:
                    xml_str = str(shape._element.xml) if hasattr(shape._element, 'xml') else ""
                    if xml_str:
                        root = ET.fromstring(xml_str)
                        for elem in root.iter():
                            if 'descr' in elem.attrib and elem.attrib['descr']:
                                alt_text = elem.attrib['descr']
                                break
                            elif 'title' in elem.attrib and elem.attrib['title']:
                                alt_text = elem.attrib['title']
                                break
                except:
                    pass
        except:
            pass

        return {
            "type": "image",
            "alt_text": alt_text.strip() if alt_text else "Image",
            "hyperlink": self._extract_shape_hyperlink(shape)
        }

    def extract_table(self, table, text_processor):
        """Extract table data with cell-level text processing."""
        if not table.rows:
            return None

        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                cell_content = ""

                if hasattr(cell, 'text_frame') and cell.text_frame:
                    cell_paras = []
                    for para in cell.text_frame.paragraphs:
                        if para.text.strip():
                            para_processed = text_processor.process_paragraph(para)
                            if para_processed and para_processed['hints']['is_bullet']:
                                level = para_processed['hints']['bullet_level']
                                indent = "  " * level
                                cell_paras.append(f"{indent}• {para_processed['clean_text']}")
                            elif para_processed:
                                cell_paras.append(para_processed['clean_text'])
                    cell_content = " ".join(cell_paras)
                else:
                    cell_content = cell.text.strip() if hasattr(cell, 'text') else ""

                row_data.append(cell_content)
            table_data.append(row_data)

        return {
            "type": "table",
            "data": table_data
        }

    def extract_chart(self, shape):
        """Extract chart/diagram information for potential Mermaid conversion."""
        try:
            chart = shape.chart
            chart_data = {
                "type": "chart",
                "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
                "title": "",
                "data_points": [],
                "categories": [],
                "series": [],
                "hyperlink": self._extract_shape_hyperlink(shape)
            }

            try:
                if hasattr(chart, 'chart_title') and chart.chart_title and hasattr(chart.chart_title, 'text_frame'):
                    chart_data["title"] = chart.chart_title.text_frame.text.strip()
            except:
                pass

            try:
                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]

                    if hasattr(plot, 'categories') and plot.categories:
                        chart_data["categories"] = [cat.label for cat in plot.categories if hasattr(cat, 'label')]

                    if hasattr(plot, 'series') and plot.series:
                        for series in plot.series:
                            series_data = {
                                "name": series.name if hasattr(series, 'name') else "",
                                "values": []
                            }
                            if hasattr(series, 'values'):
                                try:
                                    series_data["values"] = [val for val in series.values if val is not None]
                                except:
                                    pass
                            chart_data["series"].append(series_data)
            except:
                pass

            return chart_data

        except Exception:
            return {
                "type": "chart",
                "chart_type": "unknown",
                "title": "Chart",
                "data_points": [],
                "categories": [],
                "series": [],
                "hyperlink": self._extract_shape_hyperlink(shape)
            }

    def _create_non_text_content_block(self, shape, shape_info):
        """
        Create content blocks for shapes without text content.
        FIXED: More defensive shape type checking to avoid enum errors.
        """
        try:
            shape_name = getattr(shape, 'name', '')
            shape_type = shape.shape_type
            shape_type_name = str(shape_type).split('.')[-1] if hasattr(shape_type, '__str__') else 'unknown'

            if shape_type_name == 'LINE':
                return {"type": "line", "line_type": "simple",
                        "description": f"Line shape{f': {shape_name}' if shape_name else ''}"}
            elif shape_type_name == 'CONNECTOR':
                return {"type": "line", "line_type": "connector",
                        "description": f"Connector{f': {shape_name}' if shape_name else ''}"}
            elif shape_type_name == 'FREEFORM':
                return {"type": "line", "line_type": "freeform",
                        "description": f"Freeform shape{f': {shape_name}' if shape_name else ''}"}
            elif shape_type_name == 'AUTO_SHAPE':
                auto_shape_type = shape_info.get("auto_shape_type", "unknown")
                if self._is_arrow_shape(auto_shape_type):
                    return {
                        "type": "arrow",
                        "arrow_type": auto_shape_type,
                        "description": f"Arrow ({auto_shape_type}){f': {shape_name}' if shape_name else ''}"
                    }
                else:
                    return {
                        "type": "shape",
                        "shape_subtype": "auto_shape",
                        "description": f"Shape ({auto_shape_type}){f': {shape_name}' if shape_name else ''}"
                    }
            else:
                return {
                    "type": "shape",
                    "shape_subtype": shape_type_name.lower(),
                    "description": f"Shape ({shape_type_name}){f': {shape_name}' if shape_name else ''}"
                }
        except Exception as e:
            print(f"DEBUG: Error creating non-text content block: {e}")
            return {"type": "shape", "shape_subtype": "unknown", "description": "Unknown shape"}

    def _get_shape_analysis_info(self, shape):
        """
        Get basic shape information for diagram analysis and debugging.
        FIXED: More defensive property access to avoid errors.
        """
        shape_info = {
            "shape_type": "unknown",
            "auto_shape_type": None,
            "position": {
                "top": 0,
                "left": 0,
                "width": 0,
                "height": 0
            }
        }

        try:
            if hasattr(shape, 'shape_type'):
                shape_type = shape.shape_type
                shape_info["shape_type"] = str(shape_type).split('.')[-1] if hasattr(shape_type,
                                                                                     '__str__') else 'unknown'
        except Exception as e:
            print(f"DEBUG: Error getting shape type: {e}")
            shape_info["shape_type"] = "unknown"

        try:
            if hasattr(shape, 'auto_shape_type'):
                auto_shape_type = shape.auto_shape_type
                shape_info["auto_shape_type"] = str(auto_shape_type).split('.')[-1] if hasattr(auto_shape_type,
                                                                                               '__str__') else None
        except Exception as e:
            print(f"DEBUG: Error getting auto_shape_type: {e}")
            pass

        try:
            shape_info["position"] = {
                "top": getattr(shape, 'top', 0),
                "left": getattr(shape, 'left', 0),
                "width": getattr(shape, 'width', 0),
                "height": getattr(shape, 'height', 0)
            }
        except Exception as e:
            print(f"DEBUG: Error getting shape position: {e}")
            pass

        return shape_info

    def _is_arrow_shape(self, auto_shape_type):
        """
        Determine if an auto shape is an arrow type.
        FIXED: More defensive type checking.
        """
        if not auto_shape_type:
            return False

        try:
            auto_shape_str = str(auto_shape_type).upper()

            arrow_types = [
                "LEFT_ARROW", "DOWN_ARROW", "UP_ARROW", "RIGHT_ARROW",
                "LEFT_RIGHT_ARROW", "UP_DOWN_ARROW", "QUAD_ARROW",
                "LEFT_RIGHT_UP_ARROW", "BENT_ARROW", "U_TURN_ARROW",
                "CURVED_LEFT_ARROW", "CURVED_RIGHT_ARROW",
                "CURVED_UP_ARROW", "CURVED_DOWN_ARROW",
                "STRIPED_RIGHT_ARROW", "NOTCHED_RIGHT_ARROW",
                "BLOCK_ARC"
            ]

            return any(arrow_type in auto_shape_str for arrow_type in arrow_types)
        except Exception as e:
            print(f"DEBUG: Error checking arrow shape: {e}")
            return False

    def _extract_shape_hyperlink(self, shape):
        """Extract shape-level hyperlinks with URL normalization."""
        try:
            if hasattr(shape, 'click_action') and shape.click_action:
                if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink:
                    if shape.click_action.hyperlink.address:
                        return self._fix_url(shape.click_action.hyperlink.address)
        except:
            pass
        return None

    def _fix_url(self, url):
        """Normalize URLs to handle common PowerPoint URL formatting issues."""
        if not url:
            return url

        if '@' in url and not url.startswith('mailto:'):
            return f"mailto:{url}"

        if not url.startswith(('http://', 'https://', 'mailto:', 'tel:', 'ftp://', '#')):
            if url.startswith('www.') or any(
                    domain in url.lower() for domain in ['.com', '.org', '.net', '.edu', '.gov', '.io']):
                return f"https://{url}"

        return url

