"""
Enhanced PowerPoint Processor Module

This module provides enhanced PowerPoint to Markdown conversion that preserves
formatting like bold, italic, bullet points, etc.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
import re
import os


class EnhancedPowerPointProcessor:
    """Enhanced PowerPoint processor that preserves formatting during conversion."""
    
    def __init__(self):
        self.slide_content = []
        self.current_slide = 1
        
    def process_pptx(self, file_path):
        """
        Process a PowerPoint file and extract formatted content.
        
        Args:
            file_path (str): Path to the PowerPoint file
            
        Returns:
            str: Markdown formatted content
        """
        try:
            prs = Presentation(file_path)
            markdown_content = ""
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides, 1):
                self.current_slide = slide_idx
                slide_markdown = self.process_slide(slide, slide_idx)
                if slide_markdown.strip():
                    markdown_content += slide_markdown + "\n\n---\n\n"
            
            # Remove the last separator
            if markdown_content.endswith("\n\n---\n\n"):
                markdown_content = markdown_content[:-7]
                
            return markdown_content.strip()
            
        except Exception as e:
            raise Exception(f"Error processing PowerPoint file: {str(e)}")
    
    def process_slide(self, slide, slide_number):
        """Process a single slide and extract its content."""
        slide_content = f"# Slide {slide_number}\n\n"
        
        # Group shapes by their vertical position for better reading order
        shapes_with_positions = []
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                shapes_with_positions.append((shape.top, shape.left, shape))
        
        # Sort by top position first, then left position
        shapes_with_positions.sort(key=lambda x: (x[0], x[1]))
        
        # Process shapes in reading order
        for _, _, shape in shapes_with_positions:
            shape_content = self.process_shape(shape)
            if shape_content.strip():
                slide_content += shape_content + "\n\n"
        
        return slide_content
    
    def process_shape(self, shape):
        """Process a single shape and extract its formatted content."""
        content = ""
        
        try:
            # Handle different shape types
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                content = self.process_text_frame(shape.text_frame, "Text Box")
            elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                content = self.process_text_frame(shape.text_frame, "Content")
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                content = self.process_table(shape.table)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                content = self.process_group_shape(shape)
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                content = self.process_text_frame(shape.text_frame)
            elif hasattr(shape, 'text') and shape.text:
                # Fallback for shapes with direct text access
                content = self.clean_text(shape.text)
                
        except Exception as e:
            # If we can't process the shape, try to get basic text
            try:
                if hasattr(shape, 'text') and shape.text:
                    content = self.clean_text(shape.text)
            except:
                pass
        
        return content
    
    def process_group_shape(self, group_shape):
        """Process shapes within a group."""
        content = ""
        
        try:
            for shape in group_shape.shapes:
                shape_content = self.process_shape(shape)
                if shape_content.strip():
                    content += shape_content + "\n"
        except:
            pass
            
        return content
    
    def process_text_frame(self, text_frame, context=""):
        """Process a text frame and extract formatted content."""
        if not text_frame or not text_frame.paragraphs:
            return ""
        
        content = ""
        
        # Check if this looks like a title (first text frame, short content, etc.)
        is_title = self.is_likely_title(text_frame, context)
        
        for para_idx, paragraph in enumerate(text_frame.paragraphs):
            para_content = self.process_paragraph(paragraph, is_title and para_idx == 0)
            if para_content.strip():
                content += para_content + "\n"
        
        return content
    
    def is_likely_title(self, text_frame, context=""):
        """Determine if a text frame is likely a title."""
        if not text_frame.paragraphs:
            return False
            
        # Check various indicators
        first_para = text_frame.paragraphs[0]
        text = first_para.text.strip()
        
        # Short text is more likely to be a title
        if len(text) < 100 and len(text_frame.paragraphs) == 1:
            return True
            
        # Check context clues
        if context.lower() in ["title", "heading"]:
            return True
            
        # Check if it's all caps (common for titles)
        if text.isupper() and len(text) > 3:
            return True
            
        return False
    
    def process_paragraph(self, paragraph, force_header=False):
        """Process a paragraph and extract formatted content."""
        if not paragraph.runs:
            return ""
        
        # Determine paragraph type and formatting
        para_text = paragraph.text.strip()
        if not para_text:
            return ""
        
        # Check for bullet points
        bullet_level = self.get_bullet_level(paragraph)
        if bullet_level >= 0:
            return self.format_bullet_item(paragraph, bullet_level)
        
        # Check for numbered lists
        if self.is_numbered_list(paragraph):
            return self.format_numbered_item(paragraph)
        
        # Process as regular paragraph with inline formatting
        formatted_text = self.process_runs(paragraph.runs)
        
        # Apply paragraph-level formatting
        if force_header or self.is_header_paragraph(paragraph):
            # Determine header level (default to h2 for slide content)
            header_level = 2
            return f"{'#' * header_level} {formatted_text}\n"
        
        return formatted_text
    
    def get_bullet_level(self, paragraph):
        """Determine the bullet point level of a paragraph."""
        try:
            if paragraph.level is not None and paragraph.level >= 0:
                # Check if it actually has bullet formatting
                if hasattr(paragraph, '_p') and paragraph._p is not None:
                    # Look for bullet point indicators in the XML
                    xml_str = str(paragraph._p.xml) if hasattr(paragraph._p, 'xml') else ""
                    if any(bullet in xml_str for bullet in ['buChar', 'buAutoNum', 'buFont']):
                        return paragraph.level
                        
                # Also check the text for bullet-like characters
                text = paragraph.text.strip()
                if text and text[0] in ['•', '·', '-', '*', '◦', '▪', '▫']:
                    return paragraph.level if paragraph.level is not None else 0
                    
        except:
            pass
            
        # Fallback: check if text starts with bullet-like characters
        text = paragraph.text.strip()
        if text and text[0] in ['•', '·', '-', '*', '◦', '▪', '▫']:
            return 0
            
        return -1
    
    def is_numbered_list(self, paragraph):
        """Check if paragraph is part of a numbered list."""
        text = paragraph.text.strip()
        # Look for patterns like "1.", "a)", "(1)", etc.
        numbered_patterns = [
            r'^\d+\.',  # 1. 2. 3.
            r'^\d+\)',  # 1) 2) 3)
            r'^\(\d+\)',  # (1) (2) (3)
            r'^[a-z]\.',  # a. b. c.
            r'^[a-z]\)',  # a) b) c)
            r'^[A-Z]\.',  # A. B. C.
            r'^[ivx]+\.',  # i. ii. iii. (roman numerals)
        ]
        
        for pattern in numbered_patterns:
            if re.match(pattern, text):
                return True
        return False
    
    def format_bullet_item(self, paragraph, level):
        """Format a bullet point item."""
        text = self.process_runs(paragraph.runs)
        
        # Remove leading bullet characters if present
        text = re.sub(r'^[•·\-*◦▪▫]\s*', '', text)
        
        # Create proper markdown bullet with indentation
        indent = "  " * level
        return f"{indent}- {text}"
    
    def format_numbered_item(self, paragraph):
        """Format a numbered list item."""
        text = self.process_runs(paragraph.runs)
        
        # Extract the number/letter and the content
        match = re.match(r'^([0-9a-zA-Z]+[\.\)])?\s*(.*)', text)
        if match:
            number_part = match.group(1) or ""
            content = match.group(2) or text
            
            # Try to determine the numbering style
            if re.match(r'^\d+', number_part):
                # Numeric: convert to markdown numbered list
                return f"1. {content}"
            else:
                # Letters or other: use as-is but format consistently
                return f"1. {text}"
        
        return f"1. {text}"
    
    def is_header_paragraph(self, paragraph):
        """Determine if a paragraph should be treated as a header."""
        text = paragraph.text.strip()
        
        # Short text is more likely to be a header
        if len(text) < 100:
            # Check for header-like formatting or positioning
            try:
                # Check if text is bold (common for headers)
                if paragraph.runs and len(paragraph.runs) > 0:
                    first_run = paragraph.runs[0]
                    if hasattr(first_run.font, 'bold') and first_run.font.bold:
                        return True
            except:
                pass
                
        return False
    
    def process_runs(self, runs):
        """Process text runs and apply inline formatting."""
        if not runs:
            return ""
        
        formatted_text = ""
        
        for run in runs:
            text = run.text
            if not text:
                continue
            
            # Clean up text
            text = self.clean_text(text)
            
            # Apply formatting based on run properties
            formatted_text += self.apply_run_formatting(text, run)
        
        return formatted_text
    
    def apply_run_formatting(self, text, run):
        """Apply formatting to text based on run properties."""
        if not text.strip():
            return text
        
        formatted = text
        
        try:
            font = run.font
            
            # Check for bold
            if hasattr(font, 'bold') and font.bold:
                formatted = f"**{formatted}**"
            
            # Check for italic
            if hasattr(font, 'italic') and font.italic:
                formatted = f"*{formatted}*"
            
            # Check for underline (convert to bold for markdown compatibility)
            if hasattr(font, 'underline') and font.underline:
                if not (hasattr(font, 'bold') and font.bold):  # Don't double-format
                    formatted = f"**{formatted}**"
            
            # Check for code/monospace (this is harder to detect in PowerPoint)
            if hasattr(font, 'name') and font.name:
                monospace_fonts = ['Courier', 'Consolas', 'Monaco', 'Menlo', 'Source Code Pro']
                if any(mono_font in font.name for mono_font in monospace_fonts):
                    formatted = f"`{formatted}`"
            
        except Exception:
            # If we can't access font properties, just return the text
            pass
        
        return formatted
    
    def process_table(self, table):
        """Process a table and convert to markdown format."""
        if not table.rows:
            return ""
        
        markdown_table = ""
        
        # Process table rows
        for row_idx, row in enumerate(table.rows):
            row_content = "|"
            
            for cell in row.cells:
                cell_text = ""
                if cell.text_frame:
                    # Process cell content with formatting
                    for paragraph in cell.text_frame.paragraphs:
                        para_text = self.process_runs(paragraph.runs)
                        if para_text.strip():
                            cell_text += para_text + " "
                
                # Clean up cell text
                cell_text = cell_text.strip().replace('\n', ' ').replace('|', '\\|')
                row_content += f" {cell_text} |"
            
            markdown_table += row_content + "\n"
            
            # Add separator after header row
            if row_idx == 0:
                separator = "|"
                for _ in row.cells:
                    separator += "---------|"
                markdown_table += separator + "\n"
        
        return markdown_table
    
    def clean_text(self, text):
        """Clean and normalize text."""
        if not text:
            return ""
        
        # Replace smart quotes with regular quotes
        text = text.replace('"', '"').replace('"', '"')
        text = text.replace(''', "'").replace(''', "'")
        
        # Replace em and en dashes
        text = text.replace('—', '--').replace('–', '-')
        
        # Normalize whitespace but preserve intentional line breaks
        text = re.sub(r'[ \t]+', ' ', text)
        
        return text.strip()


def convert_pptx_to_markdown_enhanced(file_path):
    """
    Convert a PowerPoint file to markdown with enhanced formatting preservation.
    
    Args:
        file_path (str): Path to the PowerPoint file
        
    Returns:
        str: Markdown formatted content
    """
    processor = EnhancedPowerPointProcessor()
    return processor.process_pptx(file_path)


# Test function to validate the processor
def test_enhanced_processor():
    """Test function for the enhanced processor."""
    # This would be called with a test PowerPoint file
    print("Enhanced PowerPoint processor loaded successfully!")
    print("Features supported:")
    print("- Bold text (**text**)")
    print("- Italic text (*text*)")
    print("- Bullet points with proper nesting")
    print("- Numbered lists")
    print("- Tables with formatting")
    print("- Headers and titles")
    print("- Smart quotes and special characters")


if __name__ == "__main__":
    test_enhanced_processor()