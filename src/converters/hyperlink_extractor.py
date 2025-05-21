"""
Hyperlink Extractor Module

This module contains functions to extract hyperlinks from various document types.
"""

from pptx import Presentation  # For PowerPoint hyperlink extraction
import fitz  # PyMuPDF for PDF processing
import re  # For regex pattern matching
from config import HYPERLINK_CONTEXT_SIZE

def fix_url(url):
    """
    Fix URLs by adding appropriate schemes if missing.

    Args:
        url (str): The URL to fix

    Returns:
        str: Properly formatted URL
    """
    if not url:
        return url

    # For email addresses
    if '@' in url and not url.startswith('mailto:'):
        return f"mailto:{url}"

    # For web URLs
    if not url.startswith(('http://', 'https://', 'mailto:', 'tel:', 'ftp://', '#')):
        if url.startswith('www.') or any(
                domain in url.lower() for domain in ['.com', '.org', '.net', '.edu', '.gov', '.io']):
            return f"https://{url}"

    return url


def extract_pptx_hyperlinks(pptx_path):
    """
    Extract hyperlinks from a PowerPoint file.

    Args:
        pptx_path (str): Path to the PowerPoint file

    Returns:
        list: List of dictionaries containing text, URL, and slide number
    """
    try:
        prs = Presentation(pptx_path)
        raw_hyperlinks = []

        def process_shape(shape, slide_num, parent_type=""):
            """Recursively process shapes, including those in groups"""
            try:
                # Check for hyperlinks in shape click actions
                if hasattr(shape, 'click_action') and shape.click_action is not None:
                    if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink is not None:
                        if shape.click_action.hyperlink.address:
                            shape_text = shape.text if hasattr(shape,
                                                               'text') else f"{parent_type} Shape on Slide {slide_num}"
                            raw_hyperlinks.append({
                                'text': shape_text,
                                'url': shape.click_action.hyperlink.address,
                                'slide': slide_num
                            })
            except Exception:
                # Just skip this shape if there's an error accessing click_action
                pass

            # Process text frames that might contain hyperlinks
            if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    # Create a unique ID for this paragraph
                    p_id = f"{parent_type}shape_{id(shape)}_p_{p_idx}"
                    paragraph_links = {}

                    for run in paragraph.runs:
                        if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                            url = run.hyperlink.address
                            if url not in paragraph_links:
                                paragraph_links[url] = []
                            paragraph_links[url].append(run.text)

                    # Add merged paragraph links
                    for url, text_parts in paragraph_links.items():
                        merged_text = "".join(text_parts)
                        raw_hyperlinks.append({
                            'text': merged_text,
                            'url': url,
                            'slide': slide_num
                        })

            # Process tables
            if hasattr(shape, 'table') and shape.table is not None:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        if hasattr(cell, 'text_frame') and cell.text_frame is not None:
                            for p_idx, paragraph in enumerate(cell.text_frame.paragraphs):
                                p_id = f"{parent_type}table_{id(shape)}_row_{r_idx}_cell_{c_idx}_p_{p_idx}"
                                paragraph_links = {}

                                for run in paragraph.runs:
                                    if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                        url = run.hyperlink.address
                                        if url not in paragraph_links:
                                            paragraph_links[url] = []
                                        paragraph_links[url].append(run.text)

                                # Add merged paragraph links
                                for url, text_parts in paragraph_links.items():
                                    merged_text = "".join(text_parts)
                                    raw_hyperlinks.append({
                                        'text': merged_text,
                                        'url': url,
                                        'slide': slide_num
                                    })

            # Recursively process group shapes
            if hasattr(shape, 'shapes') and shape.shapes is not None:
                try:
                    # For GroupShape objects
                    for i, child_shape in enumerate(shape.shapes):
                        process_shape(child_shape, slide_num, f"{parent_type}Group{i}_")
                except Exception:
                    # Skip if there's an issue accessing child shapes
                    pass

            # Handle SmartArt and Charts which might have text with hyperlinks
            if hasattr(shape, 'has_chart') and shape.has_chart:
                # Charts may have text with hyperlinks in their titles or labels
                try:
                    chart = shape.chart
                    if hasattr(chart, 'chart_title') and chart.chart_title is not None:
                        # Some charts have text properties in their titles
                        if hasattr(chart.chart_title, 'text_frame'):
                            for p_idx, paragraph in enumerate(chart.chart_title.text_frame.paragraphs):
                                for run in paragraph.runs:
                                    if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                        raw_hyperlinks.append({
                                            'text': run.text,
                                            'url': run.hyperlink.address,
                                            'slide': slide_num
                                        })
                except Exception:
                    # Skip if can't access chart data
                    pass

            # Handle SmartArt
            if hasattr(shape, 'has_smart_art') and shape.has_smart_art:
                # Some SmartArt structures may contain text with hyperlinks
                try:
                    # Process SmartArt nodes if they have text frames
                    if hasattr(shape, 'text_frame') and shape.text_frame is not None:
                        for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            for run in paragraph.runs:
                                if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                                    raw_hyperlinks.append({
                                        'text': run.text,
                                        'url': run.hyperlink.address,
                                        'slide': slide_num
                                    })
                except Exception:
                    # Skip if can't access SmartArt data
                    pass

        # Process all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            # Process all shapes in the slide
            for shape in slide.shapes:
                process_shape(shape, slide_num)

        # Final cleanup - remove duplicate links and empty text links
        hyperlinks = []
        seen_urls = {}  # Format: {(url, slide): best_text}

        for link in raw_hyperlinks:
            # Skip empty links
            text = link['text'].strip()
            if not text:
                continue

            # Keep track of the best text for each URL (prefer longer, more descriptive text)
            link_id = (link['url'], link['slide'])
            if link_id not in seen_urls or len(text) > len(seen_urls[link_id]):
                seen_urls[link_id] = text

        # Create final hyperlinks list with best text for each URL
        for (url, slide), text in seen_urls.items():
            hyperlinks.append({
                'text': text,
                'url': url,
                'slide': slide
            })

        return hyperlinks
    except Exception as e:
        print(f"Error extracting PowerPoint hyperlinks: {str(e)}")
        return []


def extract_pdf_hyperlinks(pdf_path):
    """
    Extract hyperlinks from a PDF file using PyMuPDF (fitz).

    Args:
        pdf_path (str): Path to the PDF file

    Returns:
        list: List of dictionaries containing text, URL, and page number
    """
    try:
        hyperlinks = []

        # Open the PDF file with PyMuPDF
        doc = fitz.open(pdf_path)

        # Process each page
        for page_num, page in enumerate(doc, 1):
            # Get links from the page
            links = page.get_links()

            for link in links:
                # Process based on link type
                if link["kind"] == fitz.LINK_URI:  # External URI
                    url = link["uri"]

                    # Try to get text near the link location
                    rect = fitz.Rect(link["from"])
                    # Create an expanded rectangle instead of using inflate method
                    expanded_rect = fitz.Rect(rect.x0 - 10, rect.y0 - 10, rect.x1 + 10, rect.y1 + 10)
                    words = page.get_text("words", clip=expanded_rect)  # Get words near the link

                    # Combine words to form text
                    if words:
                        text = " ".join([w[4] for w in words])
                        # Clean up the text
                        text = text.strip()
                        if not text:
                            text = f"Link on page {page_num}"
                    else:
                        text = f"Link on page {page_num}"

                    hyperlinks.append({
                        'text': text,
                        'url': url,
                        'page': page_num
                    })

                elif link["kind"] == fitz.LINK_GOTO:  # Internal link
                    # The page number this links to
                    dest_page = link["page"] + 1  # 0-based to 1-based

                    # Try to get text near the link location
                    rect = fitz.Rect(link["from"])
                    # Create an expanded rectangle instead of using inflate method
                    expanded_rect = fitz.Rect(rect.x0 - 10, rect.y0 - 10, rect.x1 + 10, rect.y1 + 10)
                    words = page.get_text("words", clip=expanded_rect)

                    if words:
                        text = " ".join([w[4] for w in words])
                        text = text.strip()
                        if not text:
                            text = f"Go to page {dest_page}"
                    else:
                        text = f"Go to page {dest_page}"

                    hyperlinks.append({
                        'text': text,
                        'url': f"#page={dest_page}",
                        'page': page_num
                    })

            # Also look for URLs in the page text
            text = page.get_text()
            if text:
                # Find URLs using regex
                url_pattern = r'https?://(?:[-\w.]|(?:%[\da-fA-F]{2}))+'
                urls = re.findall(url_pattern, text)

                for url in urls:
                    # Check if we already have this URL for this page
                    if not any(link.get('url') == url and link.get('page') == page_num for link in hyperlinks):
                        # Get some context around the URL
                        try:
                            # Find position of URL in text
                            pos = text.find(url)
                            # Get some text before and after
                            start = max(0, pos - 50)
                            end = min(len(text), pos + len(url) + 50)
                            context = text[start:end].replace(url, "").strip()

                            # If context is too long or empty, use a default
                            if len(context) > 100 or not context:
                                context = f"URL on page {page_num}"
                        except:
                            context = f"URL on page {page_num}"

                        hyperlinks.append({
                            'text': context,
                            'url': url,
                            'page': page_num
                        })

        return hyperlinks
    except Exception as e:
        print(f"Error extracting PDF hyperlinks: {str(e)}")
        return []


def format_hyperlinks_section(hyperlinks, container_name="Document"):
    """
    Format hyperlinks into a markdown section.

    Args:
        hyperlinks (list): List of hyperlink dictionaries
        container_name (str): Name of the container (Presentation, Document, etc.)

    Returns:
        str: Formatted markdown string for hyperlinks section
    """
    if not hyperlinks:
        return ""

    # Group hyperlinks by page/slide
    pages_with_links = {}
    url_tracker = {}  # Track URLs we've seen per page

    for link in hyperlinks:
        # Get page number (slide for PowerPoint, page for PDF)
        page_num = link.get('page', link.get('slide', 0))
        if not page_num:
            continue

        # Get URL and text
        url = link.get('url', '')
        text = link.get('text', '').strip()

        # Skip empty entries
        if not url or not text:
            continue

        # Initialize page tracking if needed
        if page_num not in pages_with_links:
            pages_with_links[page_num] = []
            url_tracker[page_num] = set()

        # Fix URL format
        url = fix_url(url)

        # Check for duplicates on this page
        url_key = url.lower()
        if url_key not in url_tracker[page_num]:
            pages_with_links[page_num].append({
                'text': text,
                'url': url
            })
            url_tracker[page_num].add(url_key)
        else:
            # We already have this URL on this page
            # Check if the new text is better than the existing one
            for existing_link in pages_with_links[page_num]:
                if existing_link['url'].lower() == url_key:
                    # Replace with new text if it's better
                    if (len(text) > len(existing_link['text']) and
                        existing_link['text'] in text) or (
                            existing_link['text'].startswith(('Link on page', 'URL on page', 'Go to page'))
                    ):
                        existing_link['text'] = text

    # Exit if no valid links found
    if not any(links for links in pages_with_links.values()):
        return ""

    # Build the markdown
    page_title = "Slide" if container_name == "Presentation" else "Page"
    markdown = f"\n\n## Hyperlinks in {container_name}\n"

    for page_num in sorted(pages_with_links.keys()):
        links = pages_with_links[page_num]
        if not links:
            continue

        markdown += f"\n### {page_title} {page_num}\n"
        for link in links:
            markdown += f"* [{link['text']}]({link['url']})\n"

    return markdown