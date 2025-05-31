"""
PowerPoint Processor Superfile - Complete functionality for extracting and converting PowerPoint content
Combines all features for metadata extraction, content processing, and markdown conversion
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json
import re
import os
from datetime import datetime
import xml.etree.ElementTree as ET


class PowerPointProcessor:
    """Complete PowerPoint processing class with all functionality"""

    def __init__(self):
        self.supported_formats = ['.pptx', '.ppt']

    def convert_pptx_to_markdown_enhanced(self, file_path):
        """
        Main entry point: Convert PowerPoint to structured data, then to markdown with embedded metadata
        """
        try:
            prs = Presentation(file_path)

            # Extract PowerPoint metadata first
            pptx_metadata = self.extract_pptx_metadata(prs, file_path)

            # Extract structured data
            structured_data = self.extract_presentation_data(prs)

            # Convert to basic markdown
            markdown = self.convert_structured_data_to_markdown(structured_data)

            # Add PowerPoint metadata as comments for Claude to use
            markdown_with_metadata = self.add_pptx_metadata_for_claude(markdown, pptx_metadata)

            return markdown_with_metadata
        except Exception as e:
            raise Exception(f"Error processing PowerPoint file: {str(e)}")

    def extract_pptx_metadata(self, presentation, file_path):
        """Extract comprehensive metadata from PowerPoint file"""
        metadata = {}

        try:
            # Core properties from PowerPoint
            core_props = presentation.core_properties

            # Basic file info
            metadata['filename'] = os.path.basename(file_path)
            metadata['file_size'] = os.path.getsize(file_path) if os.path.exists(file_path) else None

            # Document properties
            metadata['title'] = getattr(core_props, 'title', '') or ''
            metadata['author'] = getattr(core_props, 'author', '') or ''
            metadata['subject'] = getattr(core_props, 'subject', '') or ''
            metadata['keywords'] = getattr(core_props, 'keywords', '') or ''
            metadata['comments'] = getattr(core_props, 'comments', '') or ''
            metadata['category'] = getattr(core_props, 'category', '') or ''
            metadata['content_status'] = getattr(core_props, 'content_status', '') or ''
            metadata['language'] = getattr(core_props, 'language', '') or ''
            metadata['version'] = getattr(core_props, 'version', '') or ''

            # Dates
            metadata['created'] = getattr(core_props, 'created', None)
            metadata['modified'] = getattr(core_props, 'modified', None)
            metadata['last_modified_by'] = getattr(core_props, 'last_modified_by', '') or ''
            metadata['last_printed'] = getattr(core_props, 'last_printed', None)

            # Revision and identifier
            metadata['revision'] = getattr(core_props, 'revision', None)
            metadata['identifier'] = getattr(core_props, 'identifier', '') or ''

            # Presentation-specific info
            metadata['slide_count'] = len(presentation.slides)

            # Try to get slide size
            try:
                slide_width = presentation.slide_width
                slide_height = presentation.slide_height
                # Convert from EMUs to inches (914400 EMUs = 1 inch)
                width_inches = round(slide_width / 914400, 2)
                height_inches = round(slide_height / 914400, 2)
                metadata['slide_size'] = f"{width_inches}\" x {height_inches}\""
            except:
                metadata['slide_size'] = ''

            # Try to extract slide master themes/layouts
            try:
                slide_masters = presentation.slide_masters
                if slide_masters:
                    metadata['slide_master_count'] = len(slide_masters)
                    # Get layout names if available
                    layout_names = []
                    for master in slide_masters:
                        for layout in master.slide_layouts:
                            if hasattr(layout, 'name') and layout.name:
                                layout_names.append(layout.name)
                    metadata['layout_types'] = ', '.join(set(layout_names)) if layout_names else ''
            except:
                metadata['slide_master_count'] = 0
                metadata['layout_types'] = ''

            # Application that created the file
            try:
                app_props = presentation.app_properties if hasattr(presentation, 'app_properties') else None
                if app_props:
                    metadata['application'] = getattr(app_props, 'application', '') or ''
                    metadata['app_version'] = getattr(app_props, 'app_version', '') or ''
                    metadata['company'] = getattr(app_props, 'company', '') or ''
                    metadata['doc_security'] = getattr(app_props, 'doc_security', None)
            except:
                metadata['application'] = ''
                metadata['app_version'] = ''
                metadata['company'] = ''
                metadata['doc_security'] = None

        except Exception as e:
            print(f"Warning: Could not extract some metadata: {e}")

        return metadata

    def add_pptx_metadata_for_claude(self, markdown_content, metadata):
        """Add PowerPoint metadata as comments for Claude to incorporate"""
        # Format metadata for Claude
        metadata_comments = "\n<!-- POWERPOINT METADATA FOR CLAUDE:\n"

        if metadata.get('title'):
            metadata_comments += f"Document Title: {metadata['title']}\n"
        if metadata.get('author'):
            metadata_comments += f"Author: {metadata['author']}\n"
        if metadata.get('subject'):
            metadata_comments += f"Subject: {metadata['subject']}\n"
        if metadata.get('keywords'):
            metadata_comments += f"Keywords: {metadata['keywords']}\n"
        if metadata.get('category'):
            metadata_comments += f"Category: {metadata['category']}\n"
        if metadata.get('comments'):
            metadata_comments += f"Document Comments: {metadata['comments']}\n"
        if metadata.get('created'):
            metadata_comments += f"Created Date: {metadata['created']}\n"
        if metadata.get('modified'):
            metadata_comments += f"Last Modified: {metadata['modified']}\n"
        if metadata.get('last_modified_by'):
            metadata_comments += f"Last Modified By: {metadata['last_modified_by']}\n"
        if metadata.get('version'):
            metadata_comments += f"Version: {metadata['version']}\n"
        if metadata.get('application'):
            metadata_comments += f"Created With: {metadata['application']}\n"
        if metadata.get('company'):
            metadata_comments += f"Company: {metadata['company']}\n"
        if metadata.get('language'):
            metadata_comments += f"Language: {metadata['language']}\n"
        if metadata.get('content_status'):
            metadata_comments += f"Content Status: {metadata['content_status']}\n"

        # File info
        metadata_comments += f"Filename: {metadata.get('filename', 'unknown')}\n"
        metadata_comments += f"Slide Count: {metadata.get('slide_count', 0)}\n"
        if metadata.get('slide_size'):
            metadata_comments += f"Slide Size: {metadata['slide_size']}\n"
        if metadata.get('layout_types'):
            metadata_comments += f"Layout Types: {metadata['layout_types']}\n"

        metadata_comments += "-->\n"

        # Add metadata at the beginning
        return metadata_comments + markdown_content

    def extract_presentation_data(self, presentation):
        """Extract all content with minimal processing"""
        data = {
            "total_slides": len(presentation.slides),
            "slides": []
        }

        for slide_idx, slide in enumerate(presentation.slides, 1):
            slide_data = self.extract_slide_data(slide, slide_idx)
            data["slides"].append(slide_data)

        return data

    def extract_slide_data(self, slide, slide_number):
        """Extract slide content in reading order"""
        # Get shapes in reading order
        positioned_shapes = []
        for shape in slide.shapes:
            if hasattr(shape, 'top') and hasattr(shape, 'left'):
                positioned_shapes.append((shape.top, shape.left, shape))
            else:
                positioned_shapes.append((0, 0, shape))

        positioned_shapes.sort(key=lambda x: (x[0], x[1]))

        slide_data = {
            "slide_number": slide_number,
            "content_blocks": []
        }

        for _, _, shape in positioned_shapes:
            block = self.extract_shape_content_simple(shape)
            if block:
                slide_data["content_blocks"].append(block)

        return slide_data

    def extract_shape_content_simple(self, shape):
        """Extract shape content without complex formatting logic"""
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return self.extract_image_simple(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            return self.extract_table_simple(shape.table)
        elif hasattr(shape, 'has_chart') and shape.has_chart:
            return self.extract_chart_simple(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return self.extract_group_simple(shape)
        elif hasattr(shape, 'text_frame') and shape.text_frame:
            return self.extract_text_frame_simple(shape.text_frame, shape)
        elif hasattr(shape, 'text') and shape.text:
            return self.extract_plain_text_simple(shape)
        return None

    def extract_chart_simple(self, shape):
        """Extract chart/diagram information for potential Mermaid conversion"""
        try:
            chart = shape.chart
            chart_data = {
                "type": "chart",
                "chart_type": str(chart.chart_type) if hasattr(chart, 'chart_type') else "unknown",
                "title": "",
                "data_points": [],
                "categories": [],
                "series": [],
                "hyperlink": self.extract_shape_hyperlink(shape)
            }

            # Try to get chart title
            try:
                if hasattr(chart, 'chart_title') and chart.chart_title and hasattr(chart.chart_title, 'text_frame'):
                    chart_data["title"] = chart.chart_title.text_frame.text.strip()
            except:
                pass

            # Try to extract data for potential Mermaid conversion
            try:
                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]
                    if hasattr(plot, 'categories') and plot.categories:
                        chart_data["categories"] = [cat.label for cat in plot.categories if hasattr(cat, 'label')]

                    if hasattr(plot, 'series') and plot.series:
                        for series in plot.series:
                            series_data = {
                                "name": series.name if hasattr(series, 'name') else "",
                                "values": []
                            }
                            if hasattr(series, 'values'):
                                try:
                                    series_data["values"] = [val for val in series.values if val is not None]
                                except:
                                    pass
                            chart_data["series"].append(series_data)
            except:
                pass

            return chart_data

        except Exception:
            # Fallback for charts we can't parse
            return {
                "type": "chart",
                "chart_type": "unknown",
                "title": "Chart",
                "data_points": [],
                "categories": [],
                "series": [],
                "hyperlink": self.extract_shape_hyperlink(shape)
            }

    def extract_group_simple(self, shape):
        """Extract content from grouped shapes - potential diagrams"""
        try:
            group_data = {
                "type": "group",
                "shapes": [],
                "connections": [],
                "hyperlink": self.extract_shape_hyperlink(shape)
            }

            # Process shapes in the group
            for child_shape in shape.shapes:
                child_data = self.extract_shape_content_simple(child_shape)
                if child_data:
                    # Add position information for diagram analysis
                    if hasattr(child_shape, 'top') and hasattr(child_shape, 'left'):
                        child_data["position"] = {
                            "top": child_shape.top,
                            "left": child_shape.left,
                            "width": getattr(child_shape, 'width', 0),
                            "height": getattr(child_shape, 'height', 0)
                        }
                    group_data["shapes"].append(child_data)

            # Analyze for potential diagram patterns
            group_data["diagram_type"] = self.analyze_diagram_pattern(group_data["shapes"])

            return group_data

        except Exception:
            return None

    def analyze_diagram_pattern(self, shapes):
        """Analyze shapes to determine if they form a recognizable diagram pattern"""
        if not shapes:
            return "unknown"

        text_shapes = [s for s in shapes if s.get("type") == "text"]

        # Look for flowchart patterns
        flowchart_keywords = [
            "start", "end", "begin", "finish", "process", "decision", "if", "then", "else",
            "input", "output", "step", "stage", "phase", "flow", "next", "previous"
        ]

        # Look for organizational chart patterns
        org_keywords = [
            "manager", "director", "ceo", "cto", "team", "lead", "reports to",
            "department", "division", "head", "supervisor", "employee"
        ]

        # Look for process diagram patterns
        process_keywords = [
            "workflow", "procedure", "method", "sequence", "order", "first", "second",
            "last", "finally", "initial", "final", "stage", "milestone"
        ]

        # Look for network/system diagram patterns
        network_keywords = [
            "server", "database", "client", "network", "connection", "api", "service",
            "component", "module", "system", "interface", "protocol"
        ]

        # Analyze text content
        all_text = " ".join([
            shape.get("content", {}).get("clean_text", "") if isinstance(shape.get("content"), dict)
            else str(shape.get("content", ""))
            for shape in text_shapes
        ]).lower()

        # Count keyword matches
        flowchart_score = sum(1 for keyword in flowchart_keywords if keyword in all_text)
        org_score = sum(1 for keyword in org_keywords if keyword in all_text)
        process_score = sum(1 for keyword in process_keywords if keyword in all_text)
        network_score = sum(1 for keyword in network_keywords if keyword in all_text)

        # Determine diagram type based on highest score
        scores = {
            "flowchart": flowchart_score,
            "org_chart": org_score,
            "process": process_score,
            "network": network_score
        }

        max_score = max(scores.values())
        if max_score >= 2:  # Require at least 2 keyword matches
            return max(scores, key=scores.get)

        # Fallback: analyze shape arrangement
        if len(shapes) >= 3:
            # Check if shapes are arranged in a potentially hierarchical way
            if self.has_hierarchical_arrangement(shapes):
                return "hierarchy"
            elif self.has_linear_arrangement(shapes):
                return "sequence"
            else:
                return "diagram"

        return "unknown"

    def has_hierarchical_arrangement(self, shapes):
        """Check if shapes are arranged hierarchically (org chart, tree structure)"""
        positioned_shapes = [s for s in shapes if "position" in s]
        if len(positioned_shapes) < 3:
            return False

        # Sort by vertical position
        sorted_by_top = sorted(positioned_shapes, key=lambda x: x["position"]["top"])

        # Check if there are clear levels (groups of shapes at similar heights)
        levels = []
        current_level = []
        tolerance = 50  # EMU tolerance for same level

        for shape in sorted_by_top:
            if not current_level:
                current_level.append(shape)
            elif abs(shape["position"]["top"] - current_level[0]["position"]["top"]) <= tolerance:
                current_level.append(shape)
            else:
                levels.append(current_level)
                current_level = [shape]

        if current_level:
            levels.append(current_level)

        # Hierarchical if we have at least 2 levels with the top having fewer items
        return len(levels) >= 2 and len(levels[0]) <= len(levels[1])

    def has_linear_arrangement(self, shapes):
        """Check if shapes are arranged in a linear sequence"""
        positioned_shapes = [s for s in shapes if "position" in s]
        if len(positioned_shapes) < 3:
            return False

        # Check for primarily horizontal or vertical arrangement
        positions = [(s["position"]["left"], s["position"]["top"]) for s in positioned_shapes]

        # Calculate variance in horizontal vs vertical positioning
        lefts = [p[0] for p in positions]
        tops = [p[1] for p in positions]

        left_variance = max(lefts) - min(lefts) if lefts else 0
        top_variance = max(tops) - min(tops) if tops else 0

        # Linear if one dimension has much more variance than the other
        return (left_variance > top_variance * 2) or (top_variance > left_variance * 2)

    def extract_text_frame_simple(self, text_frame, shape):
        """Extract text with basic hints for Claude"""
        if not text_frame.paragraphs:
            return None

        block = {
            "type": "text",
            "paragraphs": [],
            "shape_hyperlink": self.extract_shape_hyperlink(shape)
        }

        for para in text_frame.paragraphs:
            if not para.text.strip():
                continue

            para_data = {
                "raw_text": para.text,  # Preserve original spacing/indentation
                "clean_text": para.text.strip(),
                "hints": {
                    "has_powerpoint_level": hasattr(para, 'level') and para.level is not None,
                    "powerpoint_level": getattr(para, 'level', None),
                    "indented": len(para.text) != len(para.text.lstrip()),
                    "starts_with_bullet": para.text.strip() and para.text.strip()[0] in '•◦▪▫‣·-*+',
                    "starts_with_number": bool(re.match(r'^\s*\d+[\.\)]\s', para.text)),
                    "short_text": len(para.text.strip()) < 100,
                    "all_caps": para.text.strip().isupper() if para.text.strip() else False
                },
                "formatted_runs": self.extract_runs_simple(para.runs)
            }

            block["paragraphs"].append(para_data)

        return block

    def extract_runs_simple(self, runs):
        """Extract formatting and hyperlinks from runs"""
        formatted_runs = []

        for run in runs:
            run_data = {
                "text": run.text,
                "bold": False,
                "italic": False,
                "hyperlink": None
            }

            # Get formatting
            try:
                font = run.font
                if hasattr(font, 'bold') and font.bold:
                    run_data["bold"] = True
                if hasattr(font, 'italic') and font.italic:
                    run_data["italic"] = True
            except:
                pass

            # Get hyperlinks
            try:
                if hasattr(run, 'hyperlink') and run.hyperlink and run.hyperlink.address:
                    run_data["hyperlink"] = self.fix_url(run.hyperlink.address)
            except:
                pass

            formatted_runs.append(run_data)

        return formatted_runs

    def extract_image_simple(self, shape):
        """Extract image info with proper alt text extraction"""
        alt_text = "Image"

        try:
            # Try multiple methods to get alt text
            if hasattr(shape, 'alt_text') and shape.alt_text:
                alt_text = shape.alt_text
            elif hasattr(shape, 'image') and hasattr(shape.image, 'alt_text') and shape.image.alt_text:
                alt_text = shape.image.alt_text
            elif hasattr(shape, '_element'):
                # Try to extract from XML
                try:
                    xml_str = str(shape._element.xml) if hasattr(shape._element, 'xml') else ""
                    if xml_str:
                        root = ET.fromstring(xml_str)
                        # Look for description attributes
                        for elem in root.iter():
                            if 'descr' in elem.attrib and elem.attrib['descr']:
                                alt_text = elem.attrib['descr']
                                break
                            elif 'title' in elem.attrib and elem.attrib['title']:
                                alt_text = elem.attrib['title']
                                break
                except:
                    pass
        except:
            pass

        return {
            "type": "image",
            "alt_text": alt_text.strip() if alt_text else "Image",
            "hyperlink": self.extract_shape_hyperlink(shape)
        }

    def extract_table_simple(self, table):
        """Extract table data"""
        if not table.rows:
            return None

        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                # Extract cell text with basic formatting
                cell_content = ""
                if hasattr(cell, 'text_frame') and cell.text_frame:
                    for para in cell.text_frame.paragraphs:
                        if para.text.strip():
                            cell_content += para.text.strip() + " "
                else:
                    cell_content = cell.text
                row_data.append(cell_content.strip())
            table_data.append(row_data)

        return {
            "type": "table",
            "data": table_data
        }

    def extract_plain_text_simple(self, shape):
        """Extract plain text from shape"""
        return {
            "type": "text",
            "paragraphs": [{
                "raw_text": shape.text,
                "clean_text": shape.text.strip(),
                "hints": {
                    "has_powerpoint_level": False,
                    "powerpoint_level": None,
                    "indented": False,
                    "starts_with_bullet": shape.text.strip() and shape.text.strip()[0] in '•◦▪▫‣·-*+',
                    "starts_with_number": bool(re.match(r'^\s*\d+[\.\)]\s', shape.text)),
                    "short_text": len(shape.text.strip()) < 100,
                    "all_caps": shape.text.strip().isupper() if shape.text.strip() else False
                },
                "formatted_runs": [{"text": shape.text, "bold": False, "italic": False, "hyperlink": None}]
            }],
            "shape_hyperlink": self.extract_shape_hyperlink(shape)
        }

    def extract_shape_hyperlink(self, shape):
        """Extract shape-level hyperlink"""
        try:
            if hasattr(shape, 'click_action') and shape.click_action:
                if hasattr(shape.click_action, 'hyperlink') and shape.click_action.hyperlink:
                    if shape.click_action.hyperlink.address:
                        return self.fix_url(shape.click_action.hyperlink.address)
        except:
            pass
        return None

    def fix_url(self, url):
        """Fix URLs by adding schemes if missing"""
        if not url:
            return url

        if '@' in url and not url.startswith('mailto:'):
            return f"mailto:{url}"

        if not url.startswith(('http://', 'https://', 'mailto:', 'tel:', 'ftp://', '#')):
            if url.startswith('www.') or any(
                    domain in url.lower() for domain in ['.com', '.org', '.net', '.edu', '.gov', '.io']):
                return f"https://{url}"

        return url

    def convert_structured_data_to_markdown(self, data):
        """Convert structured data to basic markdown (Claude will enhance)"""
        markdown_parts = []

        for slide in data["slides"]:
            # Add slide marker
            markdown_parts.append(f"\n<!-- Slide {slide['slide_number']} -->\n")

            for block in slide["content_blocks"]:
                if block["type"] == "text":
                    markdown_parts.append(self.convert_text_block_to_markdown(block))
                elif block["type"] == "table":
                    markdown_parts.append(self.convert_table_to_markdown(block))
                elif block["type"] == "image":
                    markdown_parts.append(self.convert_image_to_markdown(block))
                elif block["type"] == "chart":
                    markdown_parts.append(self.convert_chart_to_markdown(block))
                elif block["type"] == "group":
                    markdown_parts.append(self.convert_group_to_markdown(block))

        return "\n\n".join(filter(None, markdown_parts))

    def convert_text_block_to_markdown(self, block):
        """Convert text block to basic markdown"""
        lines = []

        for para in block["paragraphs"]:
            line = self.convert_paragraph_to_markdown(para)
            if line:
                lines.append(line)

        # If entire shape is a hyperlink, wrap it
        result = "\n".join(lines)
        if block.get("shape_hyperlink") and result:
            result = f"[{result}]({block['shape_hyperlink']})"

        return result

    def convert_chart_to_markdown(self, block):
        """Convert chart to markdown with Mermaid diagram suggestion"""
        chart_md = f"**Chart: {block.get('title', 'Untitled Chart')}**\n"
        chart_md += f"*Chart Type: {block.get('chart_type', 'unknown')}*\n\n"

        # Add data if available
        if block.get('categories') and block.get('series'):
            chart_md += "Data:\n"
            for series in block['series']:
                if series.get('name'):
                    chart_md += f"- {series['name']}: "
                    if series.get('values'):
                        chart_md += ", ".join(map(str, series['values'][:5]))  # Limit to first 5 values
                        if len(series['values']) > 5:
                            chart_md += "..."
                    chart_md += "\n"

        # Add comment for Claude to potentially convert to Mermaid
        chart_md += f"\n<!-- DIAGRAM_CANDIDATE: chart, type={block.get('chart_type', 'unknown')} -->\n"

        if block.get("hyperlink"):
            chart_md = f"[{chart_md}]({block['hyperlink']})"

        return chart_md

    def convert_group_to_markdown(self, block):
        """Convert grouped shapes to markdown with diagram analysis"""
        diagram_type = block.get("diagram_type", "unknown")

        # Start with diagram identification
        group_md = f"**Diagram ({diagram_type})**\n\n"

        # Convert individual shapes
        shape_content = []
        for shape in block.get("shapes", []):
            if shape.get("type") == "text" and shape.get("content"):
                content = self.convert_text_block_to_markdown(shape)
                if content:
                    shape_content.append(content)
            elif shape.get("type") == "image":
                image_md = self.convert_image_to_markdown(shape)
                if image_md:
                    shape_content.append(image_md)

        if shape_content:
            group_md += "\n".join(shape_content)

        # Add diagram conversion hint for Claude
        group_md += f"\n\n<!-- DIAGRAM_CANDIDATE: {diagram_type}, shapes={len(block.get('shapes', []))} -->\n"

        if block.get("hyperlink"):
            group_md = f"[{group_md}]({block['hyperlink']})"

        return group_md

    def convert_paragraph_to_markdown(self, para):
        """Convert paragraph to basic markdown with formatting"""
        if not para["clean_text"]:
            return ""

        # Build formatted text from runs
        formatted_text = ""
        for run in para["formatted_runs"]:
            text = run["text"]
            if not text:
                continue

            # Apply formatting
            if run["bold"] and run["italic"]:
                text = f"***{text}***"
            elif run["bold"]:
                text = f"**{text}**"
            elif run["italic"]:
                text = f"*{text}*"

            # Apply hyperlink
            if run["hyperlink"]:
                text = f"[{text}]({run['hyperlink']})"

            formatted_text += text

        # Basic structure hints for Claude
        hints = para["hints"]

        # Very simple formatting - Claude will fix the structure
        if hints["starts_with_bullet"] or (hints["has_powerpoint_level"] and hints["powerpoint_level"] is not None):
            # Simple bullet - Claude will fix the nesting
            clean_text = formatted_text
            # Remove existing bullet chars
            clean_text = re.sub(r'^[•◦▪▫‣·\-\*\+]\s*', '', clean_text.strip())
            return f"- {clean_text}"
        elif hints["starts_with_number"]:
            # Simple numbered item
            clean_text = re.sub(r'^\s*\d+[\.\)]\s*', '', formatted_text)
            return f"1. {clean_text}"
        elif hints["short_text"] and hints["all_caps"]:
            # Likely a header
            return f"## {formatted_text}"
        else:
            # Regular paragraph
            return formatted_text

    def convert_table_to_markdown(self, block):
        """Convert table to markdown"""
        if not block["data"]:
            return ""

        markdown = ""
        for i, row in enumerate(block["data"]):
            markdown += "| " + " | ".join(cell.replace("|", "\\|") for cell in row) + " |\n"

            # Add separator after header
            if i == 0:
                markdown += "| " + " | ".join("---" for _ in row) + " |\n"

        return markdown

    def convert_image_to_markdown(self, block):
        """Convert image to markdown"""
        image_md = f"![{block['alt_text']}](image)"

        if block.get("hyperlink"):
            image_md = f"[{image_md}]({block['hyperlink']})"

        return image_md

    # Additional utility methods for enhanced functionality

    def validate_file(self, file_path):
        """Validate that the file exists and is a supported PowerPoint format"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.supported_formats:
            raise ValueError(
                f"Unsupported file format: {file_ext}. Supported formats: {', '.join(self.supported_formats)}")

        return True

    def get_presentation_summary(self, presentation):
        """Get a quick summary of the presentation structure"""
        summary = {
            "total_slides": len(presentation.slides),
            "slide_details": []
        }

        for idx, slide in enumerate(presentation.slides, 1):
            slide_info = {
                "slide_number": idx,
                "shape_count": len(slide.shapes),
                "text_shapes": 0,
                "image_shapes": 0,
                "table_shapes": 0,
                "chart_shapes": 0,
                "group_shapes": 0
            }

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_info["image_shapes"] += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    slide_info["table_shapes"] += 1
                elif hasattr(shape, 'has_chart') and shape.has_chart:
                    slide_info["chart_shapes"] += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    slide_info["group_shapes"] += 1
                elif hasattr(shape, 'text_frame') or hasattr(shape, 'text'):
                    slide_info["text_shapes"] += 1

            summary["slide_details"].append(slide_info)

        return summary

    def extract_all_text(self, presentation):
        """Extract all text content from the presentation for text analysis"""
        all_text = []

        for slide_idx, slide in enumerate(presentation.slides, 1):
            slide_text = {
                "slide_number": slide_idx,
                "text_content": []
            }

            for shape in slide.shapes:
                text_content = self._extract_text_from_shape(shape)
                if text_content:
                    slide_text["text_content"].append(text_content)

            all_text.append(slide_text)

        return all_text

    def _extract_text_from_shape(self, shape):
        """Helper method to extract text from any shape type"""
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text_parts = []
            for para in shape.text_frame.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())
            return " ".join(text_parts) if text_parts else None
        elif hasattr(shape, 'text') and shape.text:
            return shape.text.strip() if shape.text.strip() else None
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # Recursively extract text from grouped shapes
            group_text = []
            for child_shape in shape.shapes:
                child_text = self._extract_text_from_shape(child_shape)
                if child_text:
                    group_text.append(child_text)
            return " ".join(group_text) if group_text else None
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            # Extract text from table cells
            table_text = []
            for row in shape.table.rows:
                row_text = []
                for cell in row.cells:
                    if hasattr(cell, 'text_frame') and cell.text_frame:
                        cell_content = []
                        for para in cell.text_frame.paragraphs:
                            if para.text.strip():
                                cell_content.append(para.text.strip())
                        if cell_content:
                            row_text.append(" ".join(cell_content))
                    elif cell.text and cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    table_text.append(" | ".join(row_text))
            return "\n".join(table_text) if table_text else None

        return None

    def export_to_json(self, presentation, file_path=None):
        """Export the extracted presentation data to JSON format"""
        structured_data = self.extract_presentation_data(presentation)
        metadata = self.extract_pptx_metadata(presentation, file_path or "unknown")

        export_data = {
            "metadata": metadata,
            "content": structured_data,
            "export_timestamp": datetime.now().isoformat(),
            "processor_version": "superfile_v1.0"
        }

        return json.dumps(export_data, indent=2, default=str)

    def process_file_complete(self, file_path, output_format="markdown"):
        """
        Complete file processing with multiple output options

        Args:
            file_path (str): Path to the PowerPoint file
            output_format (str): "markdown", "json", "text", or "summary"

        Returns:
            dict: Contains the processed content and metadata
        """
        # Validate file
        self.validate_file(file_path)

        # Load presentation
        prs = Presentation(file_path)

        # Extract all data
        metadata = self.extract_pptx_metadata(prs, file_path)
        structured_data = self.extract_presentation_data(prs)
        summary = self.get_presentation_summary(prs)
        all_text = self.extract_all_text(prs)

        result = {
            "file_path": file_path,
            "metadata": metadata,
            "summary": summary,
            "processing_timestamp": datetime.now().isoformat()
        }

        if output_format == "markdown":
            markdown = self.convert_structured_data_to_markdown(structured_data)
            result["content"] = self.add_pptx_metadata_for_claude(markdown, metadata)
        elif output_format == "json":
            result["content"] = structured_data
            result["json_export"] = self.export_to_json(prs, file_path)
        elif output_format == "text":
            result["content"] = all_text
        elif output_format == "summary":
            result["content"] = {
                "summary": summary,
                "key_points": self._extract_key_points(all_text),
                "word_count": self._count_words(all_text)
            }
        else:
            raise ValueError(f"Unsupported output format: {output_format}")

        return result

    def _extract_key_points(self, all_text):
        """Extract potential key points from text content"""
        key_points = []

        for slide_text in all_text:
            for text_content in slide_text["text_content"]:
                if text_content:
                    # Look for bullet points or numbered lists
                    lines = text_content.split('\n')
                    for line in lines:
                        line = line.strip()
                        if (line and
                                (line.startswith(('•', '-', '*', '◦', '▪')) or
                                 re.match(r'^\d+[\.\)]\s', line) or
                                 len(line) < 100)):  # Short lines might be key points
                            key_points.append({
                                "slide": slide_text["slide_number"],
                                "text": line
                            })

        return key_points

    def _count_words(self, all_text):
        """Count total words in the presentation"""
        total_words = 0

        for slide_text in all_text:
            for text_content in slide_text["text_content"]:
                if text_content:
                    words = len(text_content.split())
                    total_words += words

        return total_words


# Convenience functions for backward compatibility and ease of use

def convert_pptx_to_markdown_enhanced(file_path):
    """
    Convenience function to maintain backward compatibility
    """
    processor = PowerPointProcessor()
    return processor.convert_pptx_to_markdown_enhanced(file_path)


def process_powerpoint_file(file_path, output_format="markdown"):
    """
    Convenience function for complete file processing

    Args:
        file_path (str): Path to the PowerPoint file
        output_format (str): "markdown", "json", "text", or "summary"

    Returns:
        dict: Processed content and metadata
    """
    processor = PowerPointProcessor()
    return processor.process_file_complete(file_path, output_format)


# Example usage and testing functions

def example_usage():
    """
    Example of how to use the PowerPointProcessor
    """
    # Initialize processor
    processor = PowerPointProcessor()

    # Example file path (replace with actual file)
    file_path = "example_presentation.pptx"

    try:
        # Basic markdown conversion
        markdown_result = processor.convert_pptx_to_markdown_enhanced(file_path)
        print("Markdown conversion completed")

        # Complete processing with summary
        summary_result = processor.process_file_complete(file_path, "summary")
        print(f"Presentation has {summary_result['summary']['total_slides']} slides")
        print(f"Total words: {summary_result['content']['word_count']}")

        # JSON export
        json_result = processor.process_file_complete(file_path, "json")
        print("JSON export completed")

    except Exception as e:
        print(f"Error processing file: {e}")


if __name__ == "__main__":
    # Run example usage if script is executed directly
    print("PowerPoint Processor Superfile")
    print("=" * 40)
    print("Available methods:")
    print("- convert_pptx_to_markdown_enhanced(file_path)")
    print("- process_file_complete(file_path, output_format)")
    print("- PowerPointProcessor class for advanced usage")
    print("\nExample usage:")
    print("processor = PowerPointProcessor()")
    print("result = processor.process_file_complete('file.pptx', 'markdown')")
    print("\nFor backward compatibility:")
    print("markdown = convert_pptx_to_markdown_enhanced('file.pptx')")